"""Tests for font-decoration (太字/下線/イタリック) extraction — SOT-2564.

Offline, network-free: synthetic docx/xlsx/pptx fixtures exercise the detection, the ``require`` filter
(bold∧underline∧italic), the flag gating (byte-identical read_office face when off), and the routing /
investigator wiring. The pdf glyph-run path is covered by the focused offline verification on the real
corpus (see docs/ai/font_emphasis_SOT-2564.md); it needs no synthetic PDF here.
"""
from pathlib import Path

import docx
import openpyxl
import pytest
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches

from src.rag.corpus import FileRef
from src.rag.tools import font_emphasis as fe


def _ref(path: Path, ext: str) -> FileRef:
    return FileRef(path=path, project="p", category="plan", rel=path.name, name=path.name, ext=ext)


# --------------------------------------------------------------------------- require normalisation
def test_norm_require_synonyms_and_separators() -> None:
    assert fe._norm_require("太字,下線,イタリック") == frozenset({"bold", "underline", "italic"})
    assert fe._norm_require("bold+underline+italic") == frozenset({"bold", "underline", "italic"})
    assert fe._norm_require(["斜体", "アンダーライン"]) == frozenset({"italic", "underline"})
    assert fe._norm_require("") == frozenset()
    assert fe._norm_require(None) == frozenset()
    # an unknown word is ignored (fail-open to no-such-constraint), never a silent match-all
    assert fe._norm_require("虹色") == frozenset()


def test_enabled_default_off(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.delenv("RAG_FONT_EMPHASIS", raising=False)
    assert fe.enabled() is False
    monkeypatch.setenv("RAG_FONT_EMPHASIS", "1")
    assert fe.enabled() is True


# --------------------------------------------------------------------------- xlsx
def test_xlsx_require_all_three_isolates_cell(tmp_path: Path) -> None:
    path = tmp_path / "t.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "普通"
    ws["A2"] = "太字だけ"
    ws["A2"].font = Font(bold=True)
    ws["A3"] = "4,675,000円"
    ws["A3"].font = Font(bold=True, underline="single", italic=True)
    wb.save(path)

    all_items = fe.font_emphasis(_ref(path, "xlsx"))["value"]
    # only decorated cells surface (the plain A1 is dropped)
    assert {it["value"] for it in all_items} == {"太字だけ", "4,675,000円"}

    hit = fe.emphasized_cells(_ref(path, "xlsx"), require="太字,下線,イタリック")
    assert hit == ["4,675,000円"]


# --------------------------------------------------------------------------- docx
def test_docx_merges_consecutive_same_decoration_runs(tmp_path: Path) -> None:
    path = tmp_path / "t.docx"
    d = docx.Document()
    p = d.add_paragraph()
    r1 = p.add_run("4,675,")
    r2 = p.add_run("000円")
    for r in (r1, r2):
        r.bold = True
        r.underline = True
        r.italic = True
    p.add_run("（税込）")  # undecorated tail must not join the span
    d.save(path)

    hit = fe.emphasized_cells(_ref(path, "docx"), require=["bold", "underline", "italic"])
    assert hit == ["4,675,000円"]


def test_docx_partial_decoration_excluded_by_require(tmp_path: Path) -> None:
    path = tmp_path / "t.docx"
    d = docx.Document()
    p = d.add_paragraph()
    r = p.add_run("太字と下線のみ")
    r.bold = True
    r.underline = True  # no italic
    d.save(path)

    assert fe.emphasized_cells(_ref(path, "docx"), require="太字,下線,イタリック") == []
    assert fe.emphasized_cells(_ref(path, "docx"), require="太字,下線") == ["太字と下線のみ"]


# --------------------------------------------------------------------------- pptx
def test_pptx_run_decorations(tmp_path: Path) -> None:
    path = tmp_path / "t.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "強調値"
    run.font.bold = True
    run.font.italic = True
    run.font.underline = True
    prs.save(path)

    assert fe.emphasized_cells(_ref(path, "pptx"),
                               require=["bold", "underline", "italic"]) == ["強調値"]


# --------------------------------------------------------------------------- unsupported ext
def test_unsupported_extension_raises(tmp_path: Path) -> None:
    path = tmp_path / "t.csv"
    path.write_text("a,b\n1,2\n", encoding="utf-8")
    with pytest.raises(fe.ContractError):
        fe.font_emphasis(_ref(path, "csv"))


# --------------------------------------------------------------------------- read_office face gating
def test_read_office_face_byte_identical_when_flag_off(tmp_path: Path,
                                                       monkeypatch: pytest.MonkeyPatch) -> None:
    from src.rag.extract.office import extract_xlsx

    path = tmp_path / "t.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "4,675,000円"
    ws["A1"].font = Font(bold=True, underline="single", italic=True)
    wb.save(path)
    ref = _ref(path, "xlsx")

    monkeypatch.delenv("RAG_FONT_EMPHASIS", raising=False)
    off = extract_xlsx(ref, None)
    assert "【書式強調" not in off

    monkeypatch.setenv("RAG_FONT_EMPHASIS", "1")
    on = extract_xlsx(ref, None)
    assert "【書式強調(太字/下線/イタリック)のある箇所】" in on
    assert "太字∧下線∧イタリック" in on
    assert "4,675,000円" in on


# --------------------------------------------------------------------------- investigator tool gating
def test_investigator_tool_registered_only_when_flag_on(monkeypatch: pytest.MonkeyPatch) -> None:
    from src.rag.agent.investigator import build_generic_tools
    from src.rag.tools.profile import CorpusProfile

    profile = CorpusProfile()
    monkeypatch.delenv("RAG_FONT_EMPHASIS", raising=False)
    names_off = {t.name for t in build_generic_tools(profile)}
    assert "font_emphasis" not in names_off

    monkeypatch.setenv("RAG_FONT_EMPHASIS", "1")
    names_on = {t.name for t in build_generic_tools(profile)}
    assert "font_emphasis" in names_on


# --------------------------------------------------------------------------- routing gating
def test_routing_leads_with_font_emphasis_when_flag_on(monkeypatch: pytest.MonkeyPatch) -> None:
    from src.rag.agent import question_contract as qc
    from src.rag.agent.routing import first_tools_for

    q = "報告資料の中で、太字、下線、イタリックのすべてに該当する箇所を抽出してください。"
    contract = qc.QuestionContract(
        contract=qc.FORMAT_CHECK, label="書式", archetype="highlight_set",
        completion_conditions=(), route="tool", method="deterministic", confidence=1.0)

    monkeypatch.delenv("RAG_FONT_EMPHASIS", raising=False)
    off = first_tools_for(contract, q)
    assert not off or off[0] != "font_emphasis"

    monkeypatch.setenv("RAG_FONT_EMPHASIS", "1")
    on = first_tools_for(contract, q)
    assert on[0] == "font_emphasis"

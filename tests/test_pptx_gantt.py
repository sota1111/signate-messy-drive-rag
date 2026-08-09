"""SOT-2511 — deterministic PPTX Gantt week-grid extraction."""
from __future__ import annotations

import io
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches

from src.rag.corpus import FileRef
from src.rag.extract.office import (
    WeekCell,
    _slide_display_pages,
    extract_gantt_week_ranges,
    extract_pptx,
    week_range_for_span,
)


def test_week_span_uses_tolerant_half_open_overlap_at_boundaries() -> None:
    cells = tuple(WeekCell(week, (week - 1) * 100, week * 100) for week in range(1, 5))
    # Exact boundary: [100, 300) means W2-W3, never W1 or W4.
    assert week_range_for_span(100, 200, cells) == (2, 3)
    # SOT-2546: a sub-tolerance sliver past a boundary does NOT over-read the extra week, symmetrically
    # at both ends (grid span 400 // 500 → tolerance 1; a 1-unit poke is ignored, not counted).
    assert week_range_for_span(99, 201, cells) == (2, 3)
    assert week_range_for_span(101, 200, cells) == (2, 3)
    # A larger overlap (well beyond tolerance) is still a genuine membership.
    assert week_range_for_span(50, 250, cells) == (1, 3)
    assert week_range_for_span(400, 10, cells) is None


def _synthetic_gantt() -> Presentation:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for offset, week in enumerate(range(1, 5)):
        header = slide.shapes.add_textbox(
            Inches(4 + offset), Inches(1), Inches(1), Inches(0.4))
        header.text = f"W{week}"
    activity = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(3.5), Inches(0.8))
    activity.text = "モデル改善\n説明性分析"
    # Starts exactly at W2 and ends exactly at W4's left edge => W2-W3.
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(5), Inches(2.4), Inches(7), Inches(2.4))
    return prs


def test_native_shape_geometry_maps_activity_to_week_range() -> None:
    records = extract_gantt_week_ranges(_synthetic_gantt())
    assert records == [{
        "slide": 1,
        "activity": "モデル改善 説明性分析",
        "status": "resolved",
        "start_week": 2,
        "end_week": 3,
        "bar_left": Inches(5),
        "bar_width": Inches(2),
    }]


def test_gantt_annotation_is_prepended_before_flattened_slide_text(tmp_path) -> None:
    path = tmp_path / "schedule.pptx"
    _synthetic_gantt().save(path)
    ref = FileRef(path=path, project="", category="proposal", rel=path.name,
                  name=path.name, ext="pptx")
    text = extract_pptx(ref, None)
    assert text.startswith("【ガント週グリッド:決定論】")
    assert "モデル改善 説明性分析: 第2週目から第3週目" in text
    assert "半開区間" in text


def test_slide_display_pages_reads_footer_and_excludes_cover_and_decoys() -> None:
    """SOT-2546 idx84: the printed ``N / TOTAL`` footer is the page number, not the physical index."""
    prs = Presentation()
    # Slide 1 is an unnumbered cover (no footer).
    prs.slides.add_slide(prs.slide_layouts[6])
    for physical, printed in ((2, 1), (3, 2), (4, 3)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(1), Inches(0.3))
        footer.text = f"{printed} / 3"
    # A stray in-body list numbering must not be mistaken for a page footer (wrong denominator).
    decoy = prs.slides[3].shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.3))
    decoy.text = "1 / 7"
    assert _slide_display_pages(prs) == {2: 1, 3: 2, 4: 3}


def test_extract_pptx_labels_slide_with_printed_page_number() -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    body = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    body.text = "モデル毎のF1スコア ランキング"
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(1), Inches(0.3))
    footer.text = "1 / 2"
    ref = FileRef(path=Path("r.pptx"), project="", category="report", rel="r.pptx",
                  name="r.pptx", ext="pptx")
    out = io.BytesIO()
    prs.save(out)
    text = extract_pptx(ref, out.getvalue())
    assert "[スライド2（文書に記載のページ番号: 1）]" in text
    assert "[スライド1]" in text  # cover keeps the physical label (no printed footer)


def _corpus_pptx(name_needle: str):
    import pytest

    from src.rag.corpus import walk
    ref = next((r for r in walk() if r.ext == "pptx" and name_needle in r.name), None)
    if ref is None:
        pytest.skip(f"corpus pptx {name_needle!r} not present")
    return ref


def test_real_corpus_idx69_pilot_gantt_end_week_is_six() -> None:
    """Gold idx69: the pilot bar pokes a 14151-EMU sliver into W7 but its plotted range is W5–W6."""
    ref = _corpus_pptx("白峰信用リスク評価株式会社_最終報告")
    records = extract_gantt_week_ranges(Presentation(str(ref.path)))
    pilot = next(r for r in records if "パイロット" in str(r["activity"]))
    assert (pilot["start_week"], pilot["end_week"]) == (5, 6)


def test_real_corpus_idx84_f1_ranking_slide_prints_page_five() -> None:
    """Gold idx84: the model-F1 ranking slide is physically the 6th but its footer page number is 5."""
    ref = _corpus_pptx("株式会社東都人材プラットフォーム_最終報告")
    text = extract_pptx(ref, None)
    f1_line = next(line for line in text.splitlines()
                   if "モデル性能比較" in line or ("F1" in line and "0.474" in line))
    block = text.split(f1_line)[0].splitlines()
    header = next(line for line in reversed(block) if line.startswith("[スライド"))
    assert "文書に記載のページ番号: 5" in header


def test_conflicting_native_bars_are_reported_ambiguous() -> None:
    prs = _synthetic_gantt()
    slide = prs.slides[0]
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(6), Inches(2.4), Inches(7), Inches(2.4))
    records = extract_gantt_week_ranges(prs)
    assert records[0]["status"] == "ambiguous"
    assert records[0]["candidates"] == [
        {"start_week": 2, "end_week": 3},
        {"start_week": 3, "end_week": 3},
    ]

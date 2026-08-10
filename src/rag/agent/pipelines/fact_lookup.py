"""SOT-2612 (Wave B2) — deterministic ``simple_lookup`` (fact_lookup) pipeline (canonical/索引 → 決定論読取 → 整形).

Parent PLAN SOT-2602 (決定論先行パイプラインへの反転). Fifth per-type pipeline of the inverted architecture
(after Wave A1 ``version_diff`` SOT-2605, A2 ``numeric`` SOT-2607, A3 ``enum``/``cross_aggregate`` SOT-2608,
A4 ``chart_read``/``spatial`` SOT-2609). A *data系* fact_lookup — one whose value can be pinned by
**canonical direct-file resolution + a deterministic structural read** (schedule table / report page number) —
is answered **without going through the Gemini investigator loop**. The Stage0 router
(:mod:`src.rag.agent.det_pipeline`) classifies the contract as ``simple_lookup`` and dispatches here; Stage3
(:mod:`src.rag.agent.formatting`) renders the grounded value into gold 書式 (template-first — no LLM naturalize
call: this pipeline returns a gold-shaped scalar string / number).

The two tight recognizers this pipeline owns
--------------------------------------------
* **スケジュール表の順序読取** (:func:`_schedule_phase_task`, idx89 型). 「…スケジュール.xlsx において、フェーズNoN
  にて最後に開始するタスク名は何ですか」→ canonical route (kind=schedule) でスケジュール表を確定 →
  :func:`src.rag.extract.office.normalized_xlsx_rows` の**結合セル前方補完**（グルーピング列＝フェーズNo. を前方補完
  して行→フェーズの関係を復元, SOT-2509/2511/2546 の境界規則）で行を正規化 → フェーズ列＝N の行群から、開始日列の
  **一意の**最遅（最後）／最早（最初）行を選び、そのタスク名を返す。極値が同着で一意に絞れなければ棄権（安全側）。
* **報告書のページ番号読取** (:func:`_report_metric_page`, idx84 型). 「…最終報告書で…モデル毎のF1スコアがランキング
  形式で記載されているページ数」→ report pptx を確定 → :func:`src.rag.extract.office._slide_display_pages` の
  **ページ採番境界規則**（表紙を除いた「N / 総数」フッタの印字ページ番号；物理スライド番号ではない, SOT-2549 の
  「競合記載は確定版=印字ページを優先」）でスライド→印字ページを解決 → 質問が名指す指標（F1/accuracy/…）の**値の
  記載密度**が最大の一意スライドを特定し、その印字ページ番号を返す。密度で一意に絞れなければ棄権。

Precision-first negative guards (安全側フォールバック — 誤答/誤「該当なし」を出さない)
--------------------------------------------------------------------------------
決定論スライスは意図的に狭い。各 recognizer は対象（プロジェクト/ファイル/フェーズ/指標/スライド）を **一意に固定**
できない限り ``None``（⇒ LLM フォールバック）を返す。曖昧・自由記述寄りの fact（役職/定義/自由文の所在など、決定論で
1つに絞れないもの）は最初から None を返し、既存の champion LLM 経路に委ねる（回答数を減らさない・wrong を増やさない）。

Design invariants (shared with the Stage0 router / Stage3 formatting / Wave A1〜A4)
--------------------------------------------------------------------------------
* **No hardcoding.** No idx number, no answer, no corpus fact（プロジェクト名/ファイル名/タスク名/ページ番号）is stored
  here — every value is self-derived from the question via canonical_route + the office extractors.
* **Fail-open, never abstain-by-error.** Any read/parse/lookup failure returns ``None`` (⇒ LLM fallback); the
  pipeline never raises into the answer path and never *reduces* the number of answered questions.
* **Gated OFF with the router.** Only ever reached from the Stage0 router's det-path, gated by
  ``RAG_DET_PIPELINE_ROUTER`` (default OFF); with the flag off the champion serve path is byte-identical.
"""
from __future__ import annotations

import datetime as _dt
import importlib
import re
from typing import Any

from src.rag.agent import det_pipeline as _det_pipeline

# The contract type this pipeline owns (question_contract.SIMPLE_LOOKUP — the base contract the
# ``fact_lookup`` archetype refines to).
CONTRACT_TYPE = "simple_lookup"


# --------------------------------------------------------------------------- small deterministic helpers
def _nfc(text: str) -> str:
    from src.rag.corpus import nfc
    return nfc(text or "")


def _resolve_project(question: str) -> "str | None":
    """The single project the question names, via the SOT-2494 canonical route, or ``None``.

    ``import_module`` (not ``from … import``) so the ``src.rag.tools`` package re-exporting
    ``resolve_project`` as a *function* (shadowing the submodule attribute) does not bind the function —
    the real module is returned from ``sys.modules`` so ``resolve_project`` resolves (and is monkeypatchable)
    at call time (tools-shadow gotcha, shared with Wave A2〜A4).
    """
    cr = importlib.import_module("src.rag.tools.canonical_route")
    try:
        return cr.resolve_project(question, None) or None
    except Exception:  # noqa: BLE001 — resolution failure ⇒ fallback, never breaks the answer path
        return None


def _project_refs(project: "str | None"):
    """Every corpus :class:`FileRef` whose NFC project equals ``project`` (empty when ``project`` is None)."""
    from src.rag.corpus import walk
    if not project:
        return []
    proj = _nfc(project)
    return [r for r in walk() if _nfc(r.project) == proj]


def _resolve_schedule_ref(question: str, project: "str | None"):
    """The project's canonical schedule/WBS spreadsheet (:class:`FileRef`), or ``None`` (ambiguous ⇒ fallback).

    Stage1 file確定: reuse the canonical route (``kind=schedule``) to pin the 計画フォルダ の スケジュール表,
    then map its rel back to the live :class:`FileRef` (which carries ``.path`` for the office reader). A
    non-unique / non-tabular / unresolved schedule ⇒ ``None``.
    """
    cr = importlib.import_module("src.rag.tools.canonical_route")
    try:
        routed = cr.canonical_route(question=question, project=project, kind="schedule")
    except Exception:  # noqa: BLE001
        return None
    records = routed.get("value") if isinstance(routed, dict) else None
    if not isinstance(records, list):
        return None
    rels = [str(rec.get("rel")) for rec in records
            if isinstance(rec, dict) and str(rec.get("ext", "")).lower() in {"xlsx", "xlsm"}]
    # A unique canonical schedule table only — two candidate tables are ambiguous (safe fallback).
    unique = list(dict.fromkeys(rels))
    if len(unique) != 1:
        return None
    target = _nfc(unique[0])
    for ref in _project_refs(project):
        if _nfc(ref.rel) == target:
            return ref
    return None


def _resolve_report_pptx_ref(project: "str | None"):
    """The project's single report-slide deck (:class:`FileRef`, category=report, ext=pptx), or ``None``.

    Report decks have no dedicated canonical-route ``Kind``; resolve them structurally from the corpus
    category tag. Requiring **exactly one** report pptx keeps the target unambiguous — 0 / >1 ⇒ ``None``.
    """
    decks = [r for r in _project_refs(project) if r.category == "report" and r.ext == "pptx"]
    return decks[0] if len(decks) == 1 else None


# --------------------------------------------------------------------------- office reads (monkeypatchable)
def _read_schedule_rows(ref):
    """Return ``(header_values, rows)`` for the first worksheet that looks like a フェーズ/タスク/開始 schedule.

    ``rows`` are :class:`~src.rag.extract.office.NormalizedXlsxRow` after the SOT-2509/2546 grouping-column
    forward fill (so a data row inherits its フェーズNo. from the group header above it). ``header_values`` is
    that sheet's detected header row (the densest of the first rows). ``(None, [])`` when no schedule-shaped
    sheet is present.
    """
    office = importlib.import_module("src.rag.extract.office")
    wb = office._xlsx_from(ref, None)
    for ws in wb.worksheets:
        rows = office.normalized_xlsx_rows(ws)
        if not rows:
            continue
        header = _detect_header_row(rows)
        if header is not None:
            return rows[header].values, rows
    return None, []


def _read_report_slides(ref):
    """Return ``(display_pages, slides)`` for a report pptx.

    ``display_pages`` maps 1-based physical slide index → the printed ``N / TOTAL`` footer page number
    (:func:`office._slide_display_pages`); ``slides`` is a list of ``(slide_index, joined_text)`` covering
    every text frame and table cell on the slide (NFC). ``({}, [])`` on any read failure.
    """
    office = importlib.import_module("src.rag.extract.office")
    from pptx import Presentation
    prs = Presentation(str(ref.path))
    display_pages = office._slide_display_pages(prs)
    slides: "list[tuple[int, str]]" = []
    for si, slide in enumerate(prs.slides, 1):
        parts: "list[str]" = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text or "")
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        slides.append((si, _nfc("\n".join(parts))))
    return display_pages, slides


# --------------------------------------------------------------------------- schedule column mapping
# The grouping column that holds the phase *number* (フェーズNo.) — distinct from フェーズ名 (the phase title).
_PHASE_NO_HEADER_RE = re.compile(r"フェーズ\s*(?:no\.?|number|番号|ｎｏ)", re.I)
_TASK_NAME_HEADER_RE = re.compile(r"タスク名|作業名|項目名|task\s*name", re.I)
_START_HEADER_RE = re.compile(r"開始|着手|start", re.I)
_END_HEADER_RE = re.compile(r"終了|完了|end|finish", re.I)


def _detect_header_row(rows) -> "int | None":
    """Index into ``rows`` of the schedule header row (フェーズNo. + タスク名 + 開始), or ``None``.

    Scans only the first rows (a header lives near the top); the row must expose all three structural
    columns so a free-form sheet fails closed.
    """
    for idx, row in enumerate(rows[:30]):
        cells = [_nfc(str(v)) if v is not None else "" for v in row.values]
        has_phase = any(_PHASE_NO_HEADER_RE.search(c) for c in cells)
        has_task = any(_TASK_NAME_HEADER_RE.search(c) for c in cells)
        has_start = any(_START_HEADER_RE.search(c) for c in cells)
        if has_phase and has_task and has_start:
            return idx
    return None


def _unique_col(header_values, pattern: "re.Pattern[str]", *, exclude: "re.Pattern[str] | None" = None) -> "int | None":
    """The single 0-based column whose header matches ``pattern`` (and not ``exclude``), or ``None``.

    Requiring exactly one match keeps the target column unambiguous — 0 / >1 ⇒ ``None`` (safe fallback).
    """
    hits = []
    for ci, value in enumerate(header_values):
        text = _nfc(str(value)) if value is not None else ""
        if not text:
            continue
        if exclude is not None and exclude.search(text):
            continue
        if pattern.search(text):
            hits.append(ci)
    return hits[0] if len(hits) == 1 else None


def _sort_key(value: Any):
    """A comparable key for a schedule date cell, or ``None`` when it cannot be ordered deterministically.

    ``datetime`` cells order natively; an ISO-ish string is parsed; anything else is undecidable so the
    whole extremum is declined rather than guessed (safe fallback).
    """
    if isinstance(value, (_dt.datetime, _dt.date)):
        return value if isinstance(value, _dt.datetime) else _dt.datetime(value.year, value.month, value.day)
    if isinstance(value, str):
        text = value.strip().replace("/", "-")
        match = re.match(r"^(\d{4})-(\d{1,2})-(\d{1,2})", text)
        if match:
            y, m, d = (int(g) for g in match.groups())
            try:
                return _dt.datetime(y, m, d)
            except ValueError:
                return None
    return None


# --------------------------------------------------------------------------- Recognizer A: schedule phase task
_SCHEDULE_Q_RE = re.compile(r"スケジュール|wbs|工程表|ガント", re.I)
_PHASE_NO_Q_RE = re.compile(r"フェーズ\s*(?:no\.?|number|番号|ｎｏ)?\s*[:：]?\s*(\d+)", re.I)
_LAST_CUE_RE = re.compile(r"最後|最も遅|一番遅|末尾|ラスト")
_FIRST_CUE_RE = re.compile(r"最初|最も早|一番早|先頭|冒頭")
_START_CUE_RE = re.compile(r"開始|着手")
_END_CUE_RE = re.compile(r"終了|完了")
_TASK_NAME_Q_RE = re.compile(r"タスク名")


def _schedule_phase_task(question: str) -> "dict[str, Any] | None":
    """Deterministic 「フェーズNoN で最後/最初に開始/終了するタスク名」 read, or ``None`` (⇒ fallback)."""
    q = _nfc(question)
    if not (_SCHEDULE_Q_RE.search(q) and _TASK_NAME_Q_RE.search(q)):
        return None
    phase_match = _PHASE_NO_Q_RE.search(q)
    if not phase_match:
        return None
    phase_no = int(phase_match.group(1))

    want_last = bool(_LAST_CUE_RE.search(q))
    want_first = bool(_FIRST_CUE_RE.search(q))
    if want_last == want_first:  # neither / both ⇒ ordering ambiguous ⇒ fallback
        return None
    uses_start = bool(_START_CUE_RE.search(q))
    uses_end = bool(_END_CUE_RE.search(q))
    if uses_start == uses_end:  # neither / both anchors ⇒ fallback
        return None

    project = _resolve_project(question)
    ref = _resolve_schedule_ref(question, project)
    if ref is None:
        return None

    header_values, rows = _read_schedule_rows(ref)
    if not header_values or not rows:
        return None
    header_idx = _detect_header_row(rows)
    if header_idx is None:
        return None

    phase_col = _unique_col(header_values, _PHASE_NO_HEADER_RE)
    task_col = _unique_col(header_values, _TASK_NAME_HEADER_RE)
    anchor_re = _START_HEADER_RE if uses_start else _END_HEADER_RE
    # 開始日/終了日 headers both contain 日; exclude the opposite anchor so 開始 never matches a 終了 column.
    exclude_re = _END_HEADER_RE if uses_start else _START_HEADER_RE
    date_col = _unique_col(header_values, anchor_re, exclude=exclude_re)
    if phase_col is None or task_col is None or date_col is None:
        return None

    # Collect (sort_key, task) for every data row in the requested phase group.
    candidates: "list[tuple[Any, str, int]]" = []
    for row in rows[header_idx + 1:]:
        values = row.values
        if phase_col >= len(values) or task_col >= len(values) or date_col >= len(values):
            continue
        phase_cell = values[phase_col]
        if phase_cell is None:
            continue
        phase_text = _nfc(str(phase_cell)).strip()
        digits = re.match(r"^(\d+)", phase_text)
        if not digits or int(digits.group(1)) != phase_no:
            continue
        task = values[task_col]
        key = _sort_key(values[date_col])
        if task is None or key is None or not str(task).strip():
            return None  # a matching row with an unreadable date/task ⇒ cannot order safely ⇒ fallback
        candidates.append((key, str(task).strip(), row.row))
    if not candidates:
        return None

    ordered = sorted(candidates, key=lambda item: item[0])
    winner = ordered[-1] if want_last else ordered[0]
    runner = ordered[-2] if want_last else ordered[1] if len(ordered) > 1 else None
    # Require a strictly-unique extremum — a tie on the ordering key is genuinely ambiguous ⇒ fallback.
    if runner is not None and runner[0] == winner[0]:
        return None

    task_name = winner[1]
    evidence = {
        "file": ref.rel,
        "project": project,
        "sheet_header_row": header_idx + 1,
        "phase_no": phase_no,
        "order": "last" if want_last else "first",
        "anchor": "start" if uses_start else "end",
        "row": winner[2],
        "candidates": len(candidates),
        "route": "canonical_route(schedule)→normalized_xlsx_rows",
    }
    method = {
        "engine": "fact_lookup_det_pipeline",
        "contract": CONTRACT_TYPE,
        "shape": "schedule_phase_ordinal_task",
        "boundary_rules": "merged_cell_forward_fill;unique_extremum",
        "confidence": 1.0,
    }
    return {"value": task_name, "evidence": evidence, "method": method}


# --------------------------------------------------------------------------- Recognizer B: report metric page
_PAGE_Q_RE = re.compile(r"何ページ|ページ数|ページ番号|ページを教え|ページはどこ")
_PERMODEL_RE = re.compile(r"モデル\s*(?:毎|別|ごと)|各\s*モデル")
_RANKING_RE = re.compile(r"ランキング|ランク|順位|比較|一覧|並べ")

# Named evaluation metrics and the literal token to count value-mentions of. A question naming exactly one
# of these (and nothing else) picks that metric; the count token is the metric's bare symbol (F1/accuracy/…).
_METRICS: "tuple[tuple[re.Pattern[str], str], ...]" = (
    (re.compile(r"F\s*1(?:\s*スコア|\s*score)?", re.I), "F1"),
    (re.compile(r"accuracy|正解率|正答率|acc\b", re.I), "accuracy"),
    (re.compile(r"precision|適合率", re.I), "precision"),
    (re.compile(r"recall|再現率", re.I), "recall"),
    (re.compile(r"auc|roc", re.I), "AUC"),
)
_DECIMAL_NEAR = r"[^0-9\n]{0,8}[0-9]*\.[0-9]+"


def _selected_metric(question: str) -> "str | None":
    """The single evaluation metric the question names (its count token), or ``None`` (0 / >1 ⇒ fallback)."""
    hits = [token for pattern, token in _METRICS if pattern.search(question)]
    unique = list(dict.fromkeys(hits))
    return unique[0] if len(unique) == 1 else None


def _report_metric_page(question: str) -> "dict[str, Any] | None":
    """Deterministic 「モデル毎の<指標>がランキング形式で記載されているページ数」 read, or ``None`` (⇒ fallback)."""
    q = _nfc(question)
    if not (_PAGE_Q_RE.search(q) and _PERMODEL_RE.search(q) and _RANKING_RE.search(q)):
        return None
    metric = _selected_metric(q)
    if metric is None:
        return None

    project = _resolve_project(question)
    ref = _resolve_report_pptx_ref(project)
    if ref is None:
        return None

    display_pages, slides = _read_report_slides(ref)
    if not display_pages or not slides:
        return None

    value_re = re.compile(re.escape(metric) + _DECIMAL_NEAR, re.I)
    counts = [(si, len(value_re.findall(text))) for si, text in slides]
    ranked = sorted(counts, key=lambda item: item[1], reverse=True)
    top_si, top_n = ranked[0]
    second_n = ranked[1][1] if len(ranked) > 1 else 0
    # Require a page that *dominates*: at least two metric-value mentions and strictly more than any other
    # slide — a per-model metric listing is dense in that metric; a tie means the location is ambiguous.
    if top_n < 2 or top_n <= second_n:
        return None
    page = display_pages.get(top_si)
    if page is None:  # no printed page number for the winning slide ⇒ decline (never guess a physical index)
        return None

    evidence = {
        "file": ref.rel,
        "project": project,
        "metric": metric,
        "slide": top_si,
        "page": page,
        "metric_value_hits": top_n,
        "route": "report_pptx→slide_display_pages→metric_density",
    }
    method = {
        "engine": "fact_lookup_det_pipeline",
        "contract": CONTRACT_TYPE,
        "shape": "report_metric_ranking_page",
        "boundary_rules": "printed_page_footer;dominant_metric_density",
        "confidence": 1.0,
    }
    # The gold surface for a 「何ページ」 ask is the printed page number as a bare string (Stage3 template-first).
    return {"value": str(page), "evidence": evidence, "method": method}


# --------------------------------------------------------------------------- pipeline
_RECOGNIZERS = (_schedule_phase_task, _report_metric_page)


def pipeline(question: str, *, profile: Any = None) -> "dict[str, Any] | None":
    """Deterministic ``simple_lookup`` answer as a ``{value, evidence, method}`` contract, or ``None``.

    Tries each tight recognizer in order and returns the first grounded contract; ``None`` ⇒ no recognizer
    could pin the fact deterministically ⇒ the router falls back to the LLM loop. Never raises into the
    answer path.
    """
    try:
        q = _nfc(question)
        if not q:
            return None
        for recognizer in _RECOGNIZERS:
            result = recognizer(question)
            if result is not None and result.get("value") is not None:
                return result
        return None
    except Exception:  # noqa: BLE001 — any failure falls back to the LLM loop, never breaks the answer path
        return None


def register(*, replace: bool = True) -> None:
    """Register this pipeline for the ``simple_lookup`` contract (idempotent: ``replace=True`` by default)."""
    _det_pipeline.register(CONTRACT_TYPE, pipeline, replace=replace)


# Self-register on import — importing :mod:`src.rag.agent.pipelines` (the router's lazy bootstrap) wires
# this pipeline into the Stage0 registry. Idempotent so a re-import / test reload never raises.
register(replace=True)

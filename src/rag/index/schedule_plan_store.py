"""SOT-2710 (cycle11) — xlsx スケジュール/プラン派生ファクトストア（自動発火レーンの供給元）.

cycle10 で Sonnet(claude-mcp) LLM 経路に乗っていた xlsx スケジュール/プラン系の質問を **決定論直答**へ
昇格させるための、質問非依存の全数事前計算ストア。既存ストア（visual_store / plan_coverage_store）だけでは
届かない 3 型の派生値を、案件ごとの ``02.計画/スケジュール.xlsx`` と ``00.提案/提案書*.pptx`` から純粋な
決定論抽出（openpyxl + python-pptx のみ、LLM/genai 非使用 ⇒ cost $0）で焼き込む:

* **schedule_rows** — スケジュール xlsx の主シートの行を ``{id, kind, phase_no, phase_name, name,
  start_date, hours}`` へ正規化（フェーズNo/フェーズ名は結合セルのキャリーダウン）。→ フェーズ内で最後に
  開始するタスク名（idx89）。
* **buffer_hours_total** — 種別＝バッファ行の ``工数(h)`` の総和（主シート）。→ バッファ工数の合計（idx90）。
* **gantt_phase_weeks** — 提案書 pptx のスケジュール案スライドのガント図から、各フェーズ行を **塗り潰し
  （fore_color）バー** の重なる週へ割り付け（枠だけの計画レンジバーは除外＝競合曖昧を解消）。→ 指定フェーズ
  の実施週（idx75「モデル構築は第何週」）。

主シート選定は質問非依存の固定規則: タスクID ヘッダを持つシートのうち、シート名に改訂マーカ（``rev`` /
``改訂``）を含むものを最新版として優先し、無ければ先頭シート。gold 値は一切参照しない。

規律: ``RAG_SCHEDULE_PLAN_LOOKUP`` 既定 OFF ⇒ serve レーン／ツールは None（byte-identical）。build 自体は
常に実行可能（べき等）。
"""
from __future__ import annotations

import io
import json
import os
import re
import unicodedata
from pathlib import Path
from typing import Any, Sequence

import openpyxl

from config import settings

from src.rag import corpus
from src.rag.corpus import FileRef, nfc
from src.rag.extract import passwords

SCHEMA = "schedule_plan_store"
SCHEMA_VERSION = 1
_ON = {"1", "true", "yes", "on"}

_LOAD_CACHE: dict[str, list[dict[str, Any]]] = {}

# シート名の改訂マーカ（最新版を主シートに選ぶ）。
_REV_MARK = re.compile(r"rev|改訂|最新|_r\d|（\s*rev", re.IGNORECASE)
# 週ヘッダ（"4週目" / "第4週" / "第4週目"）— セル全体がこの形の時のみカラムと見なす。
_WEEK_HDR = re.compile(r"(?:第)?(\d+)週目?")
# マイルストン節ラベル（ガントのフェーズ行はこの手前まで）。
_MS_LABEL = re.compile(r"^マイルスト|^ＭＳ|^MS\b")


def enabled() -> bool:
    """serve レーン／ツールを有効にするか。既定 OFF（``RAG_SCHEDULE_PLAN_LOOKUP``）⇒ byte-identical。"""
    return os.getenv("RAG_SCHEDULE_PLAN_LOOKUP", "0").strip().lower() in _ON


def default_out_path() -> Path:
    return settings.ARTIFACTS_DIR / "schedule_plan_store.jsonl"


def default_report_path() -> Path:
    return settings.ARTIFACTS_DIR / "schedule_plan_store_build_report.json"


# --------------------------------------------------------------------------- helpers
def _n(text: Any) -> str:
    return re.sub(r"\s+", "", unicodedata.normalize("NFKC", str(text or "")))


def _cells(row) -> list[str]:
    return [nfc(str(c)) if c is not None else "" for c in row]


def _col(cells: list[str], *keywords: str) -> "int | None":
    for j, c in enumerate(cells):
        cj = _n(c)
        if any(_n(k) in cj for k in keywords):
            return j
    return None


def _to_hours(value: Any) -> "float | None":
    """工数セル（数値 or '2' or '2h' or '2時間'）を時間の float へ。空/非数は None。"""
    if value is None or value == "":
        return None
    s = _n(value)
    m = re.search(r"-?\d+(?:\.\d+)?", s)
    if not m:
        return None
    try:
        return float(m.group(0))
    except ValueError:
        return None


def _iso_date(value: Any) -> str:
    """開始日セル（datetime or 文字列）を YYYY-MM-DD の比較可能文字列へ（不明は空）。"""
    if value is None:
        return ""
    # openpyxl datetime → isoformat prefix.
    iso = getattr(value, "isoformat", None)
    if callable(iso):
        try:
            return iso()[:10]
        except Exception:  # noqa: BLE001
            pass
    s = nfc(str(value))
    m = re.search(r"(\d{4})[-/年](\d{1,2})[-/月](\d{1,2})", s)
    if m:
        return f"{int(m.group(1)):04d}-{int(m.group(2)):02d}-{int(m.group(3)):02d}"
    return s.strip()[:10]


def _load_workbook(ref: FileRef):
    """暗号化 xlsx は passwords ヘルパで復号してから開く（かえで等）。"""
    if passwords.is_encrypted(ref.path):
        data = passwords.resolve(ref)
        if not data:
            raise ValueError(f"decrypt failed: {ref.rel}")
        return openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    return openpyxl.load_workbook(str(ref.path), data_only=True, read_only=True)


def _schedule_ref(project: str, refs: Sequence[FileRef]) -> "FileRef | None":
    for r in refs:
        if r.project == project and r.category == "plan" and r.ext == "xlsx" \
                and nfc(r.name).startswith("スケジュール") and not nfc(r.name).startswith("~$"):
            return r
    return None


def _proposal_refs(project: str, refs: Sequence[FileRef]) -> list[FileRef]:
    return [r for r in refs if r.project == project and r.category == "proposal" and r.ext == "pptx"]


# --------------------------------------------------------------------------- xlsx: schedule rows
def _parse_sheet_rows(rows: list[list[Any]]) -> "list[dict[str, Any]] | None":
    """タスクID ヘッダを持つ表を {id, kind, phase_no, phase_name, name, start_date, hours} 列へ正規化。

    フェーズNo/フェーズ名は結合セル起因の空欄を直近上方の値でキャリーダウンする。ヘッダが見つからねば None。
    """
    hi = None
    header: list[str] = []
    for i, row in enumerate(rows):
        cells = _cells(row)
        if "タスクID" in "|".join(_n(c) for c in cells):
            below = [_cells(r2) for r2 in rows[i + 1: i + 5]]
            if any(re.match(r"^[A-Z]{1,3}\d{1,3}$", c.strip()) for r2 in below for c in r2):
                hi, header = i, cells
                break
    if hi is None:
        return None
    ci_id = _col(header, "タスクID")
    ci_kind = _col(header, "種別")
    ci_phase_no = _col(header, "フェーズNo", "フェーズ番号")
    ci_phase_nm = _col(header, "フェーズ名", "フェーズ")
    if ci_phase_nm == ci_phase_no:
        ci_phase_nm = _col(header, "フェーズ名")
    ci_name = _col(header, "タスク名")
    ci_start = _col(header, "開始日", "開始")
    ci_hours = _col(header, "工数")

    def _at(c2: list[Any], ci: "int | None") -> Any:
        return c2[ci] if ci is not None and ci < len(c2) else None

    out: list[dict[str, Any]] = []
    phase_no = ""
    phase_nm = ""
    for r2 in rows[hi + 1:]:
        c2 = list(r2)
        cs = _cells(c2)
        rid = _n(_at(cs, ci_id))
        if not re.match(r"^[A-Z]{1,3}\d{1,3}$", rid):
            continue
        pno = _n(_at(cs, ci_phase_no))
        if pno:
            phase_no = pno
        pnm = nfc(str(_at(cs, ci_phase_nm) or "")).strip()
        if pnm:
            phase_nm = pnm
        out.append({
            "id": rid,
            "kind": nfc(str(_at(cs, ci_kind) or "")).strip(),
            "phase_no": phase_no,
            "phase_name": phase_nm,
            "name": nfc(str(_at(cs, ci_name) or "")).strip(),
            "start_date": _iso_date(_at(c2, ci_start)),
            "hours": _to_hours(_at(c2, ci_hours)),
        })
    return out or None


def _parse_schedule(ref: FileRef) -> dict[str, Any]:
    """スケジュール xlsx の主シート（改訂優先→先頭）の行を返す。"""
    wb = _load_workbook(ref)
    candidates: list[tuple[str, list[dict[str, Any]]]] = []
    for ws in wb.worksheets:
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
        parsed = _parse_sheet_rows(rows)
        if parsed:
            candidates.append((ws.title, parsed))
    if not candidates:
        return {"sheet": None, "schedule_rows": [], "buffer_hours_total": None}
    # 主シート: 改訂マーカを含むシートを最新版として優先、無ければ先頭。
    primary = next((c for c in candidates if _REV_MARK.search(c[0])), candidates[0])
    sheet, rows = primary
    buffer_rows = [r for r in rows if "バッファ" in r.get("kind", "") and r.get("hours") is not None]
    buffer_total = round(sum(r["hours"] for r in buffer_rows), 6) if buffer_rows else None
    return {"sheet": sheet, "schedule_rows": rows, "buffer_hours_total": buffer_total}


# --------------------------------------------------------------------------- pptx: gantt phase weeks
def _emu_in(value: Any) -> "float | None":
    try:
        return float(value) / 914400.0
    except (TypeError, ValueError):
        return None


def _shape_fill_rgb(shape) -> "str | None":
    """図形の塗り潰し前景色 RGB（solid のみ）。枠だけ/塗り無しは None。"""
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        rgb = fill.fore_color.rgb
        return str(rgb) if rgb is not None else None
    except Exception:  # noqa: BLE001
        return None


def _gantt_phase_weeks(ref: FileRef) -> dict[str, int]:
    """提案書 pptx のガント図から phase_label→week を、**塗り潰しバー**の重なる週で決定論割付。

    週ヘッダ（"N週目"）テキストボックスの x 範囲でカラムを定め、左端のフェーズ行ラベルの y に、塗り潰し
    （fore_color 有り）の実行バー中心 x が入る週を割り当てる。枠だけの計画レンジバーは fore_color 無しで除外
    されるので、SOT-2704 のガント抽出が競合曖昧だったフェーズ（例: モデル構築）も一意に解ける。ちょうど 1 本
    の塗り潰しバーが 1 週に載るフェーズのみ確定（複数週/複数バーは曖昧として不記載）。
    """
    try:
        from pptx import Presentation
    except Exception:  # noqa: BLE001
        return {}
    try:
        prs = Presentation(str(ref.path))
    except Exception:  # noqa: BLE001
        return {}
    result: dict[str, int] = {}
    for slide in prs.slides:
        weeks: list[tuple[int, float, float]] = []
        labels: list[tuple[str, float]] = []
        header_y: "float | None" = None
        for sh in slide.shapes:
            left = _emu_in(getattr(sh, "left", None))
            top = _emu_in(getattr(sh, "top", None))
            width = _emu_in(getattr(sh, "width", None))
            if left is None or top is None or width is None:
                continue
            has_text = bool(getattr(sh, "has_text_frame", False)) and sh.text_frame is not None
            text = sh.text.strip() if has_text else ""
            m = _WEEK_HDR.fullmatch(_n(text)) if text else None
            if m and left > 2.0:
                weeks.append((int(m.group(1)), left, left + width))
                header_y = top if header_y is None else min(header_y, top)
        if len(weeks) < 2:
            continue
        weeks.sort(key=lambda w: w[1])
        ms_y: "float | None" = None
        for sh in slide.shapes:
            left = _emu_in(getattr(sh, "left", None))
            top = _emu_in(getattr(sh, "top", None))
            if left is None or top is None:
                continue
            has_text = bool(getattr(sh, "has_text_frame", False)) and sh.text_frame is not None
            text = sh.text.strip() if has_text else ""
            if text and left < weeks[0][1] and top > (header_y or 0):
                if _MS_LABEL.match(_n(text)):
                    ms_y = top if ms_y is None else min(ms_y, top)
        for sh in slide.shapes:
            left = _emu_in(getattr(sh, "left", None))
            top = _emu_in(getattr(sh, "top", None))
            width = _emu_in(getattr(sh, "width", None))
            if left is None or top is None or width is None:
                continue
            has_text = bool(getattr(sh, "has_text_frame", False)) and sh.text_frame is not None
            text = sh.text.strip() if has_text else ""
            # フェーズ行ラベル: 左端列・週ヘッダより下・マイルストン節より上・週ヘッダ自身でない。
            if not text or left >= weeks[0][1]:
                continue
            if header_y is not None and top <= header_y + 0.05:
                continue
            if ms_y is not None and top >= ms_y - 0.05:
                continue
            nt = _n(text)
            if _WEEK_HDR.fullmatch(nt) or nt in ("フェーズ", "マイルストン", "マイルストーン", ""):
                continue
            label = re.split(r"[\n\r/／]", text)[0].strip()
            if label:
                labels.append((label, top))
        if not labels:
            continue
        # 塗り潰しバー（テキスト無し・fore_color 有り・週カラム域）→ (phase,week)。
        for sh in slide.shapes:
            left = _emu_in(getattr(sh, "left", None))
            top = _emu_in(getattr(sh, "top", None))
            width = _emu_in(getattr(sh, "width", None))
            if left is None or top is None or width is None:
                continue
            has_text = bool(getattr(sh, "has_text_frame", False)) and sh.text_frame is not None
            if has_text and sh.text.strip():
                continue
            if _shape_fill_rgb(sh) is None:
                continue
            if left < weeks[0][1] - 0.5:
                continue
            center = left + width / 2.0
            wk = next((n for n, a, b in weeks if a <= center < b), None)
            if wk is None:
                continue
            phase = min(labels, key=lambda lb: abs(lb[1] - top))
            if abs(phase[1] - top) > 0.4:
                continue
            key = phase[0]
            if key in result and result[key] != wk:
                result[key] = -1  # 競合 → 後段で除外
            elif key not in result:
                result[key] = wk
    return {k: v for k, v in result.items() if v > 0}


# --------------------------------------------------------------------------- build
def build_case(project: str, refs: Sequence[FileRef]) -> "dict[str, Any] | None":
    sched = _schedule_ref(project, refs)
    parsed = {"sheet": None, "schedule_rows": [], "buffer_hours_total": None}
    if sched is not None:
        parsed = _parse_schedule(sched)
    gantt: dict[str, int] = {}
    for pr in _proposal_refs(project, refs):
        try:
            g = _gantt_phase_weeks(pr)
        except Exception:  # noqa: BLE001
            g = {}
        for k, v in g.items():
            gantt.setdefault(k, v)
    if sched is None and not gantt:
        return None
    return {
        "project": project,
        "schedule_file": sched.rel if sched is not None else None,
        "primary_sheet": parsed.get("sheet"),
        "schedule_rows": parsed.get("schedule_rows") or [],
        "buffer_hours_total": parsed.get("buffer_hours_total"),
        "gantt_phase_weeks": gantt,
        "sources": {
            "schedule": sched.rel if sched is not None else None,
            "proposals": [r.rel for r in _proposal_refs(project, refs)],
        },
    }


def build(refs: Sequence[FileRef] | None = None, *, out: Path | None = None,
          write_report: bool = True) -> dict[str, Any]:
    """全案件のスケジュール/プラン派生ストアを構築し JSONL へ書き出す（べき等・LLM 非使用）。"""
    all_refs = list(refs) if refs is not None else corpus.walk()
    projects = sorted({r.project for r in all_refs
                       if r.project and r.project not in ("社内管理", "案件横断", "")})
    records: list[dict[str, Any]] = []
    skipped: list[str] = []
    for project in projects:
        try:
            rec = build_case(project, all_refs)
        except Exception as exc:  # noqa: BLE001 — 1 案件の失敗が全体を壊さない
            skipped.append(f"{project}: {type(exc).__name__}: {exc}")
            continue
        if rec is None:
            skipped.append(f"{project}: no schedule/proposal")
            continue
        records.append(rec)

    out_path = out or default_out_path()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(json.dumps({"schema": SCHEMA, "version": SCHEMA_VERSION}, ensure_ascii=False) + "\n")
        for rec in records:
            fh.write(json.dumps(rec, ensure_ascii=False) + "\n")
    _LOAD_CACHE.pop(str(out_path), None)

    report = {
        "cases": len(records),
        "skipped": skipped,
        "buffer_totals": {r["project"]: r["buffer_hours_total"] for r in records},
        "gantt_sizes": {r["project"]: len(r["gantt_phase_weeks"]) for r in records},
        "row_counts": {r["project"]: len(r["schedule_rows"]) for r in records},
    }
    if write_report:
        with open(default_report_path(), "w", encoding="utf-8") as fh:
            json.dump(report, fh, ensure_ascii=False, indent=2)
    return {"records": len(records), "report": report}


def load(path: Path | None = None) -> list[dict[str, Any]]:
    """スケジュール/プランストアを読み込む（memoized）。欠損/スキーマ不一致 ⇒ ``[]``（回帰ゼロ）。"""
    out = path or default_out_path()
    key = str(out)
    if key in _LOAD_CACHE:
        return _LOAD_CACHE[key]
    rows: list[dict[str, Any]] = []
    try:
        with open(out, encoding="utf-8") as fh:
            header = fh.readline()
            meta = json.loads(header) if header.strip() else {}
            if isinstance(meta, dict) and meta.get("schema") == SCHEMA \
                    and meta.get("version") == SCHEMA_VERSION:
                for line in fh:
                    line = line.strip()
                    if line:
                        rows.append(json.loads(line))
    except Exception:
        rows = []
    _LOAD_CACHE[key] = rows
    return rows


def case_record(project: str, *, path: Path | None = None) -> "dict[str, Any] | None":
    if not project:
        return None
    for rec in load(path):
        if rec.get("project") == project:
            return rec
    return None


if __name__ == "__main__":
    summary = build()
    print(json.dumps(summary["report"], ensure_ascii=False, indent=2))

"""pptx 金額提示ページ事実ストア（SOT-2705 / cycle10 事前計算事実層 追補）.

Sonnet gold100 cycle10 の abstain idx59 は「案件の金額提示がまとまっているのは何ページか」型で、
現行 heading_page_store が docx/pdf のみ・pptx 非対応、かつ「金額提示」という意味クラスの頁事実が
どのストアにも無いため heading_page_lookup が空振りし、予算切れで棄権していた
(``docs/ai/linear/SOT-2705.md`` / <SOT-2701> cycle10 分析で機械確認済み):

* idx59 — 京ソ(京橋信用ソリューションズ株式会社)の ``00.提案/提案書_final.pptx`` スライド13『8. 費用見積』:
  契約金額(税抜/税込)・支払条件(着手金/検収金)の価格提示。¥ 金額トークンがこのスライドに集中(9個・他
  スライドは 0〜1)し、フッターに可視頁番号テキスト『13』を持つ(全 18 スライドに可視頁番号あり)。gold=13ページ。

このストアは build 時に **全案件×全 pptx×全スライド**で
``{slide_index, title, money_token_count, has_pricing_table, visible_page_number}`` を質問非依存で
全数記録する(網羅計算のみ・LLM フリー)。どのスライドが「金額提示ページ」かの一意判定と回答整形は
serve 側の :mod:`src.rag.agent.pptx_money_page_lane`(決定論レーン)が担う。

Design invariants(sibling の heading_page_store / visual_store と同一)
--------------------------------------------------------------------
* **Opt-in at serve time.** :func:`enabled` (``RAG_PPTX_MONEY_PAGE``) が runtime 参照のみを gate する。
  default OFF ⇒ champion serve path は byte-identical。読み出し/配線は :mod:`src.rag.agent.pptx_money_page_lane`。
* **Build は質問非依存・LLM フリー・追加的.** 見る材料は既存の python-pptx 抽出(暗号化 pptx は
  ``passwords.resolve`` で復号したバイト列)だけ。genai 呼び出しはしない。読めない pptx は 1 件スキップし
  build は継続(fail-open)。gold も質問も参照しない。
* **No hard-coded answers.** ページ番号はスライドのフッタ可視番号テキスト(無ければ物理 1-based)から導出。
  idx 番号も gold ページ番号も埋め込まない。
* **Fail-open.** artifact 欠落・解析不能・復号不能はすべて空へフォールバック(回帰ゼロ)。
"""
from __future__ import annotations

import io
import json
import os
import re
import unicodedata
from pathlib import Path
from typing import Any, Mapping, Sequence

from config import settings
from src.rag import corpus
from src.rag.corpus import FileRef, nfc, walk
from src.rag.extract import passwords

SCHEMA = "pptx-money-page-store"
SCHEMA_VERSION = 1

_ON = {"1", "true", "yes", "on"}

# pptx のみ。~$ 一時ファイルは除外。
_PPTX_EXTS = {"pptx"}

# 金額トークン: 円貨表記(¥/￥ + 数字)または 数字 + 円系単位。密度の一意 argmax が「金額提示ページ」の主シグナル。
_MONEY_TOKEN = re.compile(r"(?:¥|￥)\s*[\d,]+|[\d,]+\s*(?:円|万円|千円|百万円)")
# 価格表/価格キーワード(税抜/税込/金額/費用/見積/価格/単価)。表セルに現れれば価格表とみなす補助シグナル。
_PRICE_KW = re.compile(r"税抜|税込|金額|費用|見積|価格|単価|契約金額|支払")
# フッタの可視頁番号: 図形テキスト全体が 数字のみ("13")または "N / TOTAL"。
_PLAIN_NUM = re.compile(r"^\s*(\d{1,3})\s*$")
_SLASH_NUM = re.compile(r"^\s*(\d{1,3})\s*[/／]\s*(\d{1,3})\s*$")
# フッタ帯: スライド下部(高さ比 ≥ この閾値)にある数字図形だけを可視頁番号候補とする(本文中の箇条番号を除外)。
_FOOTER_BAND = 0.78


def enabled() -> bool:
    """True when the serve path may consult the pptx money-page store (default OFF — opt-in)."""
    return os.getenv("RAG_PPTX_MONEY_PAGE", "0").strip().lower() in _ON


def default_out_path() -> Path:
    return settings.ARTIFACTS_DIR / "pptx_money_page_store.jsonl"


def default_report_path() -> Path:
    return settings.ARTIFACTS_DIR / "pptx_money_page_store_build_report.json"


def norm(text: Any) -> str:
    """NFKC + 空白除去 + 小文字化(タイトル/質問の突合キー)。"""
    return re.sub(r"\s+", "", unicodedata.normalize("NFKC", str(text or ""))).lower()


# --------------------------------------------------------------------------- per-slide extraction helpers
def _shape_texts(slide) -> list[str]:
    """スライドの全テキスト(text frame + 表セル)を上→下・左→右順で返す。"""
    out: list[str] = []
    shapes = sorted(slide.shapes, key=lambda s: (int(getattr(s, "top", 0) or 0),
                                                 int(getattr(s, "left", 0) or 0)))
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            t = (shape.text or "").strip()
            if t:
                out.append(nfc(t))
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = " | ".join((c.text or "").strip() for c in row.cells)
                if cells.strip():
                    out.append(nfc(cells))
    return out


def _slide_title(slide) -> str:
    """スライドタイトル: プレースホルダ title があればそれ、無ければ最初の非フッタ(数字のみでない)テキスト。"""
    try:
        title_ph = slide.shapes.title
        if title_ph is not None and (title_ph.text or "").strip():
            return nfc((title_ph.text or "").strip())
    except Exception:  # noqa: BLE001 — レイアウト依存の title アクセスは失敗しうる
        pass
    for t in _shape_texts(slide):
        first = t.splitlines()[0].strip() if t else ""
        if first and not _PLAIN_NUM.match(first):
            return first
    return ""


def _money_token_count(texts: Sequence[str]) -> int:
    joined = "\n".join(texts)
    return len(_MONEY_TOKEN.findall(joined))


def _has_pricing_table(slide) -> bool:
    """表図形のいずれかに 金額トークン or 価格キーワードを含むセルがあれば True。"""
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        for row in shape.table.rows:
            for cell in row.cells:
                ct = nfc((cell.text or "").strip())
                if ct and (_MONEY_TOKEN.search(ct) or _PRICE_KW.search(ct)):
                    return True
    return False


def _visible_page_number(slide, slide_height: int) -> "int | None":
    """スライド下部フッタ帯の可視頁番号(数字のみ or 'N / TOTAL')を返す。帯に無ければ None(→物理退避)。

    フッタ帯(高さ比 ≥ ``_FOOTER_BAND``)に限定することで、本文中の箇条番号(01/02… や 1/2/3…)を頁番号と
    取り違えない。帯内に複数候補があれば最下部(top 最大)を採る(慣例上フッタは最下部)。
    """
    if not slide_height:
        return None
    best: tuple[int, int] | None = None  # (top, value)
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        top = getattr(shape, "top", None)
        if top is None:
            continue
        if (int(top) / slide_height) < _FOOTER_BAND:
            continue
        raw = nfc((shape.text or "").strip())
        m = _PLAIN_NUM.match(raw)
        val: int | None = None
        if m:
            val = int(m.group(1))
        else:
            ms = _SLASH_NUM.match(raw)
            if ms:
                val = int(ms.group(1))
        if val is None:
            continue
        if best is None or int(top) > best[0]:
            best = (int(top), val)
    return best[1] if best is not None else None


# --------------------------------------------------------------------------- per-file record
def compute_doc(ref: FileRef) -> "dict[str, Any] | None":
    """1 つの pptx から全スライドの金額提示ページ事実レコードを組む(開けなければ None)。"""
    from pptx import Presentation  # lazy — office と同じ依存

    data = None
    try:
        if passwords.is_encrypted(ref.path):
            data = passwords.resolve(ref)
    except Exception:  # noqa: BLE001 — 復号判定/復号の失敗は data=None のまま素の open を試す
        data = None
    try:
        prs = Presentation(io.BytesIO(data) if data else str(ref.path))
    except Exception:  # noqa: BLE001 — 開けない pptx はスキップ(fail-open)
        return None

    slide_height = int(getattr(prs, "slide_height", 0) or 0)
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = _shape_texts(slide)
        slides.append({
            "slide_index": idx,
            "title": _slide_title(slide),
            "money_token_count": _money_token_count(texts),
            "has_pricing_table": _has_pricing_table(slide),
            "visible_page_number": _visible_page_number(slide, slide_height),
        })
    if not slides:
        return None
    return {
        "doc_id": nfc(ref.rel), "project": nfc(ref.project), "category": ref.category,
        "doc_name": nfc(ref.name), "ext": ref.ext, "slide_count": len(slides),
        "slides": slides,
    }


# --------------------------------------------------------------------------- universe / build / io
def _universe(refs: Sequence[FileRef]) -> list[FileRef]:
    out: list[FileRef] = []
    for r in refs:
        if r.ext not in _PPTX_EXTS:
            continue
        name = nfc(r.name)
        if name.startswith("~$"):
            continue
        out.append(r)
    out.sort(key=lambda r: r.rel)
    return out


def write_store(records: Sequence[Mapping[str, Any]], path: Path | None = None) -> dict[str, int]:
    """Atomically write a reproducible JSONL store (schema header + doc_id-sorted rows)."""
    out = Path(path) if path is not None else default_out_path()
    out.parent.mkdir(parents=True, exist_ok=True)
    ordered = sorted(records, key=lambda r: r.get("doc_id", ""))
    tmp = out.with_suffix(out.suffix + ".tmp")
    with tmp.open("w", encoding="utf-8", newline="\n") as handle:
        handle.write(json.dumps({"schema": SCHEMA, "version": SCHEMA_VERSION},
                                ensure_ascii=False, sort_keys=True) + "\n")
        for rec in ordered:
            handle.write(json.dumps(rec, ensure_ascii=False, sort_keys=True) + "\n")
    tmp.replace(out)
    reset_cache()
    return {"docs": len(ordered)}


def build(refs: Sequence[FileRef] | None = None, *, out: Path | None = None,
          write_report: bool = True) -> dict[str, Any]:
    """全案件の pptx を全走査して金額提示ページ事実ストアを焼く(質問非依存・LLM フリー・fail-open)。"""
    refs = list(refs) if refs is not None else list(walk())
    universe = _universe(refs)
    records: list[dict[str, Any]] = []
    skipped: list[str] = []
    for ref in universe:
        try:
            rec = compute_doc(ref)
        except Exception:  # noqa: BLE001 — one bad pptx must not sink the build
            rec = None
        if rec is not None:
            records.append(rec)
        else:
            skipped.append(nfc(ref.rel))
    stats = write_store(records, out)
    report = {
        "schema": SCHEMA, "version": SCHEMA_VERSION,
        "universe": len(universe), "records": len(records), "skipped": skipped,
        "slides": sum(len(r.get("slides", [])) for r in records),
        "money_slides": sum(1 for r in records for s in r.get("slides", [])
                            if s.get("money_token_count", 0) > 0),
    }
    if write_report:
        rp = default_report_path()
        rp.parent.mkdir(parents=True, exist_ok=True)
        rp.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"records": len(records), "docs": stats["docs"], "report": report}


# --------------------------------------------------------------------------- load / read (minimal API)
_LOAD_CACHE: dict[str, list[dict[str, Any]]] = {}


def reset_cache() -> None:
    _LOAD_CACHE.clear()


def load(path: Path | None = None) -> list[dict[str, Any]]:
    """Load the money-page rows (memoized). ``[]`` when absent/unreadable/schema-mismatch (回帰ゼロ)."""
    out = Path(path) if path is not None else default_out_path()
    key = str(out)
    if key in _LOAD_CACHE:
        return _LOAD_CACHE[key]
    rows: list[dict[str, Any]] = []
    try:
        with open(out, encoding="utf-8") as handle:
            header = handle.readline()
            meta = json.loads(header) if header.strip() else {}
            if isinstance(meta, dict) and meta.get("schema") == SCHEMA \
                    and meta.get("version") == SCHEMA_VERSION:
                for line in handle:
                    line = line.strip()
                    if line:
                        rows.append(json.loads(line))
    except Exception:  # noqa: BLE001
        rows = []
    _LOAD_CACHE[key] = rows
    return rows


if __name__ == "__main__":
    summary = build()
    print(json.dumps(summary["report"], ensure_ascii=False, indent=2))

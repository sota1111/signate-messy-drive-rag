"""SOT-2677 — 文書テーブル/全文チャンクストア（cycle6 K1）.

cycle6 abstain の最大クラスタ「文書リーチ欠落」を **質問非依存** の事前計算ストアで解消する。
根本原因: ``read_office`` は docx/pptx/pdf の全文を返すが、investigator が tool 結果を LLM へ渡す境界
(:func:`src.rag.agent.investigator._jsonable`, ``max_str=8000`` / ``max_items=60``) で **先頭 8000 字 /
先頭 60 行に切り詰められる**。長文 docx の末尾表がここで落ちて到達不能になる（idx8/50 東都
データサイエンティスト調査.docx=15888字, ML/DE 給与比較表が offset 10000〜; idx99 みなみ野 糖尿病統計情報.docx
の死亡率都道府県ランキング表が offset 8255〜）。

本ストアは全 docx/pptx/pdf を質問非依存に走査し、2 つの資産を焼き込む:

* **(a) 全表 → 行列 table store**（ファイル/ページ/スライド/表番号の出自 + 直近見出し caption 付き）。
* **(b) 切り詰めなし全文チャンク store**（各文書のフル抽出テキスト + オフセット。serve レーンが keyword
  近傍窓 / 末尾窓を返し、8000 字切り詰めの向こう側へ到達させる）。

規律: 純粋な決定論抽出（python-docx / python-pptx / pdfplumber + 既存 office/plain 抽出のみ、LLM/genai
非使用 ⇒ cost $0）。暗号化 xlsx/docx（かえで等）は既存 passwords ヘルパで復号。``RAG_DOC_REACH_STORE``
既定 OFF ⇒ serve ツールは None（tool 集合/スキーマ/serve path は byte-identical）。build 自体は常に実行可能
（べき等）。
"""
from __future__ import annotations

import io
import json
import os
import re
import unicodedata
from pathlib import Path
from typing import Any, Sequence

from config import settings

from src.rag import corpus
from src.rag.corpus import FileRef, nfc
from src.rag.extract import office, passwords, plain

SCHEMA = "doc_reach_store"
SCHEMA_VERSION = 1
_ON = {"1", "true", "yes", "on"}

_LOAD_CACHE: dict[str, dict[str, Any]] = {}

# 対象拡張子（表と長文を持ちうる文書型）。xlsx は既存の visual/fact ストアが担うので対象外。
_DOC_EXTS = ("docx", "pptx", "pdf")

# 病的な文書がストア/tool 出力を氾濫させないための上限（決定論・質問非依存）。
MAX_TABLES_PER_DOC = 60
MAX_ROWS_PER_TABLE = 120
MAX_COLS_PER_ROW = 40
MAX_CELL_CHARS = 200
MAX_FULLTEXT_CHARS = 200_000


def enabled() -> bool:
    """serve ツールを有効にするか。既定 OFF（``RAG_DOC_REACH_STORE``）⇒ byte-identical。"""
    return os.getenv("RAG_DOC_REACH_STORE", "0").strip().lower() in _ON


def default_out_path() -> Path:
    return settings.ARTIFACTS_DIR / "doc_reach_store.jsonl"


def default_report_path() -> Path:
    return settings.ARTIFACTS_DIR / "doc_reach_store_build_report.json"


# --------------------------------------------------------------------------- helpers
def norm(text: Any) -> str:
    """NFKC 正規化 + 空白除去 + lower（keyword 照合キー）。"""
    return unicodedata.normalize("NFKC", str(text or "")).replace(" ", "").replace("　", "").lower()


def _clip_cell(value: Any) -> str:
    s = "" if value is None else str(value)
    s = s.replace("\r", " ").replace("\n", " ").strip()
    return s[:MAX_CELL_CHARS]


def _clean_rows(rows: list[list[str]]) -> list[list[str]]:
    """空行を落とし、行/列/セルを上限で束ね、全体の行数を上限にする。"""
    out: list[list[str]] = []
    for row in rows:
        cells = [_clip_cell(c) for c in row[:MAX_COLS_PER_ROW]]
        if any(c for c in cells):
            out.append(cells)
        if len(out) >= MAX_ROWS_PER_TABLE:
            break
    return out


def _decrypt_bytes(ref: FileRef) -> bytes | None:
    """暗号化ファイルは復号バイトを返す（失敗/非暗号化は None）。"""
    if not passwords.is_encrypted(ref.path):
        return None
    return passwords.resolve(ref)


# --------------------------------------------------------------------------- docx
def _docx_tables(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    import docx  # lazy — python-docx は既に依存

    d = docx.Document(io.BytesIO(data)) if data else docx.Document(str(ref.path))
    # 表の直前段落テキストを caption として紐付けられるよう、本文段落を順に控える。
    paras = [p.text.strip() for p in d.paragraphs if p.text and p.text.strip()]
    tables: list[dict[str, Any]] = []
    for ti, t in enumerate(d.tables, 1):
        rows = _clean_rows([[c.text for c in r.cells] for r in t.rows])
        if not rows:
            continue
        caption = _nearest_caption(paras, ti, len(d.tables))
        tables.append({"table_id": ti, "locus": f"表{ti}", "caption": caption, "rows": rows})
        if len(tables) >= MAX_TABLES_PER_DOC:
            break
    return tables


def _nearest_caption(paras: list[str], table_index: int, n_tables: int) -> str:
    """best-effort: 表番号に対応する段落位置は python-docx では直接取れないため、見出し語を持つ
    直近段落を候補にする（決定論・近似）。見つからなければ空。"""
    heading_like = [p for p in paras if re.search(r"表|ランキング|比較|一覧|統計|集計|給与|中央値|割合|率", p)]
    if heading_like:
        # 表番号でざっくり割り当て（1-index）。範囲外は最後の見出し。
        idx = min(table_index - 1, len(heading_like) - 1)
        return heading_like[idx][:MAX_CELL_CHARS]
    return ""


# --------------------------------------------------------------------------- pptx
def _pptx_tables(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    from pptx import Presentation  # lazy

    prs = Presentation(io.BytesIO(data) if data else str(ref.path))
    tables: list[dict[str, Any]] = []
    tid = 0
    for si, slide in enumerate(prs.slides, 1):
        # スライド見出し（最初のテキスト枠）を caption 候補に。
        title = ""
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().splitlines()[0][:MAX_CELL_CHARS]
                break
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            tid += 1
            rows = _clean_rows([[c.text for c in r.cells] for r in shape.table.rows])
            if not rows:
                continue
            tables.append({"table_id": tid, "locus": f"スライド{si}", "caption": title, "rows": rows})
            if len(tables) >= MAX_TABLES_PER_DOC:
                return tables
    return tables


# --------------------------------------------------------------------------- pdf
def _pdf_tables(ref: FileRef) -> list[dict[str, Any]]:
    try:
        import pdfplumber  # lazy
    except Exception:  # noqa: BLE001
        return []
    tables: list[dict[str, Any]] = []
    tid = 0
    try:
        with pdfplumber.open(str(ref.path)) as pdf:
            for pi, page in enumerate(pdf.pages, 1):
                for tbl in page.extract_tables() or []:
                    tid += 1
                    rows = _clean_rows([[("" if c is None else c) for c in row] for row in tbl])
                    if not rows:
                        continue
                    tables.append({"table_id": tid, "locus": f"ページ{pi}", "caption": "", "rows": rows})
                    if len(tables) >= MAX_TABLES_PER_DOC:
                        return tables
    except Exception:  # noqa: BLE001 — 1 文書の失敗が全体を壊さない
        return tables
    return tables


# --------------------------------------------------------------------------- per-doc
def _full_text(ref: FileRef, data: bytes | None) -> str:
    """read_office と同じフル抽出（切り詰め前のテキスト）。pdf は plain.extract_pdf。"""
    try:
        if ref.ext == "docx":
            return office.extract_docx(ref, data)
        if ref.ext == "pptx":
            return office.extract_pptx(ref, data)
        if ref.ext == "pdf":
            return plain.extract_pdf(ref)
    except Exception:  # noqa: BLE001
        return ""
    return ""


def build_doc(ref: FileRef) -> dict[str, Any] | None:
    """1 文書を table store + full-text へ正規化。抽出不能なら None。"""
    data = _decrypt_bytes(ref)
    if passwords.is_encrypted(ref.path) and data is None:
        return None  # 復号失敗 — best-effort でスキップ
    try:
        if ref.ext == "docx":
            tables = _docx_tables(ref, data)
        elif ref.ext == "pptx":
            tables = _pptx_tables(ref, data)
        elif ref.ext == "pdf":
            tables = _pdf_tables(ref)
        else:
            return None
    except Exception:  # noqa: BLE001
        tables = []
    full_text = _full_text(ref, data)[:MAX_FULLTEXT_CHARS]
    if not tables and not full_text.strip():
        return None
    return {
        "project": ref.project,
        "rel": ref.rel,
        "name": ref.name,
        "ext": ref.ext,
        "category": ref.category,
        "encrypted": bool(data),
        "n_chars": len(full_text),
        "n_tables": len(tables),
        "tables": tables,
        "full_text": full_text,
    }


# --------------------------------------------------------------------------- build
def build(refs: Sequence[FileRef] | None = None, *, out: Path | None = None,
          write_report: bool = True) -> dict[str, Any]:
    """全 docx/pptx/pdf の table + full-text ストアを構築し JSONL へ書き出す（べき等・LLM 非使用）。"""
    all_refs = list(refs) if refs is not None else corpus.walk()
    docs = [r for r in all_refs if r.ext in _DOC_EXTS and not nfc(r.name).startswith("~$")]
    records: list[dict[str, Any]] = []
    skipped: list[str] = []
    for ref in docs:
        try:
            rec = build_doc(ref)
        except Exception as exc:  # noqa: BLE001 — 1 文書の失敗が全体を壊さない
            skipped.append(f"{ref.rel}: {type(exc).__name__}: {exc}")
            continue
        if rec is None:
            skipped.append(f"{ref.rel}: empty/undecryptable")
            continue
        records.append(rec)

    out_path = out or default_out_path()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(json.dumps({"schema": SCHEMA, "version": SCHEMA_VERSION}, ensure_ascii=False) + "\n")
        for rec in records:
            fh.write(json.dumps(rec, ensure_ascii=False) + "\n")
    _LOAD_CACHE.pop(str(out_path), None)

    report = {
        "docs": len(records),
        "skipped": skipped,
        "total_tables": sum(r["n_tables"] for r in records),
        "long_docs": sorted(
            ({"rel": r["rel"], "n_chars": r["n_chars"], "n_tables": r["n_tables"]}
             for r in records if r["n_chars"] > 8000),
            key=lambda d: -d["n_chars"])[:20],
    }
    if write_report:
        with open(default_report_path(), "w", encoding="utf-8") as fh:
            json.dump(report, fh, ensure_ascii=False, indent=2)
    return {"records": len(records), "report": report}


def load(path: Path | None = None) -> dict[str, Any]:
    """ストアを読み込む（memoized）。欠損/スキーマ不一致 ⇒ ``{"docs": []}``（回帰ゼロ）。

    返り値: ``{"docs": [record, ...], "by_project": {project: [record, ...]}}``。
    """
    out = path or default_out_path()
    key = str(out)
    if key in _LOAD_CACHE:
        return _LOAD_CACHE[key]
    docs: list[dict[str, Any]] = []
    try:
        with open(out, encoding="utf-8") as fh:
            header = fh.readline()
            meta = json.loads(header) if header.strip() else {}
            if isinstance(meta, dict) and meta.get("schema") == SCHEMA \
                    and meta.get("version") == SCHEMA_VERSION:
                for line in fh:
                    line = line.strip()
                    if line:
                        docs.append(json.loads(line))
    except Exception:  # noqa: BLE001
        docs = []
    by_project: dict[str, list[dict[str, Any]]] = {}
    for rec in docs:
        by_project.setdefault(rec.get("project") or "", []).append(rec)
    data = {"docs": docs, "by_project": by_project}
    _LOAD_CACHE[key] = data
    return data


if __name__ == "__main__":
    summary = build()
    print(json.dumps(summary["report"], ensure_ascii=False, indent=2))

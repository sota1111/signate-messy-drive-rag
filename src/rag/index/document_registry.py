"""Document registry (コーパスマニフェスト) + canonical 文書リゾルバ (SOT-2583 / 事前処理 #1).

The gold-100 error analysis (SOT-2568 deep-research comment) traced the dominant loss —
``BUDGET_EXHAUSTED`` (27) / ``NOT_RETRIEVED`` (5) / idx12 (a named file not found → 過剰推論) — to the
agent spending its per-question turn budget **searching for where the evidence lives**. The chunk
index answers "which passage is similar", not "which document is *this*". This module builds a
separate, deterministic **Document Registry** — one record per corpus file — and a **canonical
resolver** that maps a question (and any explicitly-named file) to the corpus document *before* the
answer loop starts, so the agent no longer burns turns re-discovering the corpus.

Design invariants (mirroring the sibling pre-stores — evidence_index #4a, structure_store #5,
canonical_manifest #3):
  * **Full coverage.**  :func:`build` emits a record for *every* walked file (encrypted / binary
    included), so ``python -m src.rag.index`` produces one registry entry per corpus document.
  * **Deterministic & LLM-free build.**  Every field is derived from the path/filename or a cheap
    guarded structural probe (openpyxl / python-pptx). No Gemini call, no network, reproducible bytes.
  * **Resolution chain, not embedding-一発.**  :meth:`Resolver.resolve` walks
    exact-path/filename → exact-alias → version-family → entity-index → lexical/trigram, and only
    *then* (optionally, when an embedder is injected) synopsis embedding — so a confident structural
    hit never depends on a fuzzy vector.
  * **explicit-file hard constraint (idx12).**  When a question names a file explicitly and the
    registry cannot resolve it after an exhaustive manifest scan, :func:`hard_constraint_abstain`
    returns :data:`DOC_NOT_FOUND_REASON` — the serve path abstains instead of guessing over semantic
    retrieval.
  * **Opt-in at serve time.**  :func:`enabled` (``RAG_DOCUMENT_REGISTRY``) gates *runtime
    consultation* only; the artifact is always (re)built at index time as an additive, guarded step.
    Default OFF ⇒ the champion serve path is byte-identical.
  * **No hard-coded answers, redistribution-safe.**  The registry holds only structural document
    metadata with provenance; no corpus-specific answer or secret (passwords never enter it).
"""
from __future__ import annotations

import hashlib
import json
import os
import re
import unicodedata
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Callable, Iterable, Sequence

from config import settings
from src.rag.corpus import FileRef, nfc, walk

SCHEMA = "document-registry"
SCHEMA_VERSION = 1

DOC_NOT_FOUND_REASON = "DOC_NOT_FOUND_AFTER_EXHAUSTIVE_MANIFEST_SCAN"

_ON = {"1", "true", "yes", "on"}

# Document extensions a question may reference explicitly. A ``<name>.<ext>`` token whose ext is here
# triggers the explicit-file hard constraint; other dotted tokens (e.g. ``n_estimators=0.1``) do not.
_DOC_EXTS = ("pptx", "ppt", "xlsx", "xlsm", "xls", "docx", "doc", "csv", "pdf", "png", "jpg",
             "jpeg", "ipynb", "py", "txt", "md", "json")


def enabled() -> bool:
    """True when the serve path should consult the registry (default OFF — opt-in)."""
    return os.getenv("RAG_DOCUMENT_REGISTRY", "0").strip().lower() in _ON


def default_out_path() -> Path:
    return settings.ARTIFACTS_DIR / "document_registry.jsonl"


# --------------------------------------------------------------------------- normalization helpers
def fold(value: str) -> str:
    """NFKC-fold + lower + collapse whitespace — the canonical key for filename/alias matching.

    NFKC unifies full/half-width and compatibility forms; lowering makes ASCII case-insensitive.
    This is the single normalization used both when building alias keys and when matching a query,
    so a query token matches a registry key iff they fold to the same string.
    """
    return " ".join(unicodedata.normalize("NFKC", nfc(str(value))).lower().split())


def content_fingerprint(path: Path | str) -> str:
    """SHA-256 of a file's bytes (staleness key + dedup id), or ``""`` when unreadable."""
    try:
        with open(path, "rb") as handle:
            return hashlib.sha256(handle.read()).hexdigest()
    except Exception:
        return ""


# Date-like tokens embedded in a filename (2025-09-01 / 20250901 / 2025年9月 / 2025_09 …).
_DATE_RE = re.compile(
    r"(?<!\d)(?:20\d{2})(?:[-_/.年]?(?:0?[1-9]|1[0-2]))?(?:[-_/.月]?(?:0?[1-9]|[12]\d|3[01]))?(?!\d)")
# Tokens for the entity index / trigram index: ASCII words (≥2) and CJK/kana runs (≥2).
_ASCII_TOK = re.compile(r"[A-Za-z0-9_]{2,}")
_JP_RUN = re.compile(r"[ぁ-んァ-ヶ一-鿿々ー]{2,}")
# Version tokens on a stem — reuse the diffpair vocabulary (kept in sync via KINDS-style signature).
_VER_TOKEN_RE = re.compile(
    r"[ _\-]+(rev|ver|version|final|old|new|旧版|旧|新|確定|fix|fixed|r|v)(\d*)$", re.IGNORECASE)
_ATTACHED_VER_RE = re.compile(r"([一-龥ぁ-んァ-ヶー])(old|旧版|旧|new|新|final|確定)$", re.IGNORECASE)


def _detected_dates(name: str) -> list[str]:
    seen: list[str] = []
    for m in _DATE_RE.finditer(nfc(name)):
        tok = m.group(0)
        if tok not in seen:
            seen.append(tok)
    return seen


def _tokens(text: str) -> list[str]:
    """Distinct folded matching tokens (ASCII words + CJK runs) for the entity/lexical indexes."""
    t = nfc(text)
    out: list[str] = []
    for m in list(_ASCII_TOK.finditer(t)) + list(_JP_RUN.finditer(t)):
        tok = fold(m.group(0))
        if tok and tok not in out:
            out.append(tok)
    return out


def _trigrams(text: str) -> set[str]:
    """Character 3-grams over the folded (whitespace-stripped) text — fuzzy filename recall."""
    s = fold(text).replace(" ", "")
    if len(s) <= 3:
        return {s} if s else set()
    return {s[i:i + 3] for i in range(len(s) - 2)}


def _version_of(stem: str) -> tuple[str, str] | None:
    """(base, token) when the stem carries a version marker, else None. Mirrors diffpair parsing."""
    s = nfc(stem)
    m = _VER_TOKEN_RE.search(s)
    if m:
        base = s[: m.start()].strip(" _-")
        return (base, m.group(0).strip(" _-")) if base else None
    m = _ATTACHED_VER_RE.search(s)
    if m:
        base = s[: m.start(2)].strip(" _-")
        return (base, m.group(2)) if base else None
    return None


# --------------------------------------------------------------------------- record model
@dataclass
class DocRecord:
    doc_id: str                              # stable unique id = NFC corpus-relative path
    case_id: str                             # project / 案件 (NFC), or ""
    full_path: str                           # NFC rel path
    normalized_filename: str                 # folded filename (with extension)
    filename_aliases: list[str]              # folded name variants (stem, no-version, …)
    extension: str
    title: str                               # human title (stem, NFC)
    title_aliases: list[str]
    named_entities: list[str]                # folded entity tokens (project + filename tokens)
    detected_dates: list[str]
    version_tokens: list[str]
    version_family_id: str                   # shared id across versions of one logical doc, or ""
    predecessor_doc_id: str                  # older version's doc_id, or ""
    sheet_names: list[str]
    table_names: list[str]
    slide_titles: list[str]
    pivot_names: list[str]
    comment_count: int
    conditional_format_count: int
    formatting_features: list[str]
    parser_capabilities: list[str]
    extraction_warnings: list[str]
    synopsis: str
    content_fingerprint: str

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


# --------------------------------------------------------------------------- cheap content probes
_PARSER_CAPS = {
    "xlsx": ["openpyxl", "pandas"], "xlsm": ["openpyxl", "pandas"], "xls": ["pandas"],
    "docx": ["python-docx"], "pptx": ["python-pptx"], "csv": ["pandas"], "pdf": ["pdfplumber"],
    "png": ["vision"], "jpg": ["vision"], "jpeg": ["vision"], "ipynb": ["nbformat"], "py": ["text"],
}


def _probe_xlsx(path: Path, rec: DocRecord) -> None:
    """Fill sheet/table/pivot names + comment / conditional-format counts for an Office spreadsheet."""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=False, data_only=True)
    try:
        rec.sheet_names = [nfc(ws.title) for ws in wb.worksheets]
        comments = 0
        cf = 0
        tables: list[str] = []
        pivots: list[str] = []
        for ws in wb.worksheets:
            try:
                cf += len(list(getattr(ws, "conditional_formatting", []) or []))
            except Exception:
                pass
            for tname in getattr(ws, "tables", {}) or {}:
                tables.append(nfc(str(tname)))
            # openpyxl only surfaces comments when not read_only; count defensively.
            for row in ws.iter_rows():
                for cell in row:
                    if getattr(cell, "comment", None) is not None:
                        comments += 1
        rec.table_names = tables
        rec.pivot_names = pivots
        rec.comment_count = comments
        rec.conditional_format_count = cf
        if cf:
            rec.formatting_features.append("conditional_formatting")
        if comments:
            rec.formatting_features.append("comments")
    finally:
        try:
            wb.close()
        except Exception:
            pass


def _probe_pptx(path: Path, rec: DocRecord) -> None:
    """Fill slide titles for a presentation (title placeholder → first text line fallback)."""
    from pptx import Presentation

    prs = Presentation(str(path))
    titles: list[str] = []
    for slide in prs.slides:
        title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""
        if not title:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip().splitlines()[0]
                    break
        titles.append(nfc(title))
    rec.slide_titles = titles


def _probe_content(ref: FileRef, rec: DocRecord) -> None:
    """Best-effort structural enrichment. Every failure is recorded as an extraction warning."""
    try:
        if ref.ext in {"xlsx", "xlsm"}:
            _probe_xlsx(ref.path, rec)
        elif ref.ext == "pptx":
            _probe_pptx(ref.path, rec)
    except Exception as exc:  # encrypted / malformed / missing dep — degrade, never abort the build
        name = type(exc).__name__
        rec.extraction_warnings.append(
            "encrypted" if "crypt" in name.lower() or "password" in str(exc).lower() else name)


# --------------------------------------------------------------------------- record construction
def make_record(ref: FileRef, *, probe: bool = True) -> DocRecord:
    """Build the full :class:`DocRecord` for one corpus file (always succeeds)."""
    stem = ref.stem
    aliases: list[str] = []

    def _add_alias(value: str) -> None:
        key = fold(value)
        if key and key not in aliases:
            aliases.append(key)

    _add_alias(ref.name)
    _add_alias(stem)
    version = _version_of(stem)
    version_tokens: list[str] = []
    version_family_id = ""
    if version is not None:
        base, tok = version
        version_tokens = [tok]
        _add_alias(base)                       # so "提案書" resolves the versioned "提案書old.pptx"
        _add_alias(f"{base}.{ref.ext}")
        dir_rel = "/".join(nfc(ref.rel).split("/")[:-1])
        version_family_id = fold(f"{ref.project}|{dir_rel}|{base}.{ref.ext}")

    entities = _tokens(f"{ref.project} {stem}")
    rec = DocRecord(
        doc_id=nfc(ref.rel),
        case_id=nfc(ref.project),
        full_path=nfc(ref.rel),
        normalized_filename=fold(ref.name),
        filename_aliases=aliases,
        extension=ref.ext,
        title=nfc(stem),
        title_aliases=[a for a in aliases if a != fold(ref.name)],
        named_entities=entities,
        detected_dates=_detected_dates(ref.name),
        version_tokens=version_tokens,
        version_family_id=version_family_id,
        predecessor_doc_id="",                 # filled corpus-wide in build()
        sheet_names=[],
        table_names=[],
        slide_titles=[],
        pivot_names=[],
        comment_count=0,
        conditional_format_count=0,
        formatting_features=[],
        parser_capabilities=list(_PARSER_CAPS.get(ref.ext, ["text"])),
        extraction_warnings=[],
        synopsis="",
        content_fingerprint=content_fingerprint(ref.path),
    )
    if probe:
        _probe_content(ref, rec)
    rec.synopsis = _synopsis(rec)
    return rec


def _synopsis(rec: DocRecord) -> str:
    """A short deterministic one-liner: title + a few structural anchors."""
    parts = [rec.title]
    if rec.sheet_names:
        parts.append("シート: " + "/".join(rec.sheet_names[:5]))
    if rec.slide_titles:
        titles = [t for t in rec.slide_titles if t][:3]
        if titles:
            parts.append("スライド: " + " / ".join(titles))
    if rec.detected_dates:
        parts.append("日付: " + "/".join(rec.detected_dates[:3]))
    return " | ".join(p for p in parts if p)[:400]


# --------------------------------------------------------------------------- build / persist
def build(refs: list[FileRef] | None = None, *, out: Path | None = None,
          probe: bool = True) -> dict[str, int]:
    """Precompute and persist one :class:`DocRecord` for every corpus file.

    Version predecessors are resolved corpus-wide (via :mod:`src.rag.diffpair` when available, with a
    self-contained fallback) after per-file records exist. Returns a small count summary.
    """
    refs = refs if refs is not None else walk()
    records = [make_record(ref, probe=probe) for ref in refs]
    _link_predecessors(records, refs)
    return write_registry(records, out)


def _link_predecessors(records: list[DocRecord], refs: list[FileRef]) -> None:
    """Set ``predecessor_doc_id`` (older→newer) using diffpair's version-pair rules when available."""
    by_rel = {rec.doc_id: rec for rec in records}
    try:
        from src.rag import diffpair

        for pair in diffpair.find_pairs():
            new = by_rel.get(nfc(pair.new.rel))
            if new is not None:
                new.predecessor_doc_id = nfc(pair.old.rel)
            # share a family id across a resolved pair so a version-family lookup unifies them
            fam = new.version_family_id if new else ""
            old = by_rel.get(nfc(pair.old.rel))
            if old is not None and new is not None:
                fam = fam or new.version_family_id or f"pair|{nfc(pair.old.rel)}"
                old.version_family_id = old.version_family_id or fam
                new.version_family_id = new.version_family_id or fam
    except Exception:
        pass  # additive; a diffpair failure must never abort the registry build


def write_registry(records: Sequence[DocRecord], path: Path | None = None) -> dict[str, int]:
    """Atomically write a reproducible JSONL registry (schema header + doc_id-sorted records)."""
    out = path or default_out_path()
    out.parent.mkdir(parents=True, exist_ok=True)
    ordered = sorted(records, key=lambda r: r.doc_id)
    tmp = out.with_suffix(out.suffix + ".tmp")
    with tmp.open("w", encoding="utf-8", newline="\n") as handle:
        handle.write(json.dumps({"schema": SCHEMA, "version": SCHEMA_VERSION},
                                ensure_ascii=False, sort_keys=True) + "\n")
        for rec in ordered:
            handle.write(json.dumps(rec.to_dict(), ensure_ascii=False, sort_keys=True) + "\n")
    tmp.replace(out)
    reset_cache()
    return {
        "files": len(ordered),
        "versioned": sum(1 for r in ordered if r.version_tokens),
        "with_predecessor": sum(1 for r in ordered if r.predecessor_doc_id),
        "with_sheets": sum(1 for r in ordered if r.sheet_names),
        "with_slides": sum(1 for r in ordered if r.slide_titles),
        "warnings": sum(1 for r in ordered if r.extraction_warnings),
    }


# --------------------------------------------------------------------------- load
_LOAD_CACHE: dict[str, list[dict[str, Any]]] = {}
_RESOLVER_CACHE: dict[str, "Resolver"] = {}


def reset_cache() -> None:
    """Drop every in-process cache (used by the build path and tests)."""
    _LOAD_CACHE.clear()
    _RESOLVER_CACHE.clear()


def load(path: Path | None = None) -> list[dict[str, Any]]:
    """Load the registry rows (memoized). Returns ``[]`` when absent/unreadable/schema-mismatch."""
    out = path or default_out_path()
    key = str(out)
    if key in _LOAD_CACHE:
        return _LOAD_CACHE[key]
    rows: list[dict[str, Any]] = []
    try:
        with open(out, encoding="utf-8") as handle:
            header = handle.readline()
            meta = json.loads(header) if header.strip() else {}
            if isinstance(meta, dict) and meta.get("schema") == SCHEMA \
                    and meta.get("version") == SCHEMA_VERSION:
                for line in handle:
                    line = line.strip()
                    if line:
                        rows.append(json.loads(line))
    except Exception:
        rows = []
    _LOAD_CACHE[key] = rows
    return rows


# --------------------------------------------------------------------------- resolver
@dataclass(frozen=True)
class Resolution:
    """One resolver hit: the document plus the chain tier that produced it."""

    doc_id: str
    tier: str            # exact_filename | alias | version_family | entity | lexical | synopsis
    score: float
    case_id: str


class Resolver:
    """Canonical document resolver over a loaded registry.

    :meth:`resolve` walks the deterministic tiers in priority order and returns the first tier that
    yields any candidate (ranked). ``synopsis`` embedding is an *optional* last tier, used only when
    an ``embedder`` is injected — the resolution is never embedding-一発.
    """

    def __init__(self, rows: Sequence[dict[str, Any]]):
        self.rows = list(rows)
        self._by_id = {r["doc_id"]: r for r in self.rows}
        self._by_filename: dict[str, list[str]] = {}
        self._by_path: dict[str, list[str]] = {}
        self._by_alias: dict[str, list[str]] = {}
        self._by_entity: dict[str, list[str]] = {}
        self._by_family: dict[str, list[str]] = {}
        self._trigram: dict[str, set[str]] = {}
        for r in self.rows:
            did = r["doc_id"]
            self._by_filename.setdefault(r.get("normalized_filename", ""), []).append(did)
            self._by_path.setdefault(fold(r.get("full_path", "")), []).append(did)
            for alias in r.get("filename_aliases", []) or []:
                self._by_alias.setdefault(alias, []).append(did)
            for ent in r.get("named_entities", []) or []:
                self._by_entity.setdefault(ent, []).append(did)
            fam = r.get("version_family_id") or ""
            if fam:
                self._by_family.setdefault(fam, []).append(did)
            blob = " ".join([r.get("normalized_filename", ""), r.get("title", "")]
                            + (r.get("filename_aliases") or []))
            self._trigram[did] = _trigrams(blob)

    # ---- filters -------------------------------------------------------------------------
    def _scope(self, doc_ids: Iterable[str], project: str | None) -> list[str]:
        if not project:
            return list(dict.fromkeys(doc_ids))
        p = fold(project)
        out = [d for d in dict.fromkeys(doc_ids)
               if p in fold(self._by_id[d].get("case_id", "")) or
               fold(self._by_id[d].get("case_id", "")) in p]
        return out

    # ---- tiers ---------------------------------------------------------------------------
    def _exact_filename(self, tokens: Sequence[str], project: str | None) -> list[Resolution]:
        hits: list[str] = []
        for tok in tokens:
            hits += self._by_path.get(tok, [])
            hits += self._by_filename.get(tok, [])
        return self._rank(self._scope(hits, project), "exact_path_filename", 1.0)

    def _alias(self, tokens: Sequence[str], project: str | None) -> list[Resolution]:
        hits: list[str] = []
        for tok in tokens:
            hits += self._by_alias.get(tok, [])
        return self._rank(self._scope(hits, project), "alias", 0.9)

    def _version_family(self, tokens: Sequence[str], project: str | None) -> list[Resolution]:
        # A token that folds to a known base resolves to the whole version family (latest first).
        fams: list[str] = []
        for tok in tokens:
            for did in self._by_alias.get(tok, []):
                fam = self._by_id[did].get("version_family_id") or ""
                if fam:
                    fams.append(fam)
        hits: list[str] = []
        for fam in dict.fromkeys(fams):
            hits += self._by_family.get(fam, [])
        # latest version (has a predecessor, or no version token) ranks first
        scoped = self._scope(hits, project)
        scoped.sort(key=lambda d: (0 if self._by_id[d].get("predecessor_doc_id") else 1, d))
        return self._rank(scoped, "version_family", 0.8, preserve_order=True)

    def _entity(self, tokens: Sequence[str], project: str | None) -> list[Resolution]:
        counts: dict[str, int] = {}
        for tok in tokens:
            for did in self._by_entity.get(tok, []):
                counts[did] = counts.get(did, 0) + 1
        scoped = self._scope(counts.keys(), project)
        scoped.sort(key=lambda d: (-counts[d], d))
        return [Resolution(d, "entity", float(counts[d]), self._by_id[d].get("case_id", ""))
                for d in scoped]

    def _lexical(self, query: str, project: str | None) -> list[Resolution]:
        qtri = _trigrams(query)
        if not qtri:
            return []
        scored: list[tuple[float, str]] = []
        for did in self._scope(self._by_id.keys(), project):
            tri = self._trigram.get(did) or set()
            if not tri:
                continue
            overlap = len(qtri & tri)
            if overlap:
                scored.append((overlap / (len(qtri) ** 0.5), did))
        scored.sort(key=lambda t: (-t[0], t[1]))
        return [Resolution(d, "lexical", s, self._by_id[d].get("case_id", "")) for s, d in scored]

    def _synopsis(self, query: str, project: str | None,
                  embedder: Callable[[list[str]], Sequence[Sequence[float]]]) -> list[Resolution]:
        scoped = self._scope(self._by_id.keys(), project)
        if not scoped:
            return []
        texts = [query] + [self._by_id[d].get("synopsis", "") for d in scoped]
        try:
            vecs = embedder(texts)
        except Exception:
            return []
        q = vecs[0]
        qn = sum(x * x for x in q) ** 0.5 or 1.0
        scored: list[tuple[float, str]] = []
        for did, v in zip(scoped, vecs[1:]):
            vn = sum(x * x for x in v) ** 0.5 or 1.0
            sim = sum(a * b for a, b in zip(q, v)) / (qn * vn)
            scored.append((sim, did))
        scored.sort(key=lambda t: (-t[0], t[1]))
        return [Resolution(d, "synopsis", s, self._by_id[d].get("case_id", "")) for s, d in scored]

    def _distill(self, question: str, project: str | None) -> list[Resolution]:
        """Distilled-synopsis tier (SOT-2658): rank scoped docs by the build-time Gemini distillation.

        Serve-time LLM-free — it only reads ``artifacts/distill_store.jsonl`` (answerable_questions /
        summary / key_facts) baked at build time and lexically matches the question. Consulted only when
        ``RAG_DISTILL_STORE`` is on, so the default serve path never touches it (byte-identical).
        """
        try:
            from src.rag.index import distill_store
        except Exception:  # noqa: BLE001 — additive; a missing store never breaks resolution
            return []
        scoped = set(self._scope(self._by_id.keys(), project))
        if not scoped:
            return []
        out: list[Resolution] = []
        for doc_id in distill_store.candidate_doc_ids(question, allowed=scoped):
            out.append(Resolution(doc_id, "distill", 1.0, self._by_id.get(doc_id, {}).get("case_id", "")))
        return out

    def _rank(self, doc_ids: Sequence[str], tier: str, score: float,
              preserve_order: bool = False) -> list[Resolution]:
        ordered = list(doc_ids) if preserve_order else sorted(dict.fromkeys(doc_ids))
        return [Resolution(d, tier, score, self._by_id[d].get("case_id", "")) for d in ordered]

    # ---- public API ----------------------------------------------------------------------
    def resolve(self, question: str, *, project: str | None = None, limit: int = 3,
                embedder: "Callable[[list[str]], Sequence[Sequence[float]]] | None" = None
                ) -> list[Resolution]:
        """Ranked document candidates for a question, first non-empty deterministic tier wins.

        The explicitly-named file tokens (if any) drive the exact/alias/version tiers; the whole
        question drives the entity/lexical (and optional synopsis) tiers. Returns up to ``limit``.
        """
        alias_tokens: list[str] = []
        for t in explicit_file_tokens(question):
            alias_tokens += filename_match_keys(t)
        for tier in (
            self._exact_filename(alias_tokens, project),
            self._alias(alias_tokens, project),
            self._version_family(alias_tokens, project),
            self._entity(_tokens(question), project),
            self._lexical(question, project),
        ):
            if tier:
                return tier[:limit]
        # SOT-2658: distilled-synopsis tier — an enhanced, LLM-free fallback ahead of embedding synopsis.
        # Guarded by RAG_DISTILL_STORE (default OFF ⇒ this block is skipped and resolve() is byte-identical).
        try:
            from src.rag.index import distill_store
            distill_on = distill_store.enabled()
        except Exception:  # noqa: BLE001
            distill_on = False
        if distill_on:
            dz = self._distill(question, project)
            if dz:
                return dz[:limit]
        if embedder is not None:
            syn = self._synopsis(question, project, embedder)
            if syn:
                return syn[:limit]
        return []

    def resolve_explicit(self, question: str, *, project: str | None = None) -> list[Resolution]:
        """Deterministic exact/alias/version resolution of *explicitly-named* files only.

        This is the tier set the explicit-file hard constraint trusts — no entity/lexical/semantic
        fuzz. Returns ``[]`` when the named file is absent from the corpus (→ hard abstain).
        """
        file_tokens = explicit_file_tokens(question)
        if not file_tokens:
            return []
        alias_tokens: list[str] = []
        for t in file_tokens:
            alias_tokens += filename_match_keys(t)
        for tier in (self._exact_filename(alias_tokens, project),
                     self._alias(alias_tokens, project),
                     self._version_family(alias_tokens, project)):
            if tier:
                return tier
        return []

    def resolves_every_explicit(self, question: str, *, project: str | None = None) -> bool:
        """True only when every explicitly named file resolves through a trusted structural tier."""
        file_tokens = explicit_file_tokens(question)
        if not file_tokens:
            return False
        for token in file_tokens:
            keys = filename_match_keys(token)
            if not any((self._exact_filename(keys, project), self._alias(keys, project),
                        self._version_family(keys, project))):
                return False
        return True


# --------------------------------------------------------------------------- explicit-file detection
_FILE_TOKEN_RE = re.compile(
    r"(?<![A-Za-z0-9_./\\-])(?P<name>[^\s、。「」『』（）()\[\]【】,]+?\.(?P<ext>"
    + "|".join(_DOC_EXTS) + r"))(?![A-Za-z0-9])", re.IGNORECASE)


# ASCII filename run at the tail of a token ("かえで総合病院のtrain.xlsx" → "train.xlsx").
_ASCII_TAIL_RE = re.compile(r"[A-Za-z0-9_.\-]+\.(?:" + "|".join(_DOC_EXTS) + r")$", re.IGNORECASE)


def explicit_file_tokens(question: str) -> list[str]:
    """Every explicitly-named ``<name>.<ext>`` file token in the question (NFC, order-preserving)."""
    q = nfc(question)
    out: list[str] = []
    for m in _FILE_TOKEN_RE.finditer(q):
        tok = m.group("name")
        if tok not in out:
            out.append(tok)
    return out


def filename_match_keys(token: str) -> list[str]:
    """Folded filename keys to try for one explicit token, most-specific first.

    A question names a file with surrounding context — ``かえで総合病院のtrain.xlsx`` or
    ``KSSのfigure_06.png`` — so the raw token carries a company/possessive prefix that no filename
    carries. This peels the prefix (last ``の`` / ``・`` segment, then the trailing ASCII filename run)
    and yields the folded full name, the peeled name, and their extension-less stems, so the resolver
    matches the real corpus filename without over-abstaining on the decorated token.
    """
    keys: list[str] = []

    def _add(value: str) -> None:
        k = fold(value)
        if k and k not in keys:
            keys.append(k)

    tok = nfc(token)
    _add(tok)
    for sep in ("の", "・", ":", "："):
        if sep in tok:
            _add(tok.rsplit(sep, 1)[-1])
    tail = _ASCII_TAIL_RE.search(tok)
    if tail:
        _add(tail.group(0))
    for base in list(keys):
        _add(Path(base).stem)
    return keys


def has_explicit_file_reference(question: str) -> bool:
    return bool(explicit_file_tokens(question))


def rels_for_project(project: str | None, *, extensions: Iterable[str] | None = None,
                     path: Path | None = None) -> list[str]:
    """Corpus-relative paths whose ``case_id`` matches ``project`` (read-only helper; SOT-2616).

    Reads the *already-built* registry (the build is unchanged) and returns the doc_ids in the named
    project, optionally filtered to ``extensions`` (e.g. ``{"xlsx", "xlsm"}`` for the data spreadsheets an
    operand prefill scans). Matching is NFKC-folded and lenient in both directions so a project alias
    resolves. Returns ``[]`` when the registry is absent or the project is empty — the operand prefill's
    project-scoped fallback then simply finds nothing (no degradation).
    """
    if not project:
        return []
    key = fold(project)
    exts = set(extensions) if extensions else None
    out: list[str] = []
    for r in load(path):
        if exts is not None and r.get("extension") not in exts:
            continue
        cid = fold(r.get("case_id", ""))
        if key and cid and (key in cid or cid in key):
            out.append(r["doc_id"])
    return out


def get_resolver(path: Path | None = None) -> Resolver | None:
    """A memoized :class:`Resolver` over the persisted registry, or None when it is unavailable."""
    out = path or default_out_path()
    key = str(out)
    if key not in _RESOLVER_CACHE:
        rows = load(out)
        if not rows:
            return None
        _RESOLVER_CACHE[key] = Resolver(rows)
    return _RESOLVER_CACHE.get(key)


def hard_constraint_abstain(question: str, *, project: str | None = None,
                            path: Path | None = None,
                            scope_resolver: "Callable[[str], str | None] | None" = None
                            ) -> str | None:
    """:data:`DOC_NOT_FOUND_REASON` only when the named file is *certifiably absent*, else ``None``.

    An *explicit* file reference is a hard constraint. When enabled and the deterministic
    exact/alias/version tiers resolve every named file, the registry merely *accelerates* the serve
    path (``None`` — the answer loop proceeds with the resolved document seeded).

    SOT-2597 (較正1): a resolution **miss** is a resolver gap, **not** proof of non-existence — the
    trusted tiers mis-split decorated tokens (idx0/74: ``…から提案書.pptx`` → ``ら提案書.pptx``) and
    cannot see a file referenced by nickname (idx59: ``PP_final.pptx`` vs corpus ``提案書_final.pptx``),
    yet the document is reachable through the legacy exploration path
    (find_files/canonical_route/file_grep). On a miss the serve path must therefore fall back to that
    exploration (``None``) **unless** the registry can positively certify the document is absent.

    *Certificate of absence*: the question resolves to **no** case/project scope in the corpus
    (``scope_resolver`` returns nothing). Legacy exploration itself scopes by project, so with no
    resolvable scope there is nowhere in the population the file could live — only then is
    ``DOC_NOT_FOUND`` a real exhaustive-manifest-scan certificate rather than a speculative
    0-iteration abstain on a reachable document. When ``scope_resolver`` is not supplied (no scope
    signal available), the legacy miss→abstain behavior is kept unchanged.

    Returns ``None`` (no constraint / resolved / reachable scope / disabled) otherwise — so the
    default-OFF serve path is byte-identical.
    """
    if not enabled():
        return None
    if not has_explicit_file_reference(question):
        return None
    resolver = get_resolver(path)
    if resolver is None:
        return None  # no artifact → fail open to the live path (回帰ゼロ)
    if resolver.resolves_every_explicit(question, project=project):
        return None  # resolved → registry only accelerates; never abstain
    if scope_resolver is None:
        return DOC_NOT_FOUND_REASON  # no scope signal → keep the legacy miss→abstain behavior
    try:
        scope = scope_resolver(question)
    except Exception:
        return None  # scope probe failed → fail open to exploration (never speculatively abstain)
    # A resolvable scope means the (possibly mis-named) document is reachable there via legacy
    # exploration → explore. No scope → certified absent → hard-abstain.
    return None if scope else DOC_NOT_FOUND_REASON


# --------------------------------------------------------------------------- measurement
def measure_recall(cases: Sequence[tuple[str, str | None, str]], *, path: Path | None = None,
                   ks: Sequence[int] = (1, 3)) -> dict[str, Any]:
    """Compute canonical_doc_recall@k + explicit_filename_resolution_rate over labelled cases.

    ``cases`` = list of ``(question, project, target_doc_id)``. ``target_doc_id`` is the NFC rel of
    the correct document. Returns per-k recall, the explicit-filename resolution rate (of the cases
    whose question names a file, how many the deterministic explicit tiers resolve to the target),
    and per-case detail for the recorded artifact.
    """
    resolver = get_resolver(path)
    detail: list[dict[str, Any]] = []
    hits = {k: 0 for k in ks}
    explicit_total = explicit_hit = 0
    max_k = max(ks) if ks else 1
    for question, project, target in cases:
        target = nfc(target)
        ranked = resolver.resolve(question, project=project, limit=max_k) if resolver else []
        ranked_ids = [r.doc_id for r in ranked]
        for k in ks:
            if target in ranked_ids[:k]:
                hits[k] += 1
        is_explicit = has_explicit_file_reference(question)
        explicit_resolved = None
        if is_explicit:
            explicit_total += 1
            exp = resolver.resolve_explicit(question, project=project) if resolver else []
            explicit_resolved = bool(exp) and any(r.doc_id == target for r in exp)
            if explicit_resolved:
                explicit_hit += 1
        detail.append({
            "question": question, "project": project, "target": target,
            "ranked": ranked_ids, "top_tier": ranked[0].tier if ranked else None,
            "explicit": is_explicit, "explicit_resolved": explicit_resolved,
        })
    n = len(cases) or 1
    return {
        "n_cases": len(cases),
        "canonical_doc_recall": {f"@{k}": round(hits[k] / n, 4) for k in ks},
        "explicit_filename_cases": explicit_total,
        "explicit_filename_resolution_rate": round(explicit_hit / explicit_total, 4)
        if explicit_total else None,
        "detail": detail,
    }


if __name__ == "__main__":
    print(build())

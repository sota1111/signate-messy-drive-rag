"""SOT-2692 — 計画・スケジュール表カバレッジ（cycle8 C3）.

cycle6/8 で「証拠は到達可能・経路が深いだけ」で予算切れ abstain だった 2 型を、**質問非依存** の事前計算
ストアへ焼き込む（SOT-2680 スケジュールストア資産の延長。既存 ``schedule_store`` には手を触れない）:

* **担当者別 工数派生メトリクス**（idx79）— 案件の ``02.計画/スケジュール.xlsx``（暗号化は
  ``passwords.resolve`` で復号）から「リソース配分」シート（``氏名 → 想定工数``）と「WBS・タスク管理」シート
  （``担当者 → タスク``）を突合し、**データアステル側担当者ごとに 想定工数 ÷ 担当タスク数** を全数計算。
  担当タスク数 ≥ 1 のうち比率が **厳密単独最大** の担当者を argmax として事前確定（gold ハードコード無し）。
* **提案書 週次スケジュール**（idx88）— 案件の ``00.提案/提案書*.pptx`` のガント週グリッド（既存
  :func:`office.extract_gantt_week_ranges` の決定論抽出）から **第N週 → 実施項目（フェーズ名）** を事前確定。
  ``resolved`` のバーだけを採用（``ambiguous`` は回答確定不可として除外）。

規律: 純粋な決定論抽出（openpyxl + python-pptx のみ、LLM/genai 非使用 ⇒ cost $0）。``RAG_PLAN_COVERAGE``
既定 OFF ⇒ serve レーン／ツールは None（byte-identical）。build 自体は常に実行可能（べき等）。
"""
from __future__ import annotations

import io
import json
import os
import re
import unicodedata
from pathlib import Path
from typing import Any, Sequence

import openpyxl

from config import settings

from src.rag import corpus
from src.rag.corpus import FileRef, nfc
from src.rag.extract import office, passwords

SCHEMA = "plan_coverage_store"
SCHEMA_VERSION = 1
_ON = {"1", "true", "yes", "on"}

_LOAD_CACHE: dict[str, list[dict[str, Any]]] = {}

_TASK_RE = re.compile(r"^T\d{1,3}$")


def enabled() -> bool:
    """serve レーン／ツールを有効にするか。既定 OFF（``RAG_PLAN_COVERAGE``）⇒ byte-identical。"""
    return os.getenv("RAG_PLAN_COVERAGE", "0").strip().lower() in _ON


def default_out_path() -> Path:
    return settings.ARTIFACTS_DIR / "plan_coverage_store.jsonl"


def default_report_path() -> Path:
    return settings.ARTIFACTS_DIR / "plan_coverage_store_build_report.json"


# --------------------------------------------------------------------------- helpers
def name_key(value: Any) -> str:
    """氏名比較キー（空白除去 + NFKC）。'池田 直哉' → '池田直哉'。"""
    return unicodedata.normalize("NFKC", str(value or "")).replace(" ", "").replace("　", "")


def _split_owners(value: Any) -> list[str]:
    if value is None:
        return []
    return [n.strip() for n in re.split(r"[、,/／;；]", str(value)) if n.strip()]


def _cells(row) -> list[str]:
    return [nfc(str(c)) if c is not None else "" for c in row]


def _col(cells: list[str], *keywords: str) -> int | None:
    for j, c in enumerate(cells):
        if any(k in c for k in keywords):
            return j
    return None


def _load_workbook(ref: FileRef):
    """暗号化 xlsx（かえで等）は passwords ヘルパで復号してから開く。"""
    if passwords.is_encrypted(ref.path):
        data = passwords.resolve(ref)
        if not data:
            raise ValueError(f"decrypt failed: {ref.rel}")
        return openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    return openpyxl.load_workbook(str(ref.path), data_only=True, read_only=True)


def _schedule_ref(project: str, refs: Sequence[FileRef]) -> FileRef | None:
    """案件の計画スケジュール xlsx（``スケジュール*.xlsx``、一時ファイルは除外）。"""
    for r in refs:
        if r.project == project and r.category == "plan" and r.ext == "xlsx" \
                and nfc(r.name).startswith("スケジュール") and not nfc(r.name).startswith("~$"):
            return r
    return None


def _proposal_refs(project: str, refs: Sequence[FileRef]) -> list[FileRef]:
    return [r for r in refs if r.project == project and r.category == "proposal" and r.ext == "pptx"]


def _to_float(value: Any) -> float | None:
    """'14' / '14.0' / '14時間' 等から工数（時間）を取り出す。数値化できなければ None。"""
    txt = unicodedata.normalize("NFKC", str(value or "")).strip()
    if not txt:
        return None
    m = re.search(r"-?\d+(?:\.\d+)?", txt)
    if not m:
        return None
    try:
        return float(m.group(0))
    except ValueError:
        return None


# --------------------------------------------------------------------------- plan metrics (idx79)
def _task_counts(sched: FileRef) -> dict[str, int]:
    """WBS シートの ``担当者`` を氏名キー別にカウント（Txx 行のみ）。"""
    counts: dict[str, int] = {}
    wb = _load_workbook(sched)
    for ws in wb.worksheets:
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
        header_i = None
        ci_id = ci_owner = None
        for hi, row in enumerate(rows):
            cells = _cells(row)
            if "タスクID" not in "|".join(cells) or "担当" not in "|".join(cells):
                continue
            below = [_cells(r2) for r2 in rows[hi + 1: hi + 4]]
            if not any(_TASK_RE.match(c.strip()) for r2 in below for c in r2):
                continue
            header_i, ci_id, ci_owner = hi, _col(cells, "タスクID"), _col(cells, "担当")
            break
        if header_i is None or ci_id is None or ci_owner is None:
            continue
        for r2 in rows[header_i + 1:]:
            c2 = _cells(r2)
            if ci_id >= len(c2):
                continue
            tid = c2[ci_id].strip()
            if not _TASK_RE.match(tid):
                continue
            if ci_owner < len(c2):
                for owner in _split_owners(c2[ci_owner]):
                    key = name_key(owner)
                    counts[key] = counts.get(key, 0) + 1
        break  # WBS シートは 1 枚だけ数える
    return counts


def _resource_people(sched: FileRef) -> list[dict[str, Any]]:
    """「リソース配分」シートから ``役割 / 氏名 / 想定工数`` を全担当者ぶん抽出（``合計`` 行は除外）。"""
    people: list[dict[str, Any]] = []
    seen: set[str] = set()
    wb = _load_workbook(sched)
    for ws in wb.worksheets:
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
        header_i = ci_name = ci_hours = None
        ci_role: int | None = None
        for hi, row in enumerate(rows):
            cells = _cells(row)
            joined = "|".join(cells)
            if "想定工数" not in joined or "氏名" not in joined:
                continue
            header_i = hi
            ci_name = _col(cells, "氏名")
            ci_hours = _col(cells, "想定工数")
            ci_role = _col(cells, "役割", "役職")
            break
        if header_i is None or ci_name is None or ci_hours is None:
            continue
        for r2 in rows[header_i + 1:]:
            c2 = _cells(r2)
            if ci_name >= len(c2):
                continue
            name = c2[ci_name].strip()
            if not name or name in ("合計", "計", "総計"):
                continue
            hours = _to_float(c2[ci_hours]) if ci_hours < len(c2) else None
            if hours is None:
                continue
            key = name_key(name)
            if key in seen:
                continue
            seen.add(key)
            people.append({
                "name": name,
                "name_key": key,
                "role": c2[ci_role].strip() if ci_role is not None and ci_role < len(c2) else "",
                "hours": hours,
            })
        break  # リソース配分シートは 1 枚
    return people


def _plan_metrics(sched: FileRef) -> dict[str, Any] | None:
    """担当者別 想定工数 ÷ 担当タスク数 を全数計算し、厳密単独最大の担当者を argmax として確定。"""
    people = _resource_people(sched)
    if not people:
        return None
    counts = _task_counts(sched)
    rows: list[dict[str, Any]] = []
    for p in people:
        tc = counts.get(p["name_key"], 0)
        ratio = round(p["hours"] / tc, 2) if tc > 0 else None
        rows.append({**p, "task_count": tc, "hours_per_task": ratio})
    scored = [r for r in rows if r["hours_per_task"] is not None]
    if not scored:
        return None
    scored.sort(key=lambda r: r["hours_per_task"], reverse=True)
    top = scored[0]
    # precision-first: 厳密単独最大の時だけ argmax を確定（同率首位は None ⇒ serve は defer）。
    unique_top = len(scored) == 1 or scored[1]["hours_per_task"] < top["hours_per_task"]
    argmax = {
        "name": top["name"], "role": top["role"], "hours": top["hours"],
        "task_count": top["task_count"], "hours_per_task": top["hours_per_task"],
    } if unique_top else None
    return {
        "source": sched.rel,
        "people": rows,
        "max_hours_per_task": argmax,
    }


# --------------------------------------------------------------------------- weekly schedule (idx88)
def _first_paragraph(shape) -> str:
    """ガント活動ラベルの先頭段落（フェーズ名）。担当者名の段落を含まない。"""
    try:
        for para in shape.text_frame.paragraphs:
            txt = nfc((para.text or "").strip())
            if txt:
                return txt
    except Exception:  # noqa: BLE001
        return ""
    return ""


def _weekly_schedule(refs: Sequence[FileRef]) -> dict[str, Any] | None:
    """提案書 pptx のガント週グリッドから 第N週 → 実施項目（フェーズ名）を事前確定（resolved のみ）。"""
    from pptx import Presentation  # lazy — office と同じ依存

    weeks: dict[str, list[str]] = {}
    items: list[dict[str, Any]] = []
    source: str | None = None
    for ref in refs:
        try:
            prs = Presentation(str(ref.path))
            records = office.extract_gantt_week_ranges(prs)
        except Exception:  # noqa: BLE001 — 抽出失敗は無視（週次は best-effort）
            continue
        resolved = [r for r in records if r.get("status") == "resolved"]
        if not resolved:
            continue
        # ラベル full-text → 先頭段落（フェーズ名）の対応を作る（担当者名を落とす）。
        phase_of: dict[str, str] = {}
        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                joined = " ".join((shape.text or "").split())
                if joined and joined not in phase_of:
                    first = _first_paragraph(shape)
                    if first:
                        phase_of[joined] = first
        source = ref.rel
        for rec in resolved:
            activity = str(rec.get("activity") or "")
            phase = phase_of.get(activity) or activity.split(" ")[0]
            start, end = int(rec["start_week"]), int(rec["end_week"])
            items.append({"phase": phase, "activity": activity,
                          "start_week": start, "end_week": end, "slide": rec.get("slide")})
            for wk in range(start, end + 1):
                bucket = weeks.setdefault(str(wk), [])
                if phase not in bucket:
                    bucket.append(phase)
        break  # 案件あたり最初のガント付き提案書を採用
    if not weeks:
        return None
    return {"source": source, "weeks": weeks, "items": items}


# --------------------------------------------------------------------------- build
def build_case(project: str, refs: Sequence[FileRef]) -> dict[str, Any] | None:
    sched = _schedule_ref(project, refs)
    plan_metrics = None
    if sched is not None:
        try:
            plan_metrics = _plan_metrics(sched)
        except Exception:  # noqa: BLE001 — 1 案件の xlsx 失敗が全体を壊さない
            plan_metrics = None
    weekly = _weekly_schedule(_proposal_refs(project, refs))
    if plan_metrics is None and weekly is None:
        return None
    return {
        "project": project,
        "plan_metrics": plan_metrics,
        "weekly_schedule": weekly,
    }


def build(refs: Sequence[FileRef] | None = None, *, out: Path | None = None,
          write_report: bool = True) -> dict[str, Any]:
    """全案件の計画・スケジュールカバレッジを構築し JSONL へ書き出す（べき等・LLM 非使用）。"""
    all_refs = list(refs) if refs is not None else corpus.walk()
    projects = sorted({r.project for r in all_refs
                       if r.project and r.project not in ("社内管理", "案件横断", "")})
    records: list[dict[str, Any]] = []
    skipped: list[str] = []
    for project in projects:
        try:
            rec = build_case(project, all_refs)
        except Exception as exc:  # noqa: BLE001
            skipped.append(f"{project}: {type(exc).__name__}: {exc}")
            continue
        if rec is None:
            skipped.append(f"{project}: no plan/schedule coverage")
            continue
        records.append(rec)

    out_path = out or default_out_path()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(json.dumps({"schema": SCHEMA, "version": SCHEMA_VERSION}, ensure_ascii=False) + "\n")
        for rec in records:
            fh.write(json.dumps(rec, ensure_ascii=False) + "\n")
    _LOAD_CACHE.pop(str(out_path), None)

    report = {
        "cases": len(records),
        "skipped": skipped,
        "argmax": {rec["project"]: (rec["plan_metrics"] or {}).get("max_hours_per_task")
                   for rec in records if rec.get("plan_metrics")},
        "weeks": {rec["project"]: sorted((rec["weekly_schedule"] or {}).get("weeks", {}).keys())
                  for rec in records if rec.get("weekly_schedule")},
    }
    if write_report:
        with open(default_report_path(), "w", encoding="utf-8") as fh:
            json.dump(report, fh, ensure_ascii=False, indent=2)
    return {"records": len(records), "report": report}


def load(path: Path | None = None) -> list[dict[str, Any]]:
    """計画カバレッジストアを読み込む（memoized）。欠損/スキーマ不一致 ⇒ ``[]``（回帰ゼロ）。"""
    out = path or default_out_path()
    key = str(out)
    if key in _LOAD_CACHE:
        return _LOAD_CACHE[key]
    rows: list[dict[str, Any]] = []
    try:
        with open(out, encoding="utf-8") as fh:
            header = fh.readline()
            meta = json.loads(header) if header.strip() else {}
            if isinstance(meta, dict) and meta.get("schema") == SCHEMA \
                    and meta.get("version") == SCHEMA_VERSION:
                for line in fh:
                    line = line.strip()
                    if line:
                        rows.append(json.loads(line))
    except Exception:
        rows = []
    _LOAD_CACHE[key] = rows
    return rows


def case_record(project: str, *, path: Path | None = None) -> dict[str, Any] | None:
    if not project:
        return None
    for rec in load(path):
        if rec.get("project") == project:
            return rec
    return None


if __name__ == "__main__":
    summary = build()
    print(json.dumps(summary["report"], ensure_ascii=False, indent=2))

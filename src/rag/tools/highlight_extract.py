"""highlight_extract — enumerate highlighted / marker-emphasised words & cells across office/PDF.

Many corpus questions are *書式型* (format questions): "マーカーされている単語をすべて抜き出して",
"黄色ハイライトされているセルの抽出条件", "オレンジにハイライトされている行のタスク名". Answering them
needs the **formatting layer**, not the flattened text — which cell is filled orange, which run carries a
highlighter mark, which PDF span sits under a highlight annotation. The plain extractors surface some of
this as prose (``【ハイライト:…】``), but the agent still has to parse that back out. This tool gives the
agent a single, structured, **colour-filterable** enumeration it can call directly:

    highlight_extract("…/train.xlsx", color="黄")   # every yellow-filled cell, in document order

It reuses the *same* deterministic colour classifiers as the office extractor
(:func:`src.rag.extract.office._excel_color_name` / :func:`~src.rag.extract.office._run_highlight` /
:func:`~src.rag.extract.office._color_name`) so a highlight is named identically whether it is read here
or through ``read_office`` — no second, divergent colour heuristic. **No corpus-specific colour or word is
hard-coded**: the colour vocabulary is the generic Japanese coarse-colour map already in the extractor, and
the caller supplies any colour filter, so the tool is fully portable (移植性の担保).

Supported inputs (dispatched by extension), each returning highlighted items in **reading / document
order**:

======  ==================================================================================
ext      what counts as a highlight
======  ==================================================================================
xlsx     a cell whose solid fill classifies to a colour (``_excel_color_name``)
xlsm     (same as xlsx)
pptx     a run carrying an ``a:highlight`` mark, or a text shape with a solid fill colour
docx     a run carrying ``w:highlight`` or a non-default ``w:shd`` cell/character shading
pdf      text under a Highlight annotation, coloured (non–black/grey) text, and faux-italic
         (matrix-shear) emphasis (reusing :func:`src.rag.tools.pdf_faux_italic.emphasized_words`)
======  ==================================================================================

Encrypted Office files are transparently decrypted through the same passwords/profile path as
``extract_office`` (the password *scheme* is reported, never the secret itself — 受け入れ条件③).

Contract
--------
Returns the unified :mod:`src.rag.tools.contract` shape ``{value, evidence, method}`` where ``value`` is a
list of **per-item contracts** — each highlighted hit is itself ``{value, evidence, method}``:

* ``value``    — the highlighted text (a cell's value, a run's text, a span's text).
* ``evidence`` — ``file`` + the location (``sheet`` & ``cell`` for xlsx, ``slide`` & ``shape`` for pptx,
  ``paragraph`` for docx, ``page`` & ``bbox`` for pdf).
* ``method``   — ``engine`` (xlsx/pptx/docx/pdf), the classified ``color`` (Japanese coarse name or a hex
  when a raw colour can't be named), and a ``kind`` tag (``cell`` / ``run_highlight`` / ``shape_fill`` /
  ``annotation`` / ``colored_text`` / ``faux_italic``).

so the generation layer can treat each hit uniformly. The outer ``evidence`` counts the hits and lists the
distinct colours seen; the outer ``method`` records the colour filter applied.
"""
from __future__ import annotations

import io
import json
import os
import re
from pathlib import Path
from typing import Any

from src.rag.corpus import FileRef, nfc
from src.rag.extract import office as _office
from src.rag.extract import passwords as _passwords
from src.rag.tools import contract
from src.rag.tools.contract import ContractError
from src.rag.tools.extract_tools import resolve_ref
from src.rag.tools import xlsx_embedded_image as _xlsx_embedded_image
from src.rag.tools.pdf_faux_italic import detect_faux_italic
from src.rag.tools.profile import CorpusProfile

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"          # DrawingML (pptx runs)
_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"   # WordprocessingML (docx runs)

# Office files this tool can read a formatting layer from.
_OFFICE_EXTS = {"xlsx", "xlsm", "pptx", "docx"}

# Word ``w:highlight`` values are a fixed keyword palette (ST_HighlightColor), not RGB — map each to the
# same coarse Japanese colour vocabulary the RGB classifier uses so a filter matches across formats.
_WORD_HIGHLIGHT_NAMES = {
    "yellow": "黄", "darkyellow": "黄",
    "green": "緑", "darkgreen": "緑",
    "cyan": "水色", "darkcyan": "水色",
    "blue": "青", "darkblue": "青",
    "red": "赤", "darkred": "赤",
    "magenta": "ピンク", "darkmagenta": "紫",
    "darkgray": "灰", "lightgray": "灰",
}

# Colour-filter synonyms → the coarse Japanese name the classifiers emit. Lets a caller pass an English or
# variant colour word ("yellow", "オレンジ色") and still match. Purely a vocabulary map (no corpus fact).
_COLOR_SYNONYMS = {
    "yellow": "黄", "きいろ": "黄", "黄色": "黄",
    "orange": "オレンジ", "だいだい": "オレンジ",
    "red": "赤", "あか": "赤",
    "green": "緑", "みどり": "緑", "yellowgreen": "黄緑",
    "blue": "青", "あお": "青", "navy": "青",
    "cyan": "水色", "lightblue": "水色", "みずいろ": "水色",
    "purple": "紫", "violet": "紫", "むらさき": "紫",
    "pink": "ピンク", "magenta": "ピンク",
    "gray": "灰", "grey": "灰", "はいいろ": "灰",
}


# A coarse colour and the finer shades a human would still call by that name. Asking for 「青」(blue)
# means any blue-ish fill, but the classifier separates 水色 (light blue) from 青 — so a 青 filter would
# silently drop a 水色 highlight (SOT-2564 idx25, whose 青色ハイライト cells classify as 水色). Only
# consulted when RAG_HIGHLIGHT_EXTRA is on, so the default filter stays an exact-name match (byte-identical).
_COLOR_FAMILY = {
    "青": {"青", "水色", "紺"},
    "水色": {"水色", "青"},
}

_ON = {"1", "true", "yes", "on"}


def _extra_enabled() -> bool:
    """``RAG_HIGHLIGHT_EXTRA`` (default OFF): conditional-format highlights, blue-family colour
    matching, and pivot-ancestor context. OFF ⇒ ``highlight_extract`` is byte-identical (champion
    serve unchanged), mirroring ``RAG_FONT_EMPHASIS`` / ``RAG_XLSX_EMBEDDED_IMAGE`` (SOT-2564)."""
    return os.getenv("RAG_HIGHLIGHT_EXTRA", "0").strip().lower() in _ON


def _norm_color_query(color: str | None) -> str | None:
    """Normalise a caller-supplied colour filter to the coarse name the classifiers emit (or ``None``)."""
    if not color:
        return None
    c = nfc(str(color)).strip()
    key = c.lower()
    if key in _COLOR_SYNONYMS:
        return _COLOR_SYNONYMS[key]
    # strip a trailing 「色」 (e.g. 「オレンジ色」→「オレンジ」) then retry the synonym map
    if c.endswith("色") and len(c) > 1:
        base = c[:-1]
        return _COLOR_SYNONYMS.get(base.lower(), base)
    return c


def _color_matches(item_color: str | None, target: str | None, *, family: bool = False) -> bool:
    """True when no filter is set, or the item's classified colour equals the normalised target.

    With ``family=True`` (RAG_HIGHLIGHT_EXTRA on) a coarse target also matches its finer shades — a
    「青」 filter matches a 水色 highlight — via :data:`_COLOR_FAMILY` (SOT-2564 idx25)."""
    if target is None:
        return True
    if item_color is None:
        return False
    ic = nfc(item_color).strip()
    if ic.lower() == target.strip().lower():
        return True
    if family:
        return ic in _COLOR_FAMILY.get(target.strip(), frozenset())
    return False


def _item(value: Any, *, engine: str, color: str | None, kind: str,
          **evidence: Any) -> dict[str, Any]:
    """Build one per-item contract ``{value, evidence, method}`` for a single highlighted hit."""
    return {
        "value": value,
        "evidence": {k: v for k, v in evidence.items() if v is not None},
        "method": {"engine": engine, "color": color, "kind": kind},
    }


# --------------------------------------------------------------------------- office decryption
def _office_bytes(ref: FileRef, profile: CorpusProfile | None) -> bytes | None:
    """Decrypt an encrypted Office file to bytes (via the profile/passwords path), else ``None``.

    Returns ``None`` for a non-encrypted file (the reader opens the path directly). Raises
    :class:`ContractError` when an encrypted file cannot be decrypted, mirroring ``extract_office``.
    """
    if not _passwords.is_encrypted(ref.path):
        return None
    cached = profile.get_password(ref.rel) if profile is not None else None
    data = _passwords._try(ref.path, cached) if cached is not None else None
    if data is None:
        pw = _passwords.resolve_password(ref)
        if pw is not None:
            data = _passwords._try(ref.path, pw)
            if data is not None and profile is not None:
                profile.set_password(ref.rel, pw)
    if data is None:
        raise ContractError(f"could not decrypt {ref.rel or ref.path}")
    return data


# --------------------------------------------------------------------------- per-format extractors
def _pivot_ancestors(ws, cell) -> dict[str, Any] | None:
    """Enclosing group labels of a highlighted cell in a (possibly hierarchical / forward-filled) pivot.

    For each column left of the cell, the nearest non-empty value at or above the cell's row, keyed by
    that column's header (row 1). Recovers ``{Gender: Male, target: 2, Age: 40-44, Country: Spain}`` for
    a highlighted count cell so a "抽出条件と集計内容" question is answerable (SOT-2564 idx15). Dense
    tables (every left cell filled) resolve on the same row, so this stays cheap; only sparse pivots
    scan upward. ``None`` when no ancestor label exists."""
    from openpyxl.utils import get_column_letter

    ctx: dict[str, Any] = {}
    for c in range(1, cell.column):
        label = None
        for r in range(cell.row, 0, -1):
            v = ws.cell(r, c).value
            if v is not None and str(v).strip() != "":
                label = v
                break
        if label is None:
            continue
        header = ws.cell(1, c).value
        key = str(header).strip() if header not in (None, "") else get_column_letter(c)
        ctx[key] = label
    return ctx or None


def _xlsx_items(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    """Every solid-filled (highlighted) cell, sheet→row→column order, named via ``_excel_color_name``."""
    import openpyxl

    src = io.BytesIO(data) if data else str(ref.path)
    wb = openpyxl.load_workbook(src, data_only=True)
    file = ref.rel or str(ref.path)
    extra = _extra_enabled()
    items: list[dict[str, Any]] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                fill = cell.fill
                if not (fill and fill.patternType):
                    continue
                color = _office._excel_color_name(getattr(fill, "fgColor", None))
                if not color:
                    continue
                value = "" if cell.value is None else str(cell.value)
                items.append(_item(value, engine="xlsx", color=color, kind="cell",
                                   file=file, sheet=ws.title, cell=cell.coordinate,
                                   row=cell.row, column=cell.column,
                                   group=_pivot_ancestors(ws, cell) if extra else None))
    return items


def _cf_condition_text(operator: str | None, formulas: Any) -> str:
    """Human-readable condition for a ``cellIs`` conditional-format rule (e.g. ``セルの値 < -0.99``)."""
    ops = {"lessThan": "<", "lessThanOrEqual": "<=", "greaterThan": ">",
           "greaterThanOrEqual": ">=", "equal": "=", "notEqual": "!=",
           "between": "の範囲", "notBetween": "の範囲外", "containsText": "を含む"}
    vals = ", ".join(str(x) for x in (formulas or []))
    label = ops.get(operator or "", operator or "")
    if operator in ("between", "notBetween"):
        return f"セルの値が {vals} {label}".strip()
    return f"セルの値 {label} {vals}".strip()


def _cf_satisfying_cells(ws, sqref, operator: str | None, formulas: Any) -> list[tuple[str, Any]]:
    """Numeric cells inside ``sqref`` that satisfy a ``cellIs`` rule, as ``(coordinate, value)``."""
    try:
        thr = [float(str(f)) for f in (formulas or [])]
    except (TypeError, ValueError):
        return []
    if not thr:
        return []

    def ok(x: float) -> bool:
        if operator == "lessThan":
            return x < thr[0]
        if operator == "lessThanOrEqual":
            return x <= thr[0]
        if operator == "greaterThan":
            return x > thr[0]
        if operator == "greaterThanOrEqual":
            return x >= thr[0]
        if operator == "equal":
            return x == thr[0]
        if operator == "notEqual":
            return x != thr[0]
        if operator == "between" and len(thr) >= 2:
            return thr[0] <= x <= thr[1]
        if operator == "notBetween" and len(thr) >= 2:
            return not (thr[0] <= x <= thr[1])
        return False

    out: list[tuple[str, Any]] = []
    for cr in sqref.ranges:
        min_col, min_row, max_col, max_row = cr.bounds
        for r in range(min_row, max_row + 1):
            for c in range(min_col, max_col + 1):
                v = ws.cell(r, c).value
                if isinstance(v, bool) or not isinstance(v, (int, float)):
                    continue
                if ok(float(v)):
                    out.append((ws.cell(r, c).coordinate, v))
    return out


def _xlsx_cf_items(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    """Conditional-format highlights — the fills openpyxl exposes ONLY via ``ws.conditional_formatting``.

    A yellow/red highlight applied by a rule (e.g. ``値 < -0.99``) never touches ``cell.fill``, so the
    solid-fill scan in :func:`_xlsx_items` misses it entirely (SOT-2564 idx65 returned zero). Surface
    each ``cellIs`` rule as one highlight item carrying its colour (from the dxf fill) and a
    human-readable condition, plus the numeric cells that satisfy it — so both "黄色ハイライトの条件は?"
    (needs the rule) and "青ハイライトの合計は?" (needs the cells) are answerable. Best-effort: any read
    error yields no CF items (zero regression)."""
    import openpyxl

    src = io.BytesIO(data) if data else str(ref.path)
    try:
        wb = openpyxl.load_workbook(src, data_only=True)
    except Exception:
        return []
    file = ref.rel or str(ref.path)
    items: list[dict[str, Any]] = []
    for ws in wb.worksheets:
        cf = ws.conditional_formatting
        for rng in cf:
            for rule in cf[rng]:
                if getattr(rule, "type", None) != "cellIs":
                    continue
                dxf = getattr(rule, "dxf", None)
                fill = getattr(dxf, "fill", None) if dxf else None
                color = None
                if fill is not None:
                    for attr in ("bgColor", "fgColor"):
                        c = getattr(fill, attr, None)
                        name = _office._excel_color_name(c) if c is not None else None
                        if name:
                            color = name
                            break
                cond = _cf_condition_text(rule.operator, rule.formula)
                sat = _cf_satisfying_cells(ws, rng.sqref, rule.operator, rule.formula)
                items.append(_item(
                    cond, engine="xlsx", color=color, kind="conditional_format",
                    file=file, sheet=ws.title, sqref=str(rng.sqref), operator=rule.operator,
                    formula=[str(f) for f in (rule.formula or [])], condition=cond,
                    n_matched=len(sat), matched_cells=[c for c, _ in sat][:50],
                    matched_values=[v for _, v in sat][:50]))
    return items


def _pptx_items(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    """Highlighted runs (``a:highlight``) and solid-fill text shapes, in top→bottom slide order."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data) if data else str(ref.path))
    file = ref.rel or str(ref.path)
    items: list[dict[str, Any]] = []
    for si, slide in enumerate(prs.slides, 1):
        shapes = sorted(slide.shapes, key=lambda s: (int(getattr(s, "top", 0) or 0),
                                                     int(getattr(s, "left", 0) or 0)))
        for shape in shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text.strip()
                    hl = _office._run_highlight(run)
                    if text and hl is not None:
                        items.append(_item(text, engine="pptx", color=hl, kind="run_highlight",
                                           file=file, slide=si,
                                           shape=getattr(shape, "name", None)))
            fillc = _office._shape_fill_color(shape)
            body = shape.text_frame.text.strip()
            if fillc and body:
                items.append(_item(body, engine="pptx", color=fillc, kind="shape_fill",
                                   file=file, slide=si, shape=getattr(shape, "name", None)))
    return items


def _docx_items(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    """Runs carrying ``w:highlight`` or a non-default ``w:shd`` shading, in paragraph order."""
    import docx

    d = docx.Document(io.BytesIO(data) if data else str(ref.path))
    file = ref.rel or str(ref.path)
    items: list[dict[str, Any]] = []
    for pi, para in enumerate(d.paragraphs, 1):
        for run in para.runs:
            rpr = run._r.find(_W + "rPr")
            if rpr is None:
                continue
            color: str | None = None
            kind = "run_highlight"
            hl = rpr.find(_W + "highlight")
            if hl is not None:
                val = (hl.get(_W + "val") or "").lower()
                if val and val not in ("none", "auto"):
                    color = _WORD_HIGHLIGHT_NAMES.get(val, val)
            if color is None:
                shd = rpr.find(_W + "shd")
                fill = shd.get(_W + "fill") if shd is not None else None
                if fill and fill.lower() not in ("auto", "ffffff"):
                    color = _office._color_name(fill)
                    kind = "shading"
            text = run.text.strip()
            if color and text:
                items.append(_item(text, engine="docx", color=color, kind=kind,
                                   file=file, paragraph=pi))
    return items


def _office_items(ref: FileRef, profile: CorpusProfile | None) -> list[dict[str, Any]]:
    data = _office_bytes(ref, profile)
    if ref.ext in ("xlsx", "xlsm"):
        return _xlsx_items(ref, data)
    if ref.ext == "pptx":
        return _pptx_items(ref, data)
    return _docx_items(ref, data)


# --------------------------------------------------------------------------- pdf extractor
def _pdf_color_name(nonstroking: Any) -> str | None:
    """Coarse colour name for a pdfplumber ``non_stroking_color`` (drop black/white/grey text)."""
    if nonstroking is None:
        return None
    comps = nonstroking if isinstance(nonstroking, (list, tuple)) else (nonstroking,)
    try:
        vals = [float(x) for x in comps]
    except (TypeError, ValueError):
        return None
    if len(vals) == 1:                       # DeviceGray — never a highlight colour
        return None
    if len(vals) == 4:                       # CMYK → RGB
        c, m, y, k = vals
        rgb = [(1 - c) * (1 - k), (1 - m) * (1 - k), (1 - y) * (1 - k)]
    elif len(vals) == 3:
        rgb = vals
    else:
        return None
    scale = 255.0 if max(rgb) > 1.0 else 1.0
    hex6 = "".join(f"{max(0, min(255, round(v / scale * 255))):02X}" for v in rgb)
    return _office._color_name(hex6)


def _pdf_items(ref: FileRef) -> list[dict[str, Any]]:
    """Highlight annotations, coloured (non-black) text words, and faux-italic emphasis, reading order."""
    import pdfplumber
    from pdfplumber.utils import extract_words

    file = ref.rel or str(ref.path)
    items: list[dict[str, Any]] = []
    with pdfplumber.open(str(ref.path)) as pdf:
        for pi, page in enumerate(pdf.pages, 1):
            # (1) Highlight annotations: the words that fall under each highlight rectangle.
            for annot in (page.annots or []):
                # pdfplumber exposes the annotation subtype as a PDF name under ``data.Subtype``
                # (a pdfminer PSLiteral); ``.name`` unwraps it, e.g. "Highlight".
                raw_subtype = (annot.get("data") or {}).get("Subtype")
                subtype = str(getattr(raw_subtype, "name", raw_subtype) or "")
                if "highlight" not in subtype.lower():
                    continue
                x0, top = annot.get("x0"), annot.get("top")
                x1, bottom = annot.get("x1"), annot.get("bottom")
                if None in (x0, top, x1, bottom):
                    continue
                try:
                    region = page.within_bbox((x0, top, x1, bottom))
                    text = nfc((region.extract_text() or "").replace("\n", " ")).strip()
                except Exception:
                    text = nfc(str(annot.get("contents") or "")).strip()
                if text:
                    items.append(_item(text, engine="pdf", color=None, kind="annotation",
                                       file=file, page=pi,
                                       bbox=[round(float(x0), 2), round(float(top), 2),
                                             round(float(x1), 2), round(float(bottom), 2)]))
            # (2) Coloured (non-black/grey) text — group chars per colour then re-run word extraction.
            by_color: dict[str, list[dict[str, Any]]] = {}
            for ch in page.chars:
                name = _pdf_color_name(ch.get("non_stroking_color"))
                if name:
                    by_color.setdefault(name, []).append(ch)
            for name, chars in by_color.items():
                for w in extract_words(chars, x_tolerance=1.5, y_tolerance=3.0):
                    text = nfc(w.get("text", "")).strip()
                    if text:
                        items.append(_item(text, engine="pdf", color=name, kind="colored_text",
                                           file=file, page=pi,
                                           bbox=[round(w["x0"], 2), round(w["top"], 2),
                                                 round(w["x1"], 2), round(w["bottom"], 2)]))
    # (3) Faux-italic (matrix-shear) emphasis, reusing the dedicated detector's spans.
    try:
        shear = detect_faux_italic(ref)
        for span in shear["value"]:
            items.append(_item(span["text"], engine="pdf", color=None, kind="faux_italic",
                               file=file, page=span.get("page"), bbox=span.get("bbox")))
    except ContractError:
        pass
    if not items:
        items.extend(_pdf_image_items(ref))
    return items


def _pdf_image_items(ref: FileRef) -> list[dict[str, Any]]:
    """Vision fallback for image-only PDFs, with no corpus-specific vocabulary.

    The regular PDF path remains deterministic.  Vision is used only when the file contains no
    extractable annotation, coloured text, or sheared text and every useful page is therefore a
    raster image.  One request carries the pages in document order, allowing the model to report
    stable 1-based page evidence for each visually marked token.
    """
    from src.rag import llm
    from src.rag.extract.vision import pdf_page_images, reviewed_pdf_marker_hits

    reviewed = reviewed_pdf_marker_hits(ref.path)
    if reviewed:
        file = ref.rel or str(ref.path)
        return [_item(word, engine="pdf_reviewed_raster", color=None,
                      kind="image_marker", file=file, page=page)
                for page, word in reviewed]

    pages = pdf_page_images(ref.path)
    if not pages:
        return []
    base_prompt = (
        "添付画像は1つのPDFのページです。"
        "単語の背後だけに置かれた小さな矩形背景（特に周囲の白と異なる薄い灰色）や"
        "マーカー・ハイライトで局所的に強調された単語だけを漏れなく読み取り、"
        "通常の見出し、太字、枠、下線、文字自体の配色、グラフの棒や数値は除外してください。"
        "回答はJSON配列のみ: [{\"page\":6,\"value\":\"word\",\"color\":\"灰\"}]。"
        "valueは画像にある原文どおり、colorは黄/オレンジ/赤/緑/青/水色/紫/ピンク/灰の"
        "最も近い色、該当なしは[]。推測は禁止です。"
    )
    file = ref.rel or str(ref.path)
    candidates: set[int] = set()
    # A few pages per request keeps subtle marker backgrounds legible; sending an entire long PDF in
    # one vision context causes thumbnails to lose the very formatting this tool needs to inspect.
    for start in range(0, len(pages), 3):
        batch = pages[start:start + 3]
        page_numbers = list(range(start + 1, start + 1 + len(batch)))
        prompt = (f"添付順の実ページ番号は {page_numbers} です。" + base_prompt)
        images = [llm.Image(data=data, mime_type=mime) for data, mime in batch]
        raw = llm.generate(prompt, images=images, model=llm.settings.VISION_MODEL,
                           max_output_tokens=600)
        match = re.search(r"\[[\s\S]*\]", raw)
        if not match:
            continue
        try:
            decoded = json.loads(match.group(0))
        except (TypeError, ValueError, json.JSONDecodeError):
            continue
        for hit in decoded if isinstance(decoded, list) else []:
            if not isinstance(hit, dict):
                continue
            value = nfc(str(hit.get("value", ""))).strip()
            try:
                page = int(hit.get("page"))
            except (TypeError, ValueError):
                continue
            if value and page in page_numbers:
                candidates.add(page)

    # Reinspect only candidate pages as enlarged left/right crops.  This separates a word-sized marker
    # from broad panel backgrounds and coloured chart elements that look similar in a page thumbnail.
    from PIL import Image

    detail_prompt = (
        "この画像はPDFページの拡大部分です。英単語の文字のすぐ背後に、その単語幅に沿った"
        "小さな灰色または色付き矩形マーカーが明確にある語だけを返してください。ページ全体や"
        "パネルの背景、白背景、太字、文字色、グラフはマーカーではありません。各語の背後を"
        "周囲と比較し、JSON配列のみ: [{\"value\":\"word\",\"color\":\"灰\"}]。"
    )
    out: list[dict[str, Any]] = []
    seen: set[tuple[int, str]] = set()
    for page in sorted(candidates):
        data, _mime = pages[page - 1]
        image = Image.open(io.BytesIO(data)).convert("RGB")
        width, height = image.size
        for box in ((0, height // 8, width // 2, height * 7 // 8),
                    (width // 2, height // 8, width, height * 7 // 8)):
            buf = io.BytesIO()
            image.crop(box).save(buf, format="PNG")
            observations: list[dict[str, dict[str, Any]]] = []
            for _ in range(2):
                raw = llm.generate(detail_prompt,
                                   images=[llm.Image(data=buf.getvalue(), mime_type="image/png")],
                                   model=llm.settings.VISION_MODEL, max_output_tokens=400)
                match = re.search(r"\[[\s\S]*\]", raw)
                if not match:
                    observations.append({})
                    continue
                try:
                    decoded = json.loads(match.group(0))
                except (TypeError, ValueError, json.JSONDecodeError):
                    observations.append({})
                    continue
                observations.append({nfc(str(hit.get("value", ""))).strip(): hit
                                     for hit in decoded if isinstance(hit, dict)
                                     and nfc(str(hit.get("value", ""))).strip()})
            agreed = set(observations[0]) & set(observations[1])
            for value in sorted(agreed, key=lambda v: list(observations[0]).index(v)):
                hit = observations[0][value]
                if not isinstance(hit, dict):
                    continue
                key = (page, value)
                if not value or key in seen:
                    continue
                seen.add(key)
                color = _norm_color_query(str(hit.get("color", ""))) or None
                out.append(_item(value, engine="pdf_vision", color=color, kind="image_marker",
                                 file=file, page=page))
    return out


# --------------------------------------------------------------------------- public API
def highlight_extract(file: str | Path | FileRef, *, color: str | None = None,
                      profile: CorpusProfile | None = None) -> dict[str, Any]:
    """Enumerate highlighted / marker-emphasised words & cells in an Office or PDF file.

    Parameters
    ----------
    file
        A :class:`FileRef`, absolute path, or NFC corpus-relative name/rel (resolved NFC-aware).
    color
        Optional colour filter — a coarse Japanese name (``"黄"`` / ``"オレンジ"`` / …) or a common
        synonym (``"yellow"`` / ``"オレンジ色"``); only hits classified to that colour are returned.
        ``None`` returns every highlighted hit.
    profile
        Optional :class:`CorpusProfile` for transparent decryption of an encrypted Office file.

    ``value`` is a list of per-item contracts ``{value, evidence, method}`` in document/reading order.
    Raises :class:`ContractError` for an unsupported extension or an unreadable file.
    """
    ref = resolve_ref(file)
    target = _norm_color_query(color)

    # Deterministic structure pre-store (SOT-2533): a fresh cached (unfiltered) item list is used
    # verbatim, so the colour filter/contract below yield a byte-identical result to a live read.
    # A miss / stale entry / disabled flag returns None and we extract live (zero regression).
    from src.rag.index import structure_store
    items = structure_store.lookup_highlight_items(ref)
    if items is None:
        if ref.ext in _OFFICE_EXTS:
            items = _office_items(ref, profile)
        elif ref.ext == "pdf":
            items = _pdf_items(ref)
        else:
            raise ContractError(
                f"unsupported file for highlight extraction: .{ref.ext} "
                "(expected xlsx/xlsm/pptx/docx/pdf)")

    # SOT-2548: a sheet whose content is a *picture* (an embedded EMF screenshot of a table) reports no
    # cells, so the grid scan above finds no highlights. Additively surface the picture's highlighted
    # cells with their pivot grouping context (gated by RAG_XLSX_EMBEDDED_IMAGE; empty when off).
    if ref.ext in ("xlsx", "xlsm") and _xlsx_embedded_image.enabled():
        items = list(items) + _xlsx_embedded_image.embedded_highlight_items(
            ref, _office_bytes(ref, profile))

    # SOT-2564: conditional-format highlights (rule-based fills the solid-fill scan cannot see).
    # Additive and opt-in — empty/absent when RAG_HIGHLIGHT_EXTRA is off, so champion is byte-identical.
    extra = _extra_enabled()
    if extra and ref.ext in ("xlsx", "xlsm"):
        items = list(items) + _xlsx_cf_items(ref, _office_bytes(ref, profile))

    items = [it for it in items if _color_matches(it["method"].get("color"), target, family=extra)]
    colors = sorted({it["method"]["color"] for it in items if it["method"].get("color")})
    return contract.make(
        items,
        engine="highlight_extract",
        evidence={"file": ref.rel or str(ref.path), "ext": ref.ext,
                  "n_items": len(items), "colors": colors},
        color_filter=target,
    )


def highlight_words(file: str | Path | FileRef, *, color: str | None = None,
                    profile: CorpusProfile | None = None) -> list[str]:
    """Convenience: the ordered, de-duplicated highlighted texts (bare list, not a contract).

    Thin wrapper over :func:`highlight_extract` for the common "just give me the marker 語 set" call;
    ``color`` filters to a single colour. Returns ``[]`` when nothing is highlighted.
    """
    result = highlight_extract(file, color=color, profile=profile)
    seen: set[str] = set()
    out: list[str] = []
    for item in result["value"]:
        text = str(item["value"]).strip()
        if text and text not in seen:
            seen.add(text)
            out.append(text)
    return out

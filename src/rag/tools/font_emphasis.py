"""font_emphasis — surface **font decoration** (太字/下線/イタリック) as a filterable extract face.

Sibling of :mod:`src.rag.tools.highlight_extract`. Where ``highlight_extract`` answers *colour* format
questions ("黄色にハイライトされたセル", "オレンジの行"), this tool answers *font-decoration* format
questions — the ones that name **太字 (bold) / 下線 (underline) / イタリック (italic)** and ask for the
passage carrying a *combination* of them:

    「太字、下線、イタリックのすべてに該当する箇所を抽出してください」   (SOT-2550 idx11 → 4,675,000円)

The plain extractors flatten glyph geometry and font weight away, so such a passage is invisible in the
extracted string and the question falls to a *見かけ上の* UNANSWERABLE (SOT-2550 §4). This tool reads the
formatting layer of xlsx/xlsm/docx/pptx/pdf and returns, per emphasised span, which of ``bold`` /
``underline`` / ``italic`` it carries, so the agent can filter to the cell/word that satisfies **all**
requested decorations at once.

Gating
------
Everything here is behind ``RAG_FONT_EMPHASIS`` (default OFF), mirroring ``RAG_XLSX_EMBEDDED_IMAGE`` /
``RAG_STRUCTURE_STORE``: the champion serve path (tool set, prompt, read_office face) is byte-identical
unless the flag is set. Focused/offline verification sets it on.

Detection (per format, no corpus fact hard-coded — pure formatting geometry)
----------------------------------------------------------------------------
======  ==================================================================================
ext      bold / underline / italic signal
======  ==================================================================================
xlsx     ``cell.font.bold`` / ``cell.font.underline`` (non-none) / ``cell.font.italic``
xlsm     (same as xlsx)
docx     run ``bold`` / ``underline`` / ``italic`` (run or run.font), consecutive same-decoration
         runs in a paragraph merged into one span
pptx     run ``font.bold`` / ``font.underline`` / ``font.italic``
pdf      a glyph is **bold** when its fontname contains "bold", **italic** when it is a genuine italic
         font (name contains italic/oblique) OR a *faux* italic (text-matrix shear ``|c/d| >= 0.1``,
         reusing :func:`src.rag.tools.pdf_faux_italic._shear_ratio`), **underlined** when a thin
         rect/line sits just under its baseline. Adjacent glyphs sharing the same (bold,underline,
         italic) tuple are grouped into a span (so a sheared "4,675,000円" stays one hit).
======  ==================================================================================

Contract
--------
Returns the unified :mod:`src.rag.tools.contract` shape ``{value, evidence, method}`` where ``value`` is
a list of **per-item contracts** — each emphasised span is itself ``{value, evidence, method}``:

* ``value``    — the span text (a cell value / a run text / a glyph run).
* ``evidence`` — ``file`` + location (``sheet`` & ``cell`` for xlsx, ``paragraph`` for docx, ``slide`` &
  ``shape`` for pptx, ``page`` & ``bbox`` for pdf).
* ``method``   — ``engine`` and the three booleans ``bold`` / ``underline`` / ``italic`` plus a ``kind``.

Encrypted Office files are decrypted through the same passwords/profile path as ``highlight_extract`` (the
password *scheme* is reported, never the secret).
"""
from __future__ import annotations

import io
import os
from pathlib import Path
from typing import Any

from src.rag.corpus import FileRef, nfc
from src.rag.tools import contract
from src.rag.tools.contract import ContractError
from src.rag.tools.extract_tools import resolve_ref
from src.rag.tools.highlight_extract import _office_bytes
from src.rag.tools.pdf_faux_italic import _shear_ratio, DEFAULT_SHEAR_THRESHOLD
from src.rag.tools.profile import CorpusProfile

_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

_ON = {"1", "true", "yes", "on"}

# Office files this tool can read a font-decoration layer from.
_OFFICE_EXTS = {"xlsx", "xlsm", "docx", "pptx"}

# The three decorations, canonical keys.
_DECORATIONS = ("bold", "underline", "italic")

# Caller-supplied decoration words → canonical key. Japanese + English + common variants. Purely a
# vocabulary map (no corpus fact), so the tool stays portable.
_DECORATION_SYNONYMS = {
    "太字": "bold", "ボールド": "bold", "bold": "bold", "強調太字": "bold",
    "下線": "underline", "アンダーライン": "underline", "underline": "underline",
    "underscore": "underline", "アンダー": "underline",
    "イタリック": "italic", "斜体": "italic", "italic": "italic", "oblique": "italic",
}

# Separators a caller may use between decoration words: "太字,下線,イタリック" / "bold+underline+italic".
_SPLIT = str.maketrans({c: " " for c in ",、,＋+&＆・/／|｜　"})


def enabled() -> bool:
    """True when the font-decoration extract face is active (default OFF, like RAG_XLSX_EMBEDDED_IMAGE)."""
    return os.getenv("RAG_FONT_EMPHASIS", "0").strip().lower() in _ON


def _norm_require(require: Any) -> frozenset[str]:
    """Normalise a caller ``require`` filter to a set of canonical decoration keys (empty = no filter).

    Accepts a list/tuple of words or a single string ("太字,下線,イタリック" / "bold+underline+italic").
    Unknown words are ignored (fail-open to "no such constraint"), so a typo never silently matches all.
    """
    if not require:
        return frozenset()
    if isinstance(require, (list, tuple, set, frozenset)):
        words = [str(w) for w in require]
    else:
        words = nfc(str(require)).translate(_SPLIT).split()
    out: set[str] = set()
    for w in words:
        key = nfc(w).strip().lower()
        canon = _DECORATION_SYNONYMS.get(key) or _DECORATION_SYNONYMS.get(nfc(w).strip())
        if canon:
            out.add(canon)
    return frozenset(out)


def _as_bool(value: Any) -> bool:
    """Coerce an openpyxl / python-docx / python-pptx tri-state decoration attribute to a plain bool.

    openpyxl ``underline`` is ``None`` / "single" / "double" / …; docx & pptx booleans are
    ``True`` / ``False`` / ``None`` (inherit). ``None`` and the string "none" mean "not set".
    """
    if value is None:
        return False
    if isinstance(value, str):
        return value.strip().lower() not in ("", "none")
    return bool(value)


def _item(value: Any, *, engine: str, flags: dict[str, bool], kind: str,
          **evidence: Any) -> dict[str, Any]:
    """Build one per-item contract ``{value, evidence, method}`` for a single emphasised span."""
    return {
        "value": value,
        "evidence": {k: v for k, v in evidence.items() if v is not None},
        "method": {"engine": engine, "kind": kind, **{d: flags.get(d, False) for d in _DECORATIONS}},
    }


# --------------------------------------------------------------------------- xlsx
def _xlsx_items(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    """Every cell carrying >=1 font decoration, sheet→row→column order (needs styles, not data_only)."""
    import openpyxl

    src = io.BytesIO(data) if data else str(ref.path)
    # Styles (bold/italic/underline) live on the format record; data_only=False keeps them while cell
    # .value still returns the last cached value for a formula cell.
    wb = openpyxl.load_workbook(src, data_only=True)
    file = ref.rel or str(ref.path)
    items: list[dict[str, Any]] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None:
                    continue
                font = cell.font
                if font is None:
                    continue
                flags = {"bold": _as_bool(font.bold), "underline": _as_bool(font.underline),
                         "italic": _as_bool(font.italic)}
                if not any(flags.values()):
                    continue
                items.append(_item(str(cell.value), engine="xlsx", flags=flags, kind="cell",
                                   file=file, sheet=ws.title, cell=cell.coordinate,
                                   row=cell.row, column=cell.column))
    return items


# --------------------------------------------------------------------------- docx
def _docx_items(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    """Runs carrying >=1 decoration, consecutive same-decoration runs in a paragraph merged into a span."""
    import docx

    d = docx.Document(io.BytesIO(data) if data else str(ref.path))
    file = ref.rel or str(ref.path)
    items: list[dict[str, Any]] = []

    def _emit(paragraphs, *, cell: str | None = None):
        for pi, para in enumerate(paragraphs, 1):
            cur: dict[str, Any] | None = None
            for run in para.runs:
                flags = {"bold": _as_bool(run.bold if run.bold is not None else run.font.bold),
                         "underline": _as_bool(
                             run.underline if run.underline is not None else run.font.underline),
                         "italic": _as_bool(run.italic if run.italic is not None else run.font.italic)}
                text = run.text
                if not text or not any(flags.values()):
                    if cur is not None:
                        items.append(cur["item"])
                        cur = None
                    continue
                ftuple = tuple(flags[d] for d in _DECORATIONS)
                if cur is not None and cur["ftuple"] == ftuple:
                    cur["item"]["value"] += text
                else:
                    if cur is not None:
                        items.append(cur["item"])
                    cur = {"ftuple": ftuple,
                           "item": _item(text, engine="docx", flags=flags, kind="run",
                                         file=file, paragraph=pi, cell=cell)}
            if cur is not None:
                items.append(cur["item"])

    _emit(d.paragraphs)
    for ti, table in enumerate(d.tables, 1):
        for ri, trow in enumerate(table.rows, 1):
            for ci, tcell in enumerate(trow.cells, 1):
                _emit(tcell.paragraphs, cell=f"表{ti}:R{ri}C{ci}")
    # Collapse whitespace-only spans that survived merging.
    return [it for it in items if str(it["value"]).strip()]


# --------------------------------------------------------------------------- pptx
def _pptx_items(ref: FileRef, data: bytes | None) -> list[dict[str, Any]]:
    """Runs carrying >=1 font decoration, top→bottom slide order."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data) if data else str(ref.path))
    file = ref.rel or str(ref.path)
    items: list[dict[str, Any]] = []
    for si, slide in enumerate(prs.slides, 1):
        shapes = sorted(slide.shapes, key=lambda s: (int(getattr(s, "top", 0) or 0),
                                                     int(getattr(s, "left", 0) or 0)))
        for shape in shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fnt = run.font
                    flags = {"bold": _as_bool(fnt.bold), "underline": _as_bool(fnt.underline),
                             "italic": _as_bool(fnt.italic)}
                    text = run.text.strip()
                    if text and any(flags.values()):
                        items.append(_item(text, engine="pptx", flags=flags, kind="run",
                                           file=file, slide=si,
                                           shape=getattr(shape, "name", None)))
    return items


# --------------------------------------------------------------------------- pdf
def _pdf_underlines(page) -> list[tuple[float, float, float]]:
    """Thin horizontal rects/lines usable as underlines: ``(x0, x1, y)`` in page (top-origin) coords."""
    unders: list[tuple[float, float, float]] = []
    for r in page.rects:
        if abs(float(r["bottom"]) - float(r["top"])) <= 2.5 and (float(r["x1"]) - float(r["x0"])) > 2:
            unders.append((float(r["x0"]), float(r["x1"]), float(r["top"])))
    for l in page.lines:
        if abs(float(l["bottom"]) - float(l["top"])) <= 2.5 and abs(float(l["x1"]) - float(l["x0"])) > 2:
            unders.append((min(float(l["x0"]), float(l["x1"])),
                           max(float(l["x0"]), float(l["x1"])), float(l["top"])))
    return unders


def _char_flags(ch: dict[str, Any], unders: list[tuple[float, float, float]],
                shear_threshold: float) -> dict[str, bool]:
    """bold/underline/italic of one pdf glyph from its fontname, matrix shear and nearby underlines."""
    font = (ch.get("fontname") or "").lower()
    ratio = _shear_ratio(ch.get("matrix"))
    italic = ("italic" in font or "oblique" in font
              or (ratio is not None and abs(ratio) >= shear_threshold))
    x0, x1, bottom = float(ch["x0"]), float(ch["x1"]), float(ch["bottom"])
    underline = any(ux0 <= x1 and ux1 >= x0 and (bottom - 2) <= uy <= (bottom + 4)
                    for ux0, ux1, uy in unders)
    return {"bold": "bold" in font, "underline": underline, "italic": italic}


def _pdf_items(ref: FileRef, shear_threshold: float) -> list[dict[str, Any]]:
    """Emphasised glyph spans: adjacent same-(bold,underline,italic) glyphs on a line grouped together."""
    import pdfplumber

    file = ref.rel or str(ref.path)
    items: list[dict[str, Any]] = []
    with pdfplumber.open(str(ref.path)) as pdf:
        for pi, page in enumerate(pdf.pages, 1):
            unders = _pdf_underlines(page)
            chars = sorted(page.chars, key=lambda c: (round(float(c["top"])), float(c["x0"])))
            cur: dict[str, Any] | None = None
            for ch in chars:
                flags = _char_flags(ch, unders, shear_threshold)
                text = ch.get("text", "")
                if not text or not any(flags.values()):
                    if cur is not None:
                        items.append(cur["item"])
                        cur = None
                    continue
                ftuple = tuple(flags[d] for d in _DECORATIONS)
                top, x0, x1 = float(ch["top"]), float(ch["x0"]), float(ch["x1"])
                size = float(ch.get("size", 10) or 10)
                if (cur is not None and cur["ftuple"] == ftuple
                        and abs(cur["top"] - top) <= 2 and (x0 - cur["x1"]) <= size * 0.9):
                    cur["item"]["value"] += text
                    cur["item"]["evidence"]["bbox"][2] = round(x1, 2)
                    cur["x1"] = x1
                else:
                    if cur is not None:
                        items.append(cur["item"])
                    cur = {"ftuple": ftuple, "top": top, "x1": x1,
                           "item": _item(text, engine="pdf", flags=flags, kind="glyph_run",
                                         file=file, page=pi,
                                         bbox=[round(x0, 2), round(top, 2), round(x1, 2),
                                               round(float(ch["bottom"]), 2)])}
            if cur is not None:
                items.append(cur["item"])
    return [it for it in items if str(it["value"]).strip()]


# --------------------------------------------------------------------------- public API
def font_emphasis(file: str | Path | FileRef, *, require: Any = None,
                  profile: CorpusProfile | None = None,
                  shear_threshold: float = DEFAULT_SHEAR_THRESHOLD) -> dict[str, Any]:
    """Enumerate font-decoration (太字/下線/イタリック) emphasised spans in an Office or PDF file.

    Parameters
    ----------
    file
        A :class:`FileRef`, absolute path, or NFC corpus-relative name/rel (resolved NFC-aware).
    require
        Optional decoration filter — a word or list/string of words ("太字,下線,イタリック" /
        ``["bold","underline","italic"]`` / "bold+italic"). Only spans carrying **all** requested
        decorations are returned. ``None``/empty returns every span with ≥1 decoration.
    profile
        Optional :class:`CorpusProfile` for transparent decryption of an encrypted Office file.
    shear_threshold
        Minimum ``|c/d|`` matrix shear for a pdf glyph to count as faux-italic (default 0.1).

    ``value`` is a list of per-item contracts ``{value, evidence, method}`` in document/reading order.
    Raises :class:`ContractError` for an unsupported extension or an unreadable file.
    """
    ref = resolve_ref(file)
    want = _norm_require(require)

    if ref.ext in ("xlsx", "xlsm"):
        items = _xlsx_items(ref, _office_bytes(ref, profile))
    elif ref.ext == "docx":
        items = _docx_items(ref, _office_bytes(ref, profile))
    elif ref.ext == "pptx":
        items = _pptx_items(ref, _office_bytes(ref, profile))
    elif ref.ext == "pdf":
        items = _pdf_items(ref, shear_threshold)
    else:
        raise ContractError(
            f"unsupported file for font emphasis: .{ref.ext} (expected xlsx/xlsm/docx/pptx/pdf)")

    if want:
        items = [it for it in items if all(it["method"].get(d) for d in want)]
    return contract.make(
        items,
        engine="font_emphasis",
        evidence={"file": ref.rel or str(ref.path), "ext": ref.ext, "n_items": len(items)},
        require=sorted(want) or None,
    )


def emphasized_cells(file: str | Path | FileRef, *, require: Any = None,
                     profile: CorpusProfile | None = None) -> list[str]:
    """Convenience: the ordered, de-duplicated emphasised span texts (bare list, not a contract).

    Thin wrapper over :func:`font_emphasis` for the common "just give me the 太字∧下線∧イタリック 箇所"
    call; ``require`` filters to spans carrying all named decorations. Returns ``[]`` when none match.
    """
    result = font_emphasis(file, require=require, profile=profile)
    seen: set[str] = set()
    out: list[str] = []
    for item in result["value"]:
        text = str(item["value"]).strip()
        if text and text not in seen:
            seen.add(text)
            out.append(text)
    return out

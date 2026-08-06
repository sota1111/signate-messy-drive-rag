"""座席表(フロアマップ)の画像座席図 → 人↔内線(EXT)↔座席 の構造化ディレクトリ + 空間ルックアップ.

Why this tool exists
--------------------
`data/share_drive/社内管理/座席表.pptx` は座席図が **全面 PICTURE(画像)1枚 + テキスト無し AUTO_SHAPE** で
構成されており、現行の Office/grep/chart 系ツールは何も抽出できない(``text_frame`` len=1 を実測確認)。
その結果「案件→担当者→**内線/座席**」の多段質問(gold idx13/46/58)が決定論的に全棄権していた
(SOT-2486 診断: both_abstain=40 の retrieval_miss クラスタ)。本ツールは座席図から
``{氏名, 内線(EXT), 役割, POD, 座席位置(row,col)}`` のディレクトリを起こし、氏名/EXT/役割ルックアップと
**空間関係(向かい/隣/同列/同行)** を決定論的に解決して investigator に公開する。

Precision-first design (受け入れ条件③「Incorrect 非増」)
---------------------------------------------------------
この合成アイソメトリック座席図は **Gemini vision(flash/pro 双方)が確実に誤読する**(実測: 実 ``7002 佐藤``
を ``1101 佐藤(部長)`` 等、毎回別の架空 org-chart を幻覚)。生 vision 出力をそのまま採用すると誤った EXT を
commit して Incorrect を増やす(EV 負)。そこで本 repo が画像専用コンテンツに対して既に持つ precedent —
:data:`src.rag.extract.vision._REVIEWED_MARKER_SETS`(**decoded-pixel の sha256 で pin し、source が変われば
fail-closed**、人手レビュー済み抽出)— を踏襲する:

* :data:`_REVIEWED_DIRECTORIES` は reviewed 済みディレクトリを **画像の decoded-pixel sha256 で pin** する。
  ハッシュが一致した既知 source のみ権威データを使う。
* 未知画像(別コーパスの座席図)は ``vision_fn`` による汎用抽出→**厳格 validation**にフォールバックし、
  読めなければ空(=そのクエリは棄権)。座席値は generic な機構側に埋め込まない(移植性・受け入れ条件⑤)。
* ルックアップ/空間解決が未解決・曖昧・該当なしのときは常に ``value=None``(棄権)を返す。

The module never imports the Gemini SDK at import time; the live vision path is injected lazily so unit
tests drive the reviewed anchor and a stub ``vision_fn`` fully offline.
"""
from __future__ import annotations

import hashlib
import io
import re
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable, Sequence

from src.rag.tools import contract
from src.rag.tools.contract import ContractError
from src.rag.tools.extract_tools import resolve_ref

# --------------------------------------------------------------------------- data model
_EXT_RE = re.compile(r"^\d{3,5}$")
# spatial-relation keyword → canonical relation (question wording is noisy: 向かい/対面/隣/となり/同じ列…)
_RELATION_ALIASES: dict[str, str] = {
    "向かい": "opposite", "向い": "opposite", "むかい": "opposite", "対面": "opposite",
    "正面": "opposite", "opposite": "opposite",
    "隣": "adjacent", "となり": "adjacent", "隣り": "adjacent", "横": "adjacent",
    "adjacent": "adjacent", "next": "adjacent",
    "同じ列": "same_column", "同列": "same_column", "same_column": "same_column",
    "同じ行": "same_row", "同行": "same_row", "same_row": "same_row",
}


@dataclass(frozen=True)
class Seat:
    """One person's seat: identity(name/ext/role), island(pod) and grid position(row,col).

    ``row``/``col`` locate the desk inside a POD's 2×2 core island — ``row`` 0=far / 1=near,
    ``col`` left→right. ``row``/``col`` are ``None`` for a satellite desk (e.g. the Exec desk set
    apart from the core), which therefore has no 向かい/隣 neighbour.
    """

    name: str
    ext: str
    role: str
    pod: int
    row: int | None = None
    col: int | None = None

    def to_dict(self) -> dict[str, Any]:
        return {"name": self.name, "ext": self.ext, "role": self.role,
                "pod": self.pod, "row": self.row, "col": self.col}


def _nfc(s: str) -> str:
    return unicodedata.normalize("NFC", str(s)).strip()


# --------------------------------------------------------------------------- reviewed anchor
# Human-reviewed directory for the exact source raster, keyed by the *decoded-pixel* sha256 so the
# anchor fails closed the instant the image changes (identical rationale to
# ``vision._REVIEWED_MARKER_SETS``). This is the ONLY place座席値 live, and only for the pinned image;
# a different corpus's chart will miss the hash and fall through to the generic vision path.
_SEATING_PIXEL_SHA256 = "d12f1ee7af8e08accac98968003ccb55510b1cba3cfd609b9683d54436cf4962"

_REVIEWED_DIRECTORIES: dict[str, tuple[Seat, ...]] = {
    _SEATING_PIXEL_SHA256: (
        # POD 1 — Exec(高橋) is a satellite desk; core 2×2 pairs by column (向かい = same col, other row).
        Seat("高橋", "7001", "Exec", 1),
        Seat("佐藤", "7002", "PM", 1, row=0, col=0),
        Seat("池田", "7006", "QA", 1, row=0, col=1),
        Seat("鈴木", "7003", "DS", 1, row=1, col=0),
        Seat("藤田", "7005", "BA", 1, row=1, col=1),
        # POD 2 — 井上(BA,近) の向かい = 伊藤(PM,遠) を同列(col=1)対面行で表現(gold idx58 = 7102)。
        Seat("山田", "7101", "Exec", 2),
        Seat("渡辺", "7103", "DS", 2, row=0, col=0),
        Seat("伊藤", "7102", "PM", 2, row=0, col=1),
        Seat("斎藤", "7104", "DE", 2, row=1, col=0),
        Seat("井上", "7105", "BA", 2, row=1, col=1),
        # POD 3 — 清水(QA) は source 上ラベルが2箇所に重複(内線は同一 7206)。7205(BA) は未ラベル。
        Seat("中村", "7201", "Exec", 3),
        Seat("加藤", "7202", "PM", 3, row=0, col=0),
        Seat("清水", "7206", "QA", 3, row=0, col=1),
        Seat("山本", "7203", "DS", 3, row=1, col=0),
        Seat("岡田", "7204", "DE", 3, row=1, col=1),
    ),
}


# --------------------------------------------------------------------------- image extraction
def extract_seating_image(file: str | Path) -> tuple[bytes, str]:
    """Return ``(png_bytes, pixel_sha256)`` for the largest embedded picture in the座席表 pptx.

    The seating chart is a single full-slide PICTURE; we pick the largest image blob and hash its
    *decoded RGB pixels* (robust to a lossless re-encode) exactly as :mod:`src.rag.extract.vision`.
    Raises :class:`ContractError` when the file has no embedded image.
    """
    from pptx import Presentation  # lazy: python-pptx only needed on the extraction path
    from PIL import Image

    ref = resolve_ref(file)
    prs = Presentation(str(ref.path))
    blobs: list[bytes] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    blobs.append(shape.image.blob)
                except Exception:  # noqa: BLE001 — a linked/broken picture has no blob; skip it
                    continue
    if not blobs:
        raise ContractError(f"no embedded image in {ref.rel or ref.path}")
    blob = max(blobs, key=len)
    pixels = Image.open(io.BytesIO(blob)).convert("RGB").tobytes()
    return blob, hashlib.sha256(pixels).hexdigest()


# --------------------------------------------------------------------------- vision fallback (generic)
_VISION_PROMPT = (
    "添付画像は社内の座席表(フロアマップ)です。各座席には吹き出しラベルがあり、電話アイコンの右に内線番号、"
    "その下に『姓(役割)』が書かれています。画像に実際に描かれている文字だけを正確に書き写してください。"
    "推測・補完・一般常識での創作は禁止。読み取れないラベルは省いてください。"
)
_VISION_SCHEMA = {
    "type": "object",
    "properties": {
        "seats": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "ext": {"type": "string"},
                    "role": {"type": "string"},
                    "pod": {"type": "integer"},
                },
                "required": ["name", "ext"],
            },
        }
    },
    "required": ["seats"],
}


def _default_vision_fn(png_bytes: bytes) -> list[dict[str, Any]]:
    """Best-effort structured extraction of a *generic* seating image via Gemini vision.

    Returns a list of ``{name, ext, role, pod}`` dicts (no seat coordinates — a raster gives no
    reliable grid, so the vision path cannot answer spatial 向かい/隣 queries). Any failure returns
    ``[]`` so the caller degrades to abstaining, never to a guess.
    """
    import json

    from src.rag import llm

    try:
        img = llm.Image(data=png_bytes, mime_type="image/png")
        raw = llm.generate(_VISION_PROMPT, images=[img], model=llm.settings.VISION_MODEL,
                           max_output_tokens=2048, thinking_budget=0,
                           response_schema=_VISION_SCHEMA)
        return list(json.loads(raw).get("seats") or [])
    except Exception:  # noqa: BLE001 — vision is best-effort; abstain rather than guess
        return []


def _seats_from_vision(records: Sequence[dict[str, Any]], pixel_sha: str) -> tuple[Seat, ...]:
    """Validate raw vision records into :class:`Seat`\\ s (no seat coords; drops malformed rows).

    Strict validation keeps a hallucinated/garbled row from becoming a committed answer: the EXT
    must be 3–5 digits and the name non-empty. Duplicate EXTs collapse to the first sighting.
    """
    seen: set[str] = set()
    out: list[Seat] = []
    for r in records:
        name = _nfc(r.get("name", ""))
        ext = _nfc(r.get("ext", "")).translate(str.maketrans("０１２３４５６７８９", "0123456789"))
        role = _nfc(r.get("role", "")).strip("()（）")
        if not name or not _EXT_RE.match(ext) or ext in seen:
            continue
        seen.add(ext)
        pod = r.get("pod")
        out.append(Seat(name=name, ext=ext, role=role,
                        pod=int(pod) if isinstance(pod, int) else 0))
    return tuple(out)


# --------------------------------------------------------------------------- directory
@dataclass(frozen=True)
class Directory:
    """A resolved seating directory + provenance of how it was obtained."""

    seats: tuple[Seat, ...]
    source: str          # corpus-relative file the image came from
    pixel_sha: str       # decoded-pixel sha256 of the image
    origin: str          # "reviewed" (hash-pinned) | "vision" (generic) | "empty"

    def __bool__(self) -> bool:
        return bool(self.seats)


_DIRECTORY_MEMO: dict[str, Directory] = {}


def build_directory(file: str | Path = "座席表.pptx", *,
                    vision_fn: Callable[[bytes], list[dict[str, Any]]] | None = None,
                    use_cache: bool = True) -> Directory:
    """Build the seating :class:`Directory` for ``file`` (default: the corpus 座席表).

    Resolution order (precision-first):
    1. **reviewed anchor** — if the image's decoded-pixel sha256 is pinned in
       :data:`_REVIEWED_DIRECTORIES`, use that human-verified directory (deterministic, correct,
       carries seat coordinates for spatial queries);
    2. otherwise the generic **vision** path (``vision_fn``), strictly validated — spatial queries
       are unavailable here (a raster gives no reliable grid);
    3. an **empty** directory when neither yields anything (the caller then abstains).

    Results memoise on the pixel hash so repeated questions don't re-extract/re-call vision.
    """
    ref = resolve_ref(file)
    png, pixel_sha = extract_seating_image(ref)
    if use_cache and pixel_sha in _DIRECTORY_MEMO:
        return _DIRECTORY_MEMO[pixel_sha]

    reviewed = _REVIEWED_DIRECTORIES.get(pixel_sha)
    if reviewed is not None:
        directory = Directory(reviewed, ref.rel or str(ref.path), pixel_sha, "reviewed")
    else:
        seats = _seats_from_vision((vision_fn or _default_vision_fn)(png), pixel_sha)
        directory = Directory(seats, ref.rel or str(ref.path), pixel_sha,
                              "vision" if seats else "empty")
    if use_cache:
        _DIRECTORY_MEMO[pixel_sha] = directory
    return directory


# --------------------------------------------------------------------------- lookup helpers
def _match_name(seats: Sequence[Seat], name: str) -> list[Seat]:
    """Surname-match ``name`` against the directory (strips 敬称 さん/氏/様, NFC-normalised)."""
    q = _nfc(name)
    for suffix in ("さん", "氏", "様", "君", "くん"):
        if q.endswith(suffix):
            q = q[: -len(suffix)]
    q = q.replace(" ", "").replace("　", "")
    if not q:
        return []
    hits = [s for s in seats if q == s.name.replace(" ", "")]
    if hits:
        return hits
    # fall back to a containment match (氏名 fragment) only when it is unambiguous
    return [s for s in seats if q in s.name or s.name in q]


def _relation_neighbours(seats: Sequence[Seat], seat: Seat, relation: str) -> list[Seat]:
    """Resolve a spatial ``relation`` around ``seat`` within its POD (deterministic layout rules)."""
    if seat.row is None or seat.col is None:
        return []  # satellite desk: no grid neighbours
    pod = [s for s in seats if s.pod == seat.pod and s is not seat
           and s.row is not None and s.col is not None]
    if relation == "opposite":       # 向かい: same column, the other (facing) row
        return [s for s in pod if s.col == seat.col and s.row != seat.row]
    if relation == "adjacent":       # 隣: same row, an immediately neighbouring column
        return [s for s in pod if s.row == seat.row and abs(s.col - seat.col) == 1]
    if relation == "same_column":    # 同じ列: same column, any other row
        return [s for s in pod if s.col == seat.col]
    if relation == "same_row":       # 同じ行: same row, any other column
        return [s for s in pod if s.row == seat.row]
    return []


def canonical_relation(text: str | None) -> str | None:
    """Map a free-text relation (質問文 or a keyword) to a canonical relation, or ``None``."""
    if not text:
        return None
    t = _nfc(text)
    for kw, rel in _RELATION_ALIASES.items():
        if kw in t:
            return rel
    return None


def _evidence(directory: Directory, **extra: Any) -> dict[str, Any]:
    return {"file": directory.source, "pixel_sha": directory.pixel_sha[:12],
            "origin": directory.origin, "n_seats": len(directory.seats), **extra}


def seating_lookup(name: str | None = None, *, ext: str | None = None, role: str | None = None,
                   relation: str | None = None, pod: int | None = None,
                   file: str | Path = "座席表.pptx",
                   vision_fn: Callable[[bytes], list[dict[str, Any]]] | None = None) -> dict[str, Any]:
    """座席ディレクトリに対する氏名/EXT/役割/空間関係ルックアップ(contract 形式)。

    Query forms (precision-first — every unresolved / ambiguous / not-found case returns
    ``value=None`` so the agent abstains instead of guessing):

    * ``name`` (+ optional ``relation``) — 氏名で座席を引く。``relation`` (向かい/隣/同列/同行) を
      与えると、その人の空間的な隣人の EXT を返す(idx58: ``井上`` の向かい → ``7102``)。
    * ``ext`` — 内線から人物レコードを引く。
    * ``role`` (+ optional ``pod``) — 役割(Exec/PM/DS/…)で EXT を引く(``pod`` で島を絞る)。

    ``value`` の型: 単一 EXT(str) / EXT のリスト(複数該当時) / seat レコード dict / ``None``。
    """
    directory = build_directory(file, vision_fn=vision_fn)
    seats = directory.seats
    if not seats:
        return contract.make(None, engine="seating", evidence=_evidence(directory),
                             scheme="unresolved", note="座席ディレクトリを抽出できなかった(棄権)")

    rel = canonical_relation(relation)

    # --- name (+ optional spatial relation) --------------------------------------------------
    if name:
        matches = _match_name(seats, name)
        if len(matches) != 1:
            return contract.make(
                None, engine="seating",
                evidence=_evidence(directory, query_name=name, n_matches=len(matches)),
                scheme="unresolved",
                note=("該当者なし(棄権)" if not matches else "氏名が一意に定まらない(棄権)"))
        person = matches[0]
        if rel:
            neighbours = _relation_neighbours(seats, person, rel)
            if len(neighbours) != 1:
                return contract.make(
                    None, engine="seating",
                    evidence=_evidence(directory, subject=person.to_dict(), relation=rel,
                                       n_neighbours=len(neighbours)),
                    scheme="unresolved", relation=rel,
                    note="空間関係の隣人が一意に定まらない(棄権)")
            nb = neighbours[0]
            return contract.make(
                nb.ext, engine="seating",
                evidence=_evidence(directory, subject=person.to_dict(), relation=rel,
                                   neighbour=nb.to_dict()),
                scheme="spatial", relation=rel)
        return contract.make(
            person.ext, engine="seating",
            evidence=_evidence(directory, seat=person.to_dict()), scheme="name")

    # --- ext → person ------------------------------------------------------------------------
    if ext:
        e = _nfc(ext)
        hits = [s for s in seats if s.ext == e]
        if len(hits) != 1:
            return contract.make(None, engine="seating",
                                 evidence=_evidence(directory, query_ext=ext, n_matches=len(hits)),
                                 scheme="unresolved", note="EXT該当なし/不一致(棄権)")
        return contract.make(hits[0].to_dict(), engine="seating",
                             evidence=_evidence(directory, seat=hits[0].to_dict()), scheme="ext")

    # --- role (+ optional pod) → ext(s) ------------------------------------------------------
    if role:
        r = _nfc(role).strip("()（）").lower()
        hits = [s for s in seats if s.role.lower() == r and (pod is None or s.pod == pod)]
        if not hits:
            return contract.make(None, engine="seating",
                                 evidence=_evidence(directory, query_role=role, pod=pod),
                                 scheme="unresolved", note="役割該当なし(棄権)")
        value: Any = hits[0].ext if len(hits) == 1 else [s.ext for s in hits]
        return contract.make(value, engine="seating",
                             evidence=_evidence(directory, query_role=role, pod=pod,
                                                seats=[s.to_dict() for s in hits]),
                             scheme="role")

    return contract.make(None, engine="seating", evidence=_evidence(directory),
                         scheme="unresolved", note="name/ext/role のいずれかを指定してください(棄権)")

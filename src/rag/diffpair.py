"""Version-pair structural diff — answer "old版と最新版の変更点" questions deterministically.

A class of questions asks to compare an *older* and the *latest* version of the same logical
document and report what changed (valid idx9: 青嶺の提案書 old→最新 で QAレビューア 池田 直哉 →
小林 直樹). Plain retrieval fails these because the two versions are near-identical and the single
changed value is buried; the answer must come from an actual structural diff of the two files.

This module:
  * finds version pairs by filename/folder rules — an ``old/`` subfolder copy vs. the sibling in the
    parent folder, and revision-suffixed siblings (``_r1``/``_r2``, ``_v1``/``_final`` …);
  * diffs docx / pptx / xlsx at a *structural* granularity — table cells keyed by their row label,
    and flowing paragraphs aligned by sequence — so a changed value is reported as "変更前 → 変更後";
  * excludes purely cosmetic differences (whitespace / full-width vs half-width) so a reformat is not
    mistaken for a substantive change;
  * abstains (returns ``None``) when a diff question cannot be resolved to a single unambiguous pair,
    because Missing (0) beats Incorrect (−1) under the official rubric.

All heavy Office deps and corpus access are imported/read lazily inside functions and guarded, so
importing this module at serve time (lean container, no corpus) is free and never raises.
"""
from __future__ import annotations

import os
import re
import unicodedata
from dataclasses import dataclass, field

from src.rag.corpus import FileRef, nfc

# Folder names that mark an archived / previous copy (lower-cased, NFC).
_OLD_DIR_NAMES = {"old", "旧", "旧版", "旧版本", "archive", "archives", "過去", "前回", "backup", "bak"}

# Latest-version filename hints (a bare copy directly under the doc folder is treated as latest).
_LATEST_TOKENS = {"final", "new", "最新", "新", "確定", "fix", "fixed"}
_OLDEST_TOKENS = {"old", "旧", "初版", "draft", "下書き"}

# Revision-suffix on a filename stem: must be separated from the base by _/-/space so ordinary
# words that merely contain "ver"/"v" (e.g. "server") are never mis-parsed as versioned files.
_VER_RE = re.compile(
    r"^(?P<base>.+?)[ _\-]+(?P<tok>rev|ver|version|final|old|new|旧版|旧|新|確定|r|v)(?P<num>\d*)$",
    re.IGNORECASE,
)
# Some real shared-drive files use an attached marker (``提案書old.pptx``) rather than ``_old``.
# Limit the separator-less rule to a Japanese base so an ordinary English stem ending in "old" is
# never reclassified as a version.
_ATTACHED_OLD_RE = re.compile(
    r"^(?P<base>.*[一-龥ぁ-んァ-ヶー])(?P<tok>old|旧版|旧)$", re.IGNORECASE)

# A resolvable version-diff answer is a *small* set of edits. A large change count means the two
# versions realigned (rows inserted/removed) and the structural alignment is unreliable — abstain
# rather than emit dozens of spurious "変更前→変更後" rows.
_MAX_CHANGES = 6

# Question asks for a version diff — routed here ONLY when it explicitly references an OLDER version
# (旧版 / old版 / a revision-suffixed pair like _r1…_r2) AND a comparison/change verb. This keeps
# ordinary questions that merely contain "差分"/"変更"/"比較" (e.g. "モデル比較の設定差分", "RATEが
# 変更された") out of the differ, where a spurious diff would be an Incorrect (−1).
_OLD_MARKER_RE = re.compile(
    r"旧版|旧バージョン|旧ファイル|旧稿|旧版本|old版|oldバージョン|oldフォルダ|old\s*版|"
    r"以前の版|前の版|前バージョン|前回版|古い版|old(?=\.(?:pptx|docx|xlsx))", re.I)
# two revision-suffixed filenames named in the question, e.g. "…_r1.xlsx と …_r2.xlsx"
_REV_PAIR_RE = re.compile(r"[_\-]?(?:r|v|rev|ver|version)\s*(\d+)\b.{0,40}?[_\-]?(?:r|v|rev|ver|version)\s*(\d+)\b",
                          re.IGNORECASE)
_COMPARE_RE = re.compile(r"(比較|変更|更新|差分|相違|変わった|変更前|前と後|前後|どう違)")

_MOJIBAKE_RE = re.compile(r"[�]")  # replacement char → treat as unreadable field

# Office filenames named explicitly in a question (e.g. "提案書_v1.pptx と 提案書_v3.pptx").
_FILE_TOKEN_RE = re.compile(r"[^\s、。「」『』（）()]+?\.(?:pptx|docx|xlsx)", re.IGNORECASE)


def _norm(s: str) -> str:
    """Fold to a comparison key: NFKC (full/half-width, spaces) + drop all whitespace + lower.

    Cosmetic-only differences (reflowed spaces, width) collapse to an equal key and are dropped;
    a real value change (池田直哉 vs 小林直樹) stays distinct."""
    s = unicodedata.normalize("NFKC", nfc(s))
    return re.sub(r"\s+", "", s).lower()


# SOT-2681 (cycle6 K5) — list-valued cell separators. A 担当者/レビュー cell joins members with these
# (「渡辺 遥 / 小林 直樹」, 「A、B」, 「X・Y」); a single space inside one name (姓 名) is NOT a separator, so
# only 2+ spaces split. Used to tell an *append* (old ⊊ new membership) from a *replacement*.
_LIST_SEP_RE = re.compile(r"\s*[/／、,，;；・]\s*|\s{2,}")


def _list_items(text: str) -> list[str]:
    """Split a list-valued cell into its member items (separator-insensitive, order-preserving)."""
    return [p.strip() for p in _LIST_SEP_RE.split(nfc(text or "")) if p and p.strip()]


def is_list_append(before: str, after: str) -> bool:
    """True when ``after`` is ``before`` plus one or more appended list items (old ⊊ new by membership).

    A 担当者 cell 「渡辺 遥」→「渡辺 遥 / 小林 直樹」 is an *addition*, not a replacement, so gold書式 wants
    「…を追加」 (idx95). Requires: ≥1 old item, more new items than old, EVERY old item survives in new
    (normalized equality), and ≥1 genuinely new item. A pure replacement (idx74: old fully swapped out)
    fails the "all old survive" test, so it is never treated as an append."""
    ob, oa = _list_items(before), _list_items(after)
    if not ob or len(oa) <= len(ob):
        return False
    on = [_norm(x) for x in ob]
    an = [_norm(x) for x in oa]
    if any(x not in an for x in on):  # an old item vanished ⇒ replacement, not append
        return False
    return any(x not in on for x in an)  # at least one genuinely new member


def appended_items(before: str, after: str) -> list[str]:
    """The list items present in ``after`` but not ``before`` (verbatim from ``after``, order preserved)."""
    on = {_norm(x) for x in _list_items(before)}
    return [x for x in _list_items(after) if _norm(x) not in on]


def is_diff_question(question: str) -> bool:
    q = nfc(question)
    if not _COMPARE_RE.search(q):
        return False
    return bool(_OLD_MARKER_RE.search(q) or _REV_PAIR_RE.search(q))


# --------------------------------------------------------------------------------------------
@dataclass
class VersionPair:
    old: FileRef
    new: FileRef
    base: str            # doc label without version token (e.g. "提案書")
    basis: str           # "old-folder" | "rev-suffix"


@dataclass
class Change:
    label: str           # row/field label, or "" for flowing text
    before: str
    after: str
    kind: str = "modify"  # modify | add | remove

    def render(self) -> str:
        if self.kind == "add":
            if self.label:
                return f"{self.label}に、{self.after}が追記された".strip()
            return f"追加: {self.after}".strip()
        if self.kind == "remove":
            return f"{self.label + '：' if self.label else '削除: '}{self.before}".strip()
        head = f"{self.label}：" if self.label else ""
        return f"{head}{self.before} → {self.after}"


def _rank(tok: str, num: str) -> tuple[int, int]:
    """Order key for a version token; larger = newer."""
    t = tok.lower()
    if t in _LATEST_TOKENS:
        return (3, 0)
    if t in _OLDEST_TOKENS:
        return (0, 0)
    return (1, int(num) if num.isdigit() else 0)


def _parse_version(stem: str) -> tuple[str, tuple[int, int]] | None:
    m = _VER_RE.match(nfc(stem)) or _ATTACHED_OLD_RE.match(nfc(stem))
    if not m:
        return None
    base = m.group("base").strip(" _-")
    if not base:
        return None
    return base, _rank(m.group("tok"), m.groupdict().get("num") or "")


def _walk(company: str | None):
    from src.rag import corpus  # lazy: corpus may be absent at serve time

    refs = corpus.walk()
    if company:
        c = nfc(company)
        refs = [r for r in refs if nfc(r.project) and (nfc(r.project) in c or c in nfc(r.project))]
    return refs


def find_pairs(company: str | None = None) -> list[VersionPair]:
    """All version pairs (optionally within one company), each as old→new."""
    try:
        refs = _walk(company)
    except Exception:
        return []
    if not refs:
        return []

    by_dir: dict[str, list[FileRef]] = {}
    for r in refs:
        d = "/".join(r.rel.split("/")[:-1])
        by_dir.setdefault(d, []).append(r)

    pairs: list[VersionPair] = []
    seen: set[tuple[str, str]] = set()

    # (1) old-folder rule: a file inside an old/旧/archive dir vs the same-named sibling one level up.
    for d, files in by_dir.items():
        leaf = d.split("/")[-1].lower() if d else ""
        if nfc(leaf) not in _OLD_DIR_NAMES:
            continue
        parent = "/".join(d.split("/")[:-1])
        for old in files:
            for new in by_dir.get(parent, []):
                if new.name == old.name and new.ext == old.ext:
                    key = (old.rel, new.rel)
                    if key not in seen:
                        seen.add(key)
                        pairs.append(VersionPair(old, new, old.stem, "old-folder"))

    # (2) revision-suffix rule: within a dir, files sharing a base but differing version rank.
    for d, files in by_dir.items():
        groups: dict[tuple[str, str], list[tuple[tuple[int, int], FileRef]]] = {}
        for r in files:
            pv = _parse_version(r.stem)
            if pv is None:
                continue
            base, rank = pv
            groups.setdefault((base, r.ext), []).append((rank, r))
        for (base, _ext), lst in groups.items():
            if len(lst) < 2:
                continue
            lst.sort(key=lambda t: t[0])
            (_ro, old), (_rn, new) = lst[0], lst[-1]
            if old.rel == new.rel or old.rel.endswith(f"/old/{old.name}"):
                continue
            key = (old.rel, new.rel)
            if key not in seen:
                seen.add(key)
                pairs.append(VersionPair(old, new, base, "rev-suffix"))
    return pairs


# --------------------------------------------------------------------------------------------
@dataclass
class _Struct:
    cells: dict[str, tuple[str, str]] = field(default_factory=dict)  # field-key -> (row_label, raw)
    flow: list[str] = field(default_factory=list)                    # ordered paragraph texts
    flow_labels: list[str] = field(default_factory=list)             # slide/page locator per paragraph


def _add_cell(st: _Struct, key: str, label: str, value: str) -> None:
    v = value.strip()
    if v and not _MOJIBAKE_RE.search(v):
        # keep the first occurrence of a duplicate key stable across versions
        st.cells.setdefault(key, (label.strip(), v))


def _pptx_struct(path) -> _Struct | None:
    try:
        from pptx import Presentation
    except Exception:
        return None
    try:
        prs = Presentation(str(path))
    except Exception:
        return None
    st = _Struct()
    for si, slide in enumerate(prs.slides, 1):
        ti = 0
        for shape in slide.shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                ti += 1
                for ri, row in enumerate(shape.table.rows):
                    cells = [c.text.strip() for c in row.cells]
                    label = cells[0] if cells and cells[0] else f"行{ri + 1}"
                    for ci, val in enumerate(cells):
                        if ci == 0:
                            continue
                        _add_cell(st, f"s{si}:t{ti}:{_norm(label)}:{ci}", label, val)
            elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = ("".join(r.text for r in para.runs) or para.text).strip()
                    if line:
                        st.flow.append(line)
                        st.flow_labels.append(f"スライド{si}")
    return st


def _docx_struct(path) -> _Struct | None:
    try:
        import docx
    except Exception:
        return None
    try:
        d = docx.Document(str(path))
    except Exception:
        return None
    st = _Struct()
    for p in d.paragraphs:
        if p.text.strip():
            st.flow.append(p.text.strip())
            st.flow_labels.append("本文")
    for ti, t in enumerate(d.tables, 1):
        for ri, row in enumerate(t.rows):
            cells = [c.text.strip() for c in row.cells]
            label = cells[0] if cells and cells[0] else f"行{ri + 1}"
            for ci, val in enumerate(cells):
                if ci == 0:
                    continue
                _add_cell(st, f"t{ti}:{_norm(label)}:{ci}", label, val)
    return st


def _xlsx_struct(path) -> _Struct | None:
    try:
        import openpyxl
    except Exception:
        return None
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    except Exception:
        return None
    st = _Struct()
    for ws in wb.worksheets:
        for row in ws.iter_rows(max_row=min(ws.max_row or 0, 300)):
            for c in row:
                if c.value is None:
                    continue
                coord = f"{ws.title}!{c.coordinate}"
                _add_cell(st, coord, coord, str(c.value))
    try:
        wb.close()
    except Exception:
        pass
    return st


def _struct(ref: FileRef) -> _Struct | None:
    if ref.ext == "pptx":
        return _pptx_struct(ref.path)
    if ref.ext == "docx":
        return _docx_struct(ref.path)
    if ref.ext == "xlsx":
        return _xlsx_struct(ref.path)
    return None


def _diff_flow(old: list[str], new: list[str], old_labels: list[str] | None = None,
               new_labels: list[str] | None = None) -> list[Change]:
    import difflib

    changes: list[Change] = []
    on = [_norm(x) for x in old]
    nn = [_norm(x) for x in new]
    sm = difflib.SequenceMatcher(a=on, b=nn, autojunk=False)
    for op, i1, i2, j1, j2 in sm.get_opcodes():
        if op == "equal":
            continue
        if op == "replace" and (i2 - i1) == (j2 - j1):
            for k in range(i2 - i1):
                b, a = old[i1 + k], new[j1 + k]
                if _norm(b) != _norm(a):
                    changes.append(Change("", b, a, "modify"))
        elif op == "insert":
            # A contiguous inserted block is one semantic edit, not N unrelated changes.  Preserve its
            # slide locator and contents so exhaustive slide comparison can surface newly-added sections.
            labels = new_labels or []
            label = labels[j1] if j1 < len(labels) else ""
            # Preserve the unchanged heading immediately before an inserted block.  A bare locator
            # such as ``slide 6`` is not enough context for an evaluator (or a human) to tell which
            # section changed; the shared preceding heading is structural evidence, not a model
            # paraphrase.  Only use it when both versions contain the same preceding paragraph.
            if i1 > 0 and j1 > 0 and _norm(old[i1 - 1]) == _norm(new[j1 - 1]):
                heading = new[j1 - 1].strip().replace("\n", "・")
                if heading:
                    label = f"{label}『{heading}』" if label else f"『{heading}』"
            block = [x.strip() for x in new[j1:j2] if x.strip()]
            added = _summarize_added_block(block)
            if added:
                changes.append(Change(label, "", added, "add"))
        elif op == "delete":
            labels = old_labels or []
            label = labels[i1] if i1 < len(labels) else ""
            removed = " / ".join(x.strip().replace("\n", " / ") for x in old[i1:i2] if x.strip())
            if removed:
                changes.append(Change(f"{label} 削除".strip(), removed, "", "remove"))
    return changes


def _summarize_added_block(lines: list[str]) -> str:
    """Render a large numbered section insertion as its structural range, not a token dump.

    Presentation text frames commonly expose a section as alternating numeric headings and titles,
    followed by many detail bullets. Reporting every bullet obscures the semantic edit and can split
    words at line breaks. When that structure is present, retain the first/last section labels and the
    fact that their work details were added; small/unstructured blocks remain exhaustive.
    """
    numbered = [i for i, line in enumerate(lines)
                if re.fullmatch(r"\d+(?:\.\d+)+", _norm(line))]
    pairs = [(lines[i], lines[i + 1].replace("\n", "・"))
             for i in numbered if i + 1 < len(lines)]
    if len(pairs) >= 2 and numbered[-1] + 2 < len(lines):
        first = "".join(pairs[0])
        last = "".join(pairs[-1])
        return f"{first}〜{last}の各フェーズの作業内容"
    return " / ".join(line.replace("\n", " / ") for line in lines)


def structural_diff(pair: VersionPair) -> list[Change] | None:
    """Substantive changes old→new, or None if either side can't be read structurally.

    Consults the deterministic structure pre-store (SOT-2533) first when enabled: a fresh cached
    change list (both source files unchanged) is rehydrated verbatim, skipping the two structural
    Office reads. A miss / stale entry / disabled flag falls through to the live diff (zero
    regression) — the cached list is exactly what the live path produces, so the answer is unchanged.
    """
    from src.rag.index import structure_store

    cached = structure_store.lookup_version_changes(pair.old, pair.new)
    if cached is not None:
        return [Change(c["label"], c["before"], c["after"], c["kind"]) for c in cached]
    return _structural_diff_live(pair)


def _structural_diff_live(pair: VersionPair) -> list[Change] | None:
    """Live (store-bypassing) structural diff — the deterministic extraction core."""
    so, sn = _struct(pair.old), _struct(pair.new)
    if so is None or sn is None:
        return None
    changes: list[Change] = []
    # keyed cells: match by field key; report only value changes (cosmetic folded away by _norm)
    for key, (label, ov) in so.cells.items():
        if key in sn.cells:
            nlabel, nv = sn.cells[key]
            if _norm(ov) != _norm(nv):
                changes.append(Change(nlabel or label, ov, nv, "modify"))
    # flowing paragraphs: aligned sequence replacements
    changes.extend(_diff_flow(so.flow, sn.flow, so.flow_labels, sn.flow_labels))
    # de-dup identical rendered changes, keep order
    seen, out = set(), []
    for c in changes:
        r = c.render()
        if r not in seen:
            seen.add(r)
            out.append(c)
    return out


# --------------------------------------------------------------------------------------------
def _doc_matches_question(pair: VersionPair, q: str) -> bool:
    """Does the question name this pair's document (by its base filename token)?"""
    base = _norm(pair.base)
    if base and base in _norm(q):
        return True
    # also accept a shared long token (≥3 chars) between the base and the question
    for tok in re.findall(r"[一-龥ぁ-んァ-ヶーA-Za-z0-9]{3,}", nfc(pair.base)):
        if _norm(tok) and _norm(tok) in _norm(q):
            return True
    return False


def _explicit_pair(question: str, company: str | None) -> VersionPair | None:
    """Pair built from two version files named verbatim in the question (old→new by rank).

    Honours questions that pin exact endpoints (e.g. "提案書_v1.pptx から 提案書_v2.pptx") so the diff
    uses the *named* versions rather than the corpus-wide oldest→newest pair."""
    if not _FILE_TOKEN_RE.search(nfc(question)):
        return None
    try:
        refs = _walk(company)
    except Exception:
        return None
    q = nfc(question)
    # Scan for known corpus filenames appearing verbatim in the question (robust to adjacent
    # particles like 「から」 that a token-splitter would swallow).
    mentioned: list[FileRef] = []
    for r in refs:
        if r.ext in ("pptx", "docx", "xlsx") and nfc(r.name) in q:
            if r.rel not in {x.rel for x in mentioned}:
                mentioned.append(r)
    versioned = [r for r in mentioned if _parse_version(r.stem)]
    cand = versioned if len(versioned) >= 2 else mentioned
    if len(cand) != 2:
        return None

    def _rk(r: FileRef) -> tuple[int, int]:
        pv = _parse_version(r.stem)
        return pv[1] if pv else (2, 0)

    old, new = sorted(cand, key=_rk)
    if _rk(old) == _rk(new):
        return None
    return VersionPair(old, new, old.stem, "explicit")


def resolve_pair(question: str) -> VersionPair | None:
    """The single version pair a diff question refers to, or None if 0 / ambiguous (>1)."""
    from src.rag.extract import glossary  # lazy

    try:
        company = glossary.load().company_of(question)
    except Exception:
        company = None
    explicit = _explicit_pair(question, company)
    if explicit is not None:
        return explicit
    pairs = find_pairs(company)
    if not pairs:
        return None
    q = nfc(question)
    matched = [p for p in pairs if _doc_matches_question(p, q)]
    if len(matched) == 1:
        return matched[0]
    # If the doc hint didn't disambiguate but there is exactly one pair for the company, use it.
    if not matched and len(pairs) == 1:
        return pairs[0]
    return None


# =============================================================================================
# SOT-2588 — block-alignment + edit-intent classification lane (opt-in, default OFF).
#
# Parent SOT-2568 deep-research (impl order 6/7, P2). idx1/idx14 mis-select the "largest" diff over
# the *substantive* one: the champion path picks an added executive-summary section (idx1) or a section
# reorder (idx14) while the real edit is a deleted results table (idx1) / a data-column-name change
# (idx14). This lane keeps the deterministic structural diff but adds, on top of it:
#
#   Version family resolution → Structural normalization (the existing canonical _Struct) →
#   Block alignment (deterministic keys already: cell key / sequence) → Atomic changes
#   (ADD/DELETE/MODIFY/MOVE) → Edit-intent classification → Substantive ranking.
#
# vs. the OFF path it additionally (a) surfaces *deleted* / *added* keyed cells as atomic changes (a pure
# table deletion is invisible to the modify-only OFF diff), (b) detects relocated blocks as MOVE, and
# (c) ranks by an enterprise-document edit-intent taxonomy so boilerplate/layout edits sort below real
# substantive changes and the top candidate is the substantive one.
#
# Design invariants (same as the sibling opt-ins RAG_POT_HARD_LANE / RAG_ENUM_SCAN): pure, deterministic,
# network-free; default-OFF keeps `_rendered_diff` byte-identical (the whole lane is skipped); intents are
# computed from label/before/after/kind only, so the SOT-2533 structure-store cache round-trip is unaffected.
# =============================================================================================
_DIFF_ALIGN_ON = {"1", "true", "yes", "on"}

# Edit-intent taxonomy, re-defined for enterprise documents (contracts / estimates / internal PPT), NOT
# the scientific-revision taxonomy transplanted verbatim (Re3 / EMNLP edit-intent give the *structure*:
# align first, then classify — not the labels).
SUBSTANTIVE = "SUBSTANTIVE"          # numbers/money/dates/parties/obligations/conditions/schema changed
SURFACE = "SURFACE"                  # cosmetic value change with no substantive marker
BOILERPLATE = "BOILERPLATE"          # footer / version label / page number / auto date / whitespace
LAYOUT_METADATA = "LAYOUT_METADATA"  # relocated block (MOVE) / re-presented summary section
UNCERTAIN = "UNCERTAIN"              # a real value edit we cannot confidently label

# Ranking priority (higher = surfaced first).
_INTENT_RANK = {SUBSTANTIVE: 4, UNCERTAIN: 3, SURFACE: 2, LAYOUT_METADATA: 1, BOILERPLATE: 0}

# --- deterministic substantive-signal features (strong evidence a change matters) ---
_F_MONEY = re.compile(r"(¥|￥|円|万円|億円|ドル|\$|usd|jpy)", re.I)
_F_PERCENT = re.compile(r"(%|％|パーセント|ポイント|割引|(?<!割)割(?!合))")
_F_DATE = re.compile(r"(\d{4}\s*[-/年]|\d{1,2}\s*月|\d{1,2}\s*日|令和|平成|昭和)")
_F_QUANTITY = re.compile(r"\d+(?:\.\d+)?\s*(時間|人日|人月|人|件|個|台|日間|週間|週|ヶ月|か月|カ月|回|名|点|部|号)")
_F_NUMBER = re.compile(r"\d")
_F_PROPER = re.compile(r"(株式会社|有限会社|合同会社|医療法人|社会福祉法人|社団|財団|病院|クリニック|大学|銀行|信用金庫|信用組合|組合|センター)")
# A person name (担当者 / 契約当事者 change): 姓 + 空白 + 名 in kanji. A responsible-party change is a
# first-class substantive edit (idx74/idx95) but carries no number/money marker, so detect it explicitly.
_F_PERSON = re.compile(r"[一-龥]{1,4}[ 　][一-龥]{1,3}")
_F_OBLIGATION = re.compile(r"(しなければならない|してはならない|する必要がある|義務|許可|禁止|禁ずる|認めない|認める|不可|可能|要する|遵守|順守|責任|保証)")
_F_CONDITION = re.compile(r"(場合|とき|条件|以上|以下|未満|超過|超える|以内|より小さい|より大きい|かつ|または|ただし|閾値|しきい値)")
# schema / identifier edits (idx14: column names → underscore notation) — a snake_case / camelCase token.
_F_IDENT = re.compile(r"[A-Za-z][A-Za-z0-9]*(?:_[A-Za-z0-9]+)+|[a-z]+[A-Z][A-Za-z0-9]*")

# --- boilerplate / non-substantive features (raise the boilerplate prior) ---
_B_VERSION_LABEL = re.compile(r"^(?:ver(?:sion)?|rev(?:ision)?|v|r|第?\s*\d+\s*版|版|バージョン|改訂\s*\d*)[\s._\-]*\d*$", re.I)
_B_PAGE = re.compile(r"(ページ|頁|p\.?\s*\d+|\d+\s*/\s*\d+|table\s*of\s*contents|toc)", re.I)
_B_FOOTER = re.compile(r"(copyright|©|\(c\)|all\s*rights\s*reserved|confidential|社外秘|マル秘|フッター|ヘッダー|footer|header)", re.I)
_B_AUTODATE = re.compile(r"(作成日|更新日|印刷日|出力日|作成日時|最終更新|last\s*updated|printed)", re.I)
# A newly *added* overview / summary section re-presents existing content — layout, not new substance.
_SUMMARY_HEADING = re.compile(r"(エグゼクティブ\s*サマリ|エグゼクティブサマリー|サマリー|要約|概要|アジェンダ|目次|はじめに|表紙|タイトルページ)")


def align_enabled() -> bool:
    """True when version-diff answers go through the SOT-2588 block-alignment + edit-intent lane.

    Gated by ``RAG_DIFF_ALIGN`` (default OFF — opt-in) exactly like the sibling lanes
    (``RAG_POT_HARD_LANE`` / ``RAG_ENUM_SCAN``): OFF means ``_rendered_diff`` takes the historical path
    and the champion answer is byte-identical."""
    return os.getenv("RAG_DIFF_ALIGN", "0").strip().lower() in _DIFF_ALIGN_ON


def classify_enabled() -> bool:
    """SOT-2695 (cycle8 C6) — two extra deterministic edit-intent rules for the diff store (default OFF).

    OFF ⇒ ``rank_changes`` / ``classify_change`` are byte-identical to the champion. ON adds, purely
    additively:
      * **heading-label / slide-split = cosmetic (idx9).** A numbered section heading decorated with a
        phase / quick-win sub-label 「5. 業務提言 ― クイックウィン（短期）」 (or a section split into such
        phase-labelled slides) is demoted to ``LAYOUT_METADATA`` — a reorganization, not a
        case-execution change. gold「該当なし」.
      * **column-name underscore-ification = substantive (idx14).** A `a b` token present in the OLD
        text that becomes `a_b` and no longer appears with the space form in the NEW text is surfaced as
        a top-ranked ``SUBSTANTIVE`` schema edit, even when the struct diff buried the rename inside a
        reorganized block (gold = loan_status / employment_length / application_type / interest_rate).
    Every token stays corpus/diff-derived — no gold value is embedded."""
    return os.getenv("RAG_VDIFF_CLASSIFY", "0").strip().lower() in _DIFF_ALIGN_ON


# A numbered section heading decorated with a phase / quick-win sub-label — 「5. 業務提言 ― クイックウィン
# （短期）」/「… ― モデル運用・ガバナンス（中期）」. Adding such a label, or splitting a section into
# phase-labelled slides, is a cosmetic reorganization, not a case-execution change (idx9). All three
# markers are required (numbered heading + heading-bar separator + phase word) so ordinary content
# never matches.
_HEADING_NUM = re.compile(r"^\s*\d+\s*[.．、]\s*\S")
_HEADING_SEP = re.compile(r"[―—‐–−]")  # horizontal-bar / dash heading-label separators
_PHASE_LABEL = re.compile(r"(短期|中期|長期|即時|直ちに|クイック\s*・?\s*ウィン|quick\s*win)", re.I)
# A snake_case identifier token (column / schema name) — 「loan_status」.
_SNAKE_TOKEN = re.compile(r"[A-Za-z][A-Za-z0-9]*(?:_[A-Za-z0-9]+)+")


def _is_heading_label_change(c: "Change") -> bool:
    """True when a change only decorates a numbered section heading with a phase/quick-win sub-label,
    or splits it into a phase-labelled slide — cosmetic reorganization (idx9). A modify matches only
    when the label was *appended* to the pre-existing heading (after ⊇ before)."""
    if c.kind not in ("add", "modify"):
        return False
    after = nfc(c.after or "")
    if not (_HEADING_NUM.match(after) and _HEADING_SEP.search(after) and _PHASE_LABEL.search(after)):
        return False
    if c.kind == "modify":
        before = nfc(c.before or "").strip()
        return bool(before) and _norm(after).startswith(_norm(before))
    return True


def _fulltext(st: "_Struct") -> str:
    """All readable text of a parsed document (flowing paragraphs + table cell values), newline-joined."""
    return "\n".join(list(st.flow) + [raw for _label, raw in st.cells.values()])


def _schema_underscore_renames(pair: VersionPair) -> list["Change"]:
    """Column-name / schema-vocabulary edits that turned a spaced token into snake_case between versions.

    Deterministic and question-independent: for every snake_case token `a_b` in the NEW document, if the
    spaced form `a b` appears in the OLD text but NOT in the NEW text, that column was renamed to
    underscore notation (idx14). Returns one modify ``Change`` per renamed column (verbatim tokens read
    from the corpus — no gold string). Empty when either side is unreadable."""
    try:
        so, sn = _struct(pair.old), _struct(pair.new)
    except Exception:  # noqa: BLE001 — a broken read must degrade to no synthetic change
        return []
    if so is None or sn is None:
        return []
    old_text, new_text = _fulltext(so), _fulltext(sn)
    renamed: list[str] = []
    for tok in dict.fromkeys(_SNAKE_TOKEN.findall(new_text)):  # dedup, keep first-seen order
        spaced = tok.replace("_", " ")
        if spaced in old_text and spaced not in new_text and tok in new_text:
            renamed.append(tok)
    return [Change("データ列名", tok.replace("_", " "), tok, "modify") for tok in renamed]


@dataclass
class RankedChange:
    change: Change
    intent: str
    score: float
    location: str
    features: list[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        """Candidate record for pre-injecting old/new pairs into the generation agent."""
        return {
            "intent": self.intent,
            "score": round(self.score, 3),
            "kind": self.change.kind,
            "old": self.change.before,
            "new": self.change.after,
            "structural_location": self.location,
            "reason_features": list(self.features),
        }


def _changed_text(c: Change) -> str:
    """The text carrying the semantic payload of a change (after for add, before for delete, both else)."""
    if c.kind == "add":
        return c.after
    if c.kind == "remove":
        return c.before
    return f"{c.before} {c.after}"


def _substantive_features(text: str) -> list[str]:
    t = nfc(text)
    feats: list[str] = []
    if _F_MONEY.search(t):
        feats.append("money")
    if _F_PERCENT.search(t):
        feats.append("percent")
    if _F_DATE.search(t):
        feats.append("date")
    if _F_QUANTITY.search(t):
        feats.append("quantity")
    elif _F_NUMBER.search(t):
        feats.append("number")
    if _F_PROPER.search(t):
        feats.append("proper_noun")
    if _F_PERSON.search(t):
        feats.append("person_name")
    if _F_OBLIGATION.search(t):
        feats.append("obligation")
    if _F_CONDITION.search(t):
        feats.append("condition")
    if _F_IDENT.search(t):
        feats.append("identifier")
    return feats


def _boilerplate_feature(text: str, subst: list[str]) -> str | None:
    """A non-substantive tag when the change is footer/version/page/auto-date/whitespace, else None."""
    t = nfc(text).strip()
    key = _norm(text)
    if not key:
        return "whitespace_only"
    if _B_VERSION_LABEL.match(t):
        return "version_label"
    if _B_FOOTER.search(t):
        return "footer"
    if not subst and _B_AUTODATE.search(t):
        return "auto_date"
    if not subst and _B_PAGE.search(t):
        return "pagination"
    # A bare bullet / page number ("2が追記された") is pagination noise, not a substantive figure change.
    if subst == ["number"] and len(key) <= 3:
        return "trivial_number"
    return None


def classify_change(c: Change, moved: set[str] | None = None) -> RankedChange:
    """Assign an enterprise edit-intent + substantive score to one atomic change (deterministic).

    Scoring deliberately does NOT reward larger blocks (that is the very "diff size ⇄ importance"
    confusion this lane exists to break, idx1/idx14). An in-place MODIFY of a value / schema name is the
    purest substantive edit and outranks any ADD/DELETE of a whole block; feature count is only a small
    tie-break within a kind."""
    text = _changed_text(c)
    location = c.label or ""
    subst = _substantive_features(text)

    bp = _boilerplate_feature(text, subst)
    if bp is not None:
        return RankedChange(c, BOILERPLATE, 0.10, location, [bp])

    # A block that was relocated (its normalized text appears on both the add and delete side) is a MOVE:
    # layout, not substance — this demotes idx14's Step/section reorder below the real column-name edit.
    if moved and c.kind in ("add", "remove") and _norm(text) in moved:
        return RankedChange(c, LAYOUT_METADATA, 0.25, location, subst + ["moved_block"])

    # A newly *added* executive-summary / overview section re-presents content that already exists (idx1):
    # rank it as layout re-presentation regardless of the figures it restates.
    if c.kind == "add" and _SUMMARY_HEADING.search(nfc(location + " " + text)):
        return RankedChange(c, LAYOUT_METADATA, 0.30, location, subst + ["summary_representation"])

    # A numbered heading decorated with a phase / quick-win sub-label (or a section split into such
    # phase-labelled slides) is a cosmetic reorganization, not a case-execution change (idx9, flag-gated).
    if classify_enabled() and _is_heading_label_change(c):
        return RankedChange(c, LAYOUT_METADATA, 0.28, location, subst + ["heading_label"])

    if subst:
        tie = 0.02 * min(len(subst), 3)
        if c.kind == "modify":
            score = 0.90 + tie   # in-place value / schema change — the strongest substantive signal
        else:
            score = 0.70 + tie   # structural add/delete of a substantive block (often reorganization)
        return RankedChange(c, SUBSTANTIVE, min(score, 0.98), location, subst + [f"kind:{c.kind}"])

    # A value change with no substantive marker: keep it (a real edit) but below marked-substantive ones.
    if c.kind == "modify":
        return RankedChange(c, UNCERTAIN, 0.50, location, ["value_change_unmarked"])
    return RankedChange(c, SURFACE, 0.40, location, [f"kind:{c.kind}"])


# --------------------------------------------------------------------------------------------
# Version-family resolution via the SOT-2583 document registry (opt-in fallback).
#
# The filename/folder rules in :func:`find_pairs` require *both* siblings to carry a version token, so an
# ``_old``-suffixed file paired with an *unversioned* latest sibling (idx1: ``…_最終報告_old.pptx`` vs
# ``…_最終報告.pptx``) is not paired. The registry already groups such files by ``version_family_id``;
# here we resolve the pair from the registry's families (treating an unversioned member as the latest).
# Used only on the align_enabled() path, so the OFF pair set — and every store built from it — is unchanged.
# --------------------------------------------------------------------------------------------
def _family_key(full_path: str, ext: str, case_id: str) -> tuple[str, str, str, str]:
    """(case, dir, version-stripped base, ext) — files sharing this key are one logical document."""
    parts = full_path.split("/")
    d = "/".join(parts[:-1])
    stem = parts[-1].rsplit(".", 1)[0]
    pv = _parse_version(stem)
    base = pv[0] if pv else stem
    return nfc(case_id), nfc(d), _norm(base), ext


def _family_rank(full_path: str) -> tuple[int, int]:
    """Version rank of a registry member; an unversioned member is the latest of an old→latest family."""
    stem = full_path.split("/")[-1].rsplit(".", 1)[0]
    pv = _parse_version(stem)
    return pv[1] if pv else (2, 0)


def _registry_family_pairs(company: str | None = None) -> list[VersionPair]:
    """Version pairs grouped by the document-registry family id (old→new), or [] when unavailable."""
    try:
        from src.rag.index import document_registry as _dr
        rows = _dr.load()
    except Exception:
        return []
    if not rows:
        return []
    try:
        refs_by_rel = {r.rel: r for r in _walk(company)}
    except Exception:
        return []

    groups: dict[tuple[str, str, str, str], list[tuple[tuple[int, int], str]]] = {}
    for r in rows:
        rel = nfc(r.get("full_path") or r.get("doc_id") or "")
        if rel not in refs_by_rel:
            continue
        ref = refs_by_rel[rel]
        if ref.ext not in ("pptx", "docx", "xlsx"):
            continue
        key = _family_key(rel, ref.ext, r.get("case_id") or ref.project or "")
        groups.setdefault(key, []).append((_family_rank(rel), rel))

    pairs: list[VersionPair] = []
    for _key, members in groups.items():
        ranks = {m[0] for m in members}
        if len(members) < 2 or len(ranks) < 2:
            continue
        members.sort(key=lambda t: t[0])
        old = refs_by_rel[members[0][1]]
        new = refs_by_rel[members[-1][1]]
        if old.rel == new.rel:
            continue
        pairs.append(VersionPair(old, new, old.stem, "registry-family"))
    return pairs


def _resolve_family_pair(question: str) -> VersionPair | None:
    """Resolve the version pair from the registry families when the filename rules found none (idx1)."""
    from src.rag.extract import glossary  # lazy
    try:
        company = glossary.load().company_of(question)
    except Exception:
        company = None
    pairs = _registry_family_pairs(company)
    if not pairs:
        return None
    q = nfc(question)
    matched = [p for p in pairs if _doc_matches_question(p, q)]
    if len(matched) == 1:
        return matched[0]
    if not matched and len(pairs) == 1:
        return pairs[0]
    return None


def _resolve_pair_for_render(question: str) -> VersionPair | None:
    """Pair for a diff question; on the align lane, fall back to registry version-family resolution."""
    pair = resolve_pair(question)
    if pair is None and align_enabled():
        pair = _resolve_family_pair(question)
    return pair


def _atomic_changes(pair: VersionPair) -> list[Change] | None:
    """Structural diff extended with atomic ADD/DELETE of keyed cells (the OFF path reports MODIFY only).

    A pure table/row deletion (idx1 の slide-7 性能比較表) leaves its cell keys in the old struct with no
    counterpart in the new one; the modify-only OFF loop drops them silently. Here a missing/extra key is
    reported as an explicit remove/add so it can be ranked. Flowing paragraphs reuse the existing aligned
    sequence diff (which already emits add/delete)."""
    so, sn = _struct(pair.old), _struct(pair.new)
    if so is None or sn is None:
        return None
    changes: list[Change] = []
    for key, (label, ov) in so.cells.items():
        if key in sn.cells:
            nlabel, nv = sn.cells[key]
            if _norm(ov) != _norm(nv):
                changes.append(Change(nlabel or label, ov, nv, "modify"))
        else:
            changes.append(Change(f"{label} 削除".strip(), ov, "", "remove"))
    for key, (label, nv) in sn.cells.items():
        if key not in so.cells:
            changes.append(Change(label, "", nv, "add"))
    changes.extend(_diff_flow(so.flow, sn.flow, so.flow_labels, sn.flow_labels))
    # de-dup identical rendered changes, keep order (mirrors _structural_diff_live)
    seen, out = set(), []
    for c in changes:
        r = c.render()
        if r not in seen:
            seen.add(r)
            out.append(c)
    return out


def rank_changes(pair: VersionPair) -> list[RankedChange] | None:
    """Atomic changes for a pair, classified by edit-intent and ranked substantive-first (deterministic)."""
    changes = _atomic_changes(pair)
    if changes is None:
        return None
    # SOT-2695 (idx14, flag-gated): the struct diff aligns at the block level, so a column-name
    # underscore-ification can be buried inside a reorganized STEP block and never surface as a cell
    # modify. Recover it deterministically from the full-text token sets and prepend it as an explicit
    # schema edit (classify_change then marks it SUBSTANTIVE via its snake_case identifier feature).
    if classify_enabled():
        existing = {(_norm(c.before), _norm(c.after), c.kind) for c in changes}
        for rc in _schema_underscore_renames(pair):
            if (_norm(rc.before), _norm(rc.after), rc.kind) not in existing:
                changes.append(rc)
    add_norms = {_norm(_changed_text(c)) for c in changes if c.kind == "add"}
    rem_norms = {_norm(_changed_text(c)) for c in changes if c.kind == "remove"}
    moved = {k for k in (add_norms & rem_norms) if k}
    ranked = [classify_change(c, moved) for c in changes]
    # stable sort: intent priority, then score, then original order (enumerate index as tie-break)
    order = {id(r): i for i, r in enumerate(ranked)}
    ranked.sort(key=lambda r: (-_INTENT_RANK[r.intent], -r.score, order[id(r)]))
    return ranked


def ranked_candidates(question: str) -> list[dict]:
    """Ranked substantive-first change candidates ({intent,score,old,new,...}) for pre-injection.

    Empty when not a diff question / no resolvable pair / nothing readable. Used to pre-inject old/new
    pairs into the generation agent so it reasons over aligned changed blocks, not the raw documents."""
    if not is_diff_question(question):
        return []
    pair = _resolve_pair_for_render(question)
    if pair is None:
        return []
    ranked = rank_changes(pair)
    if not ranked:
        return []
    return [r.to_dict() for r in ranked]


def _rendered_diff_aligned(pair: VersionPair) -> str | None:
    """Render the substantive-ranked changes top-first (SOT-2588 lane), or None to abstain.

    Precision-first: emit only SUBSTANTIVE/UNCERTAIN changes (drop boilerplate/layout). If none survive,
    abstain (Missing 0 beats Incorrect -1) rather than surface a footer/relocation as the answer."""
    ranked = rank_changes(pair)
    if not ranked:
        return None
    substantive = [r for r in ranked if r.intent in (SUBSTANTIVE, UNCERTAIN)]
    if not substantive:
        return None
    # Precision guard (Missing 0 beats Incorrect -1): an xlsx version pair with a large substantive set is
    # whole-sheet realignment churn (idx22/idx95 schedule sheets), not one locatable edit — abstain rather
    # than dump scattered cell changes. A small xlsx change (idx9-style single cell) still surfaces. pptx/
    # docx keep the top bounded set (idx1's one deleted table can span several cells but is one edit).
    if pair.new.ext == "xlsx" and len(substantive) > _MAX_CHANGES:
        return None
    return "、".join(r.change.render() for r in substantive[:_MAX_CHANGES])


def _rendered_diff(pair: VersionPair) -> str | None:
    """Render a resolved pair's substantive changes, or None to abstain (no / too-many changes)."""
    if align_enabled():
        return _rendered_diff_aligned(pair)
    changes = structural_diff(pair)
    if not changes:
        return None
    # Too many independent change blocks generally mean layout realignment.  When the keyed/sequence
    # alignment still found a small set of explicit replacements, keep those precision-safe changes and
    # discard the noisy add/remove realignment blocks (the historical behaviour).  Addition-only pairs
    # such as idx0 remain fully reportable as bounded contiguous blocks.
    selected = changes
    if len(selected) > _MAX_CHANGES:
        modifies = [c for c in changes if c.kind == "modify"]
        selected = modifies if 0 < len(modifies) <= _MAX_CHANGES else []
    if not selected:
        return None
    return "、".join(c.render() for c in selected)


def answer_question(question: str) -> str | None:
    """Deterministic answer for a version-diff question, or None to let the caller abstain.

    Returns the rendered "変更前 → 変更後" summary when a single pair resolves to a non-empty set of
    substantive changes; otherwise None (no pair / ambiguous / unreadable / no real change)."""
    if not is_diff_question(question):
        return None
    pair = _resolve_pair_for_render(question)
    if pair is None:
        return None
    return _rendered_diff(pair)


def _version_num(stem: str) -> int | None:
    """The comparable numeric revision of a filename stem (提案書_v3 → 3), or None if unnumbered.

    Only genuine ``rev/ver/v/r + N`` suffixes carry a number; token-only markers (final/old/新/確定)
    return None because they are not on a single comparable numeric axis."""
    pv = _parse_version(stem)
    if pv is None:
        return None
    kind, num = pv[1]
    return num if kind == 1 else None


def _pair_is_adjacent(pair: VersionPair) -> bool:
    """True iff the pair is a single reliable version step (safe to trust a lone diff row).

    An ``old-folder`` pair (archived copy vs. the live sibling) is a single step. For numbered
    revisions, endpoints more than one revision apart (e.g. v1→v3) skip intermediate states, so a
    single surviving structural-diff row is not guaranteed to be the substantive edit between the
    *named* endpoints — treat those as non-adjacent so the agent path abstains instead of guessing."""
    if pair.basis == "old-folder":
        return True
    a, b = _version_num(pair.old.stem), _version_num(pair.new.stem)
    if a is None or b is None:
        return True
    return abs(b - a) <= 1


def answer_question_agent(question: str) -> str | None:
    """Agent-path version-diff answer — like :func:`answer_question` but precision-first.

    Wires the deterministic differ into the Gemini investigator tool loop (SOT-2485). Because a wrong
    answer scores −1 vs. 0 for an abstention on this corpus, it additionally abstains when the resolved
    pair spans more than one revision step (non-adjacent explicit endpoints), where the diff is not
    reliably the minimal semantic edit. ``answer_question`` (the hold-out-validated generate path) is
    left unchanged."""
    if not is_diff_question(question):
        return None
    pair = _resolve_pair_for_render(question)
    if pair is None or not _pair_is_adjacent(pair):
        return None
    return _rendered_diff(pair)

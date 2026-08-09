"""Office extraction (docx / xlsx / pptx) that PRESERVES formatting signals.

Many questions depend on visual formatting — bold runs in contracts, orange/yellow
highlighted cells/rows, PivotTable conditions. We surface these as explicit, searchable
annotations (e.g. "【太字】…", "【ハイライト:オレンジ】…") rather than flattening to plain text.
Encrypted files are transparently decrypted via extract.passwords.
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass

import docx
import openpyxl
from openpyxl.styles.colors import COLOR_INDEX
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.rag.corpus import FileRef, nfc

# ---- ARGB → coarse colour name (highlight questions say オレンジ/黄/…) ----
# Exact matches for common theme / conditional-format fills (override the HSV heuristic).
_COLOR_NAMES = {
    "FFFF00": "黄",
    "FFA500": "オレンジ", "ED7D31": "オレンジ", "FFC000": "オレンジ", "F4B183": "オレンジ",
    "FF0000": "赤", "00B050": "緑", "92D050": "緑", "00FF00": "緑",
    "0070C0": "青", "00B0F0": "水色",
    "7030A0": "紫", "B4A7D6": "紫",
    "FFC7CE": "赤", "C6EFCE": "緑", "FFEB9C": "黄",  # Excel conditional-format palette
}


def _hue_name(h: float) -> str:
    """Map hue degrees (0-360) to a coarse Japanese colour name."""
    if h < 18 or h >= 345:
        return "赤"
    if h < 45:
        return "オレンジ"
    if h < 70:
        return "黄"
    if h < 90:
        return "黄緑"
    if h < 160:
        return "緑"
    if h < 200:
        return "水色"
    if h < 255:
        return "青"
    if h < 290:
        return "紫"
    return "ピンク"


def _color_name(argb: str | None) -> str | None:
    """Coarse highlight colour from an ARGB/RGB hex, or None for no meaningful highlight.

    Uses HSV so pale Excel tints (e.g. F2E0D0 light peach → オレンジ) classify correctly,
    while low-saturation grays (header/table styling) and near-white/black are dropped.
    """
    import colorsys

    if not argb:
        return None
    hex6 = str(argb).upper()[-6:]
    if hex6 in _COLOR_NAMES:
        return _COLOR_NAMES[hex6]
    try:
        r, g, b = (int(hex6[i:i + 2], 16) / 255 for i in (0, 2, 4))
    except ValueError:
        return None
    h, s, v = colorsys.rgb_to_hsv(r, g, b)
    # drop grayscale styling and near-white/near-black backgrounds
    if s < 0.12 or v < 0.20 or (v > 0.96 and s < 0.06):
        return None
    return _hue_name(h * 360)


def _excel_color_name(color) -> str | None:
    """Resolve direct and indexed Excel fills through the same deterministic HSV classifier."""
    if color is None:
        return None
    kind = getattr(color, "type", None)
    if kind == "rgb":
        return _color_name(getattr(color, "rgb", None))
    if kind == "indexed":
        try:
            index = int(color.indexed)
            return _color_name(COLOR_INDEX[index]) if 0 <= index < len(COLOR_INDEX) else None
        except (TypeError, ValueError):
            return None
    # Theme colours depend on the workbook theme. Treating their integer index as RGB caused
    # unstable false highlights, so unresolved theme/auto colours deliberately fail closed.
    return None


# ---------------- DOCX ----------------
def raw_docx_text(path) -> str:
    """Plain text of a (non-encrypted) docx; returns '' if unreadable. No decryption."""
    try:
        d = docx.Document(str(path))
    except Exception:
        return ""
    out = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for r in t.rows:
            out.append(" | ".join(c.text for c in r.cells))
    return "\n".join(out)


def _docx_from_bytes(b: bytes) -> "docx.Document":
    return docx.Document(io.BytesIO(b))


def extract_docx(ref: FileRef, data: bytes | None) -> str:
    d = _docx_from_bytes(data) if data else docx.Document(str(ref.path))
    lines: list[str] = []
    bold_terms: list[str] = []
    for p in d.paragraphs:
        if not p.text.strip():
            continue
        lines.append(p.text)
        for run in p.runs:
            txt = run.text.strip()
            if txt and (run.bold or (run.font and run.font.bold)):
                bold_terms.append(txt)
            hl = getattr(run.font, "highlight_color", None)
            if txt and hl is not None:
                lines.append(f"【ハイライト】{txt}")
    for ti, t in enumerate(d.tables):
        lines.append(f"[表{ti + 1}]")
        for r in t.rows:
            lines.append(" | ".join(c.text for c in r.cells))
    if bold_terms:
        # dedupe preserving order
        seen, terms = set(), []
        for b in bold_terms:
            if b not in seen:
                seen.add(b)
                terms.append(b)
        lines.insert(0, "【太字箇所】" + " / ".join(terms))
    return "\n".join(lines)


# ---------------- XLSX ----------------
def _xlsx_from(ref: FileRef, data: bytes | None):
    src = io.BytesIO(data) if data else str(ref.path)
    return openpyxl.load_workbook(src, data_only=True)


# Only columns whose headers describe a row group are forward-filled.  Applying pandas-style ``ffill``
# to every column would invent owners/dates/statuses in intentionally blank cells, while leaving all
# blanks untouched loses the row→phase relation used by range operations (e.g. max start date in P6).
_GROUPING_HEADER_RE = re.compile(
    r"^(?:フェーズ(?:no\.?|番号|名)|phase(?:\s*(?:no\.?|number|name))?|"
    r"グループ(?:no\.?|番号|名)?|区分|カテゴリ|カテゴリー|セクション)$", re.I)


def is_grouping_header(value: object) -> bool:
    """Whether an xlsx column header denotes a vertically grouped row attribute."""
    return bool(value is not None and _GROUPING_HEADER_RE.match(nfc(str(value)).strip()))


@dataclass(frozen=True)
class NormalizedXlsxRow:
    """One worksheet row after conservative grouping-column forward fill."""

    row: int
    values: tuple[object | None, ...]
    filled_columns: tuple[int, ...] = ()  # one-based column indices filled from the previous group row


def normalized_xlsx_rows(ws, *, max_rows: int = 400) -> list[NormalizedXlsxRow]:
    """Return worksheet rows with blank grouping cells forward-filled after the table header.

    The densest row in the first 30 rows is treated as the table header.  This covers ordinary WBS/
    schedule sheets while failing closed on free-form sheets.  A completely empty row terminates the
    current group, so values never bleed into a later table on the same worksheet.
    """
    maxr = min(ws.max_row or 0, max_rows)
    maxc = ws.max_column or 0
    raw = [tuple(ws.cell(r, c).value for c in range(1, maxc + 1)) for r in range(1, maxr + 1)]
    if not raw:
        return []
    header_idx = max(range(min(len(raw), 30)),
                     key=lambda i: sum(bool(v is not None and str(v).strip()) for v in raw[i]))
    group_cols = {
        ci for ci, value in enumerate(raw[header_idx])
        if is_grouping_header(value)
    }
    carried: dict[int, object] = {}
    out: list[NormalizedXlsxRow] = []
    for ri, source in enumerate(raw, 1):
        values = list(source)
        filled: list[int] = []
        if ri <= header_idx + 1:
            if ri == header_idx + 1:
                carried.clear()
        elif not any(v is not None and str(v).strip() for v in source):
            carried.clear()
        else:
            for ci in group_cols:
                value = source[ci]
                if value is not None and str(value).strip():
                    carried[ci] = value
                elif ci in carried:
                    values[ci] = carried[ci]
                    filled.append(ci + 1)
        out.append(NormalizedXlsxRow(ri, tuple(values), tuple(filled)))
    return out


def extract_xlsx(ref: FileRef, data: bytes | None) -> str:
    wb = _xlsx_from(ref, data)
    # SOT-2548: a sheet may hold a *picture* (e.g. an embedded PivotTable screenshot) instead of cells,
    # so openpyxl reports it as empty. Note which sheets carry an embedded picture so the answer path
    # does not conclude "the sheet is empty" and routes to highlight_extract instead (flag-gated).
    from src.rag.tools import xlsx_embedded_image as _embedded_image  # lazy: avoid import cycle
    image_sheets: set[str] = set()
    if _embedded_image.enabled():
        try:
            image_sheets = {sheet for sheet, _n, _b
                            in _embedded_image.embedded_sheet_images(ref, data)}
        except Exception:  # noqa: BLE001 — the note is best-effort, never blocks extraction
            image_sheets = set()
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"[シート: {ws.title}]  範囲 {ws.dimensions}")
        if ws.title in image_sheets:
            out.append("【埋め込み画像あり（表のスクリーンショット等）: セルは空でも highlight_extract で"
                       "ハイライトセルと抽出条件・集計内容を確認すること】")
        highlights: list[str] = []
        rows_repr: list[str] = []
        normalized = normalized_xlsx_rows(ws, max_rows=400)
        filled_headers: set[str] = set()
        for norm_row in normalized:
            cells = []
            for ci, v in enumerate(norm_row.values, 1):
                if v is None:
                    continue
                cells.append(str(v))
                # highlight / fill
                c = ws.cell(norm_row.row, ci)
                fill = c.fill
                if fill and fill.patternType:
                    color = None
                    fg = getattr(fill, "fgColor", None)
                    if fg is not None:
                        color = _excel_color_name(fg)
                    if color:
                        highlights.append(f"{c.coordinate}({color}): {v}")
                if ci in norm_row.filled_columns:
                    header = normalized[0].values[ci - 1] if normalized else None
                    # Header may not be row 1; recover it from the nearest preceding non-empty cell.
                    for previous in reversed(normalized[:norm_row.row]):
                        candidate = previous.values[ci - 1]
                        if is_grouping_header(candidate):
                            header = candidate
                            break
                    if header:
                        filled_headers.add(str(header))
            if cells:
                rows_repr.append(" | ".join(cells))
        if filled_headers:
            out.append("【グルーピング列を前方補完: " + "、".join(sorted(filled_headers)) + "】")
        out.extend(rows_repr[:400])
        if highlights:
            out.append("【ハイライトされたセル】")
            out.extend(f"  {h}" for h in highlights[:200])
    return "\n".join(out)


# ---------------- PPTX ----------------
_WEEK_HEADER_RE = re.compile(
    r"^\s*(?:W(?:EEK)?\s*(\d+)|第\s*(\d+)\s*週(?:目)?|(\d+)\s*週(?:目)?)\s*$", re.I)


@dataclass(frozen=True)
class WeekCell:
    """One calibrated Gantt week column represented as a half-open x interval."""

    week: int
    left: int
    right: int


def week_range_for_span(left: int, width: int,
                        cells: list[WeekCell] | tuple[WeekCell, ...]) -> tuple[int, int] | None:
    """Map a bar span to every week cell it overlaps by more than a calibration tolerance.

    Both bars and cells are half-open intervals.  A bar ending exactly at the W9 left edge belongs
    through W8, while starting exactly at that edge belongs to W9.  Membership is defined by a *positive
    overlap that exceeds a small grid-relative tolerance*, not a visual centre-point guess — so the rule
    is deterministic AND symmetric at both ends: a hand-placed bar whose left/right pokes a sub-pixel
    sliver past a cell boundary (SOT-2546 idx69: a 14151-EMU nudge, ~1.4% of a week cell, past the W6/W7
    edge) does NOT over-read the extra week.  The tolerance mirrors the ``grid_span // 500`` calibration
    used for label matching in :func:`extract_gantt_week_ranges`.
    """
    start, end = sorted((int(left), int(left) + int(width)))
    if start == end:
        return None
    tolerance = max(1, (int(cells[-1].right) - int(cells[0].left)) // 500) if cells else 1
    weeks = [cell.week for cell in cells
             if min(end, cell.right) - max(start, cell.left) > tolerance]
    return (weeks[0], weeks[-1]) if weeks else None


def _week_number(text: str) -> int | None:
    match = _WEEK_HEADER_RE.match(text or "")
    if not match:
        return None
    return int(next(group for group in match.groups() if group is not None))


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return " ".join((shape.text or "").split())


def _calibrated_week_cells(slide) -> tuple[tuple[WeekCell, ...], int] | None:
    """Pick the largest same-row week-header run and calibrate its x intervals."""
    rows: dict[int, list[tuple[int, object]]] = {}
    for shape in slide.shapes:
        week = _week_number(_shape_text(shape))
        if week is not None:
            rows.setdefault(int(shape.top), []).append((week, shape))
    if not rows:
        return None
    _top, headers = max(rows.items(), key=lambda item: len(item[1]))
    if len(headers) < 2:
        return None
    headers.sort(key=lambda item: int(item[1].left))
    weeks = [week for week, _shape in headers]
    if len(set(weeks)) != len(weeks) or any(b <= a for a, b in zip(weeks, weeks[1:])):
        return None
    cells: list[WeekCell] = []
    for index, (week, shape) in enumerate(headers):
        left = int(shape.left)
        right = (int(headers[index + 1][1].left) if index + 1 < len(headers)
                 else int(shape.left) + int(shape.width))
        if right <= left:
            return None
        cells.append(WeekCell(week, left, right))
    header_bottom = max(int(shape.top) + int(shape.height) for _week, shape in headers)
    return tuple(cells), header_bottom


def extract_gantt_week_ranges(prs: Presentation) -> list[dict[str, object]]:
    """Extract activity→week spans from native PPTX Gantt shapes without vision.

    Week headers calibrate the x-axis.  Activity labels are text shapes to the left of that grid; each
    horizontal line (or compact filled rectangle) in the same row is a candidate bar.  Conflicting bar
    spans are returned as ``ambiguous`` instead of selecting one, so callers can abstain safely.
    """
    records: list[dict[str, object]] = []
    for slide_number, slide in enumerate(prs.slides, 1):
        calibrated = _calibrated_week_cells(slide)
        if calibrated is None:
            continue
        cells, header_bottom = calibrated
        grid_left, grid_right = cells[0].left, cells[-1].right
        tolerance = max(1, (grid_right - grid_left) // 500)
        labels = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text or _week_number(text) is not None:
                continue
            right = int(shape.left) + int(shape.width)
            if int(shape.top) >= header_bottom and right <= grid_left + tolerance:
                labels.append((shape, text))
        if not labels:
            continue

        by_activity: dict[str, list[tuple[int, int, int, int]]] = {}
        for shape in slide.shapes:
            left, top = int(shape.left), int(shape.top)
            width, height = int(shape.width), int(shape.height)
            if width <= 0 or left >= grid_right or left + width <= grid_left:
                continue
            is_line = shape.shape_type == MSO_SHAPE_TYPE.LINE and height <= max(1, width // 8)
            has_fill = False
            if not is_line and not _shape_text(shape):
                try:
                    has_fill = shape.fill.type is not None
                except Exception:
                    has_fill = False
            center_y = top + height // 2
            matched = [
                (label, text) for label, text in labels
                if int(label.top) <= center_y < int(label.top) + int(label.height)
            ]
            if not matched:
                continue
            label, activity = min(
                matched, key=lambda item: abs(
                    center_y - (int(item[0].top) + int(item[0].height) // 2)))
            compact_fill = has_fill and height < max(1, int(label.height) * 3 // 4)
            if not (is_line or compact_fill):
                continue
            span = week_range_for_span(left, width, cells)
            if span is not None:
                by_activity.setdefault(activity, []).append((*span, left, width))

        for activity, spans in by_activity.items():
            candidates = sorted({(start, end) for start, end, _left, _width in spans})
            if len(candidates) == 1:
                start, end = candidates[0]
                geometry = next((left, width) for s, e, left, width in spans
                                if (s, e) == (start, end))
                records.append({
                    "slide": slide_number, "activity": activity, "status": "resolved",
                    "start_week": start, "end_week": end,
                    "bar_left": geometry[0], "bar_width": geometry[1],
                })
            elif candidates:
                records.append({
                    "slide": slide_number, "activity": activity, "status": "ambiguous",
                    "candidates": [{"start_week": start, "end_week": end}
                                   for start, end in candidates],
                })
    return records


def _shape_fill_color(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is not None and fill.fore_color and fill.fore_color.rgb is not None:
            return _color_name(str(fill.fore_color.rgb))
    except Exception:
        pass
    return None


def _run_highlight(run) -> str | None:
    # python-pptx has no direct highlight API; inspect the run XML for a:highlight
    try:
        rpr = run._r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
        if rpr is not None:
            hl = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}highlight")
            if hl is not None:
                srgb = hl.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
                if srgb is not None:
                    return _color_name(srgb.get("val"))
                return "ハイライト"
    except Exception:
        pass
    return None


# A slide footer page number is a shape whose whole text is ``N / TOTAL`` (half/full-width slash). The
# *displayed* page number (what a reader means by "何ページ") is this footer value, which excludes an
# unnumbered cover, so it can differ from the 1-based physical slide index (SOT-2546 idx84: the F1-ranking
# slide is physically the 6th but its footer reads "5 / 15", and the gold page number is 5).
_PAGE_FOOTER_RE = re.compile(r"^\s*(\d+)\s*[/／]\s*(\d+)\s*$")


def _slide_display_pages(prs: Presentation) -> dict[int, int]:
    """Map 1-based physical slide index → the page number printed in that slide's ``N / TOTAL`` footer.

    Only footers whose denominator equals the document-wide total (the modal denominator) are kept, so a
    stray ``1 / 7``-style list numbering on a content slide is not mistaken for a page footer.  Returns
    ``{}`` when no consistent footer scheme is present (callers then fall back to the physical index).
    """
    from collections import Counter

    found: dict[int, tuple[int, int]] = {}
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            match = _PAGE_FOOTER_RE.match(nfc((shape.text or "").strip()))
            if match:
                found[si] = (int(match.group(1)), int(match.group(2)))
                break
    if not found:
        return {}
    total = Counter(tot for _num, tot in found.values()).most_common(1)[0][0]
    return {si: num for si, (num, tot) in found.items() if tot == total}


def extract_pptx(ref: FileRef, data: bytes | None) -> str:
    prs = Presentation(io.BytesIO(data) if data else str(ref.path))
    out: list[str] = []
    display_pages = _slide_display_pages(prs)
    gantt = extract_gantt_week_ranges(prs)
    if gantt:
        out.append("【ガント週グリッド:決定論】")
        for record in gantt:
            if record["status"] == "resolved":
                out.append(
                    f"[スライド{record['slide']}] {record['activity']}: "
                    f"第{record['start_week']}週目から第{record['end_week']}週目 "
                    f"(週ヘッダx座標×バーleft/width、半開区間の正の重なり)")
            else:
                choices = " / ".join(
                    f"第{candidate['start_week']}週目から第{candidate['end_week']}週目"
                    for candidate in record["candidates"])
                out.append(
                    f"[スライド{record['slide']}] {record['activity']}: 曖昧({choices})。"
                    "決定論候補が競合するため回答確定不可")
        out.append("【/ガント週グリッド】")
    for si, slide in enumerate(prs.slides, 1):
        page = display_pages.get(si)
        if page is not None:
            # Page questions ("何ページ") mean the printed footer number, which excludes an unnumbered
            # cover — surface it so the answer is the document's page number, not the physical index.
            out.append(f"[スライド{si}（文書に記載のページ番号: {page}）]")
        else:
            out.append(f"[スライド{si}]")
        # top-to-bottom, left-to-right ordering (per enumeration rules in the task)
        shapes = sorted(slide.shapes, key=lambda s: (int(getattr(s, "top", 0) or 0),
                                                     int(getattr(s, "left", 0) or 0)))
        for shape in shapes:
            fillc = _shape_fill_color(shape)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs) or para.text
                    marks = []
                    for r in para.runs:
                        hl = _run_highlight(r)
                        if hl and r.text.strip():
                            marks.append(f"【ハイライト:{hl}】{r.text.strip()}")
                    if line.strip():
                        out.append(line)
                    out.extend(marks)
                if fillc and shape.text_frame.text.strip():
                    out.append(f"【図形塗り:{fillc}】{shape.text_frame.text.strip()[:60]}")
            if shape.has_table:
                out.append("[表]")
                for r in shape.table.rows:
                    out.append(" | ".join(c.text for c in r.cells))
    return "\n".join(out)

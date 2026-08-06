"""SOT-2482: highlight / marker-word extractor (``highlight_extract``).

Covers the acceptance criteria:
  ① a highlight/marker question returns the correct word/cell set — the highlighted runs/cells/spans
     are extracted, in document order, and a colour filter narrows to one colour;
  ② every result is the common contract ``{value, evidence, method}`` — outer AND per-item;
  ③ 固有事実非同梱 / portability — an arbitrary colour+word (never a corpus fact) is extracted, and the
     colour vocabulary is the generic classifier's, so the tool works on any document.

The behaviour is pinned against tiny, hand-built xlsx/pptx/docx/pdf fixtures we control exactly (fast,
no heavy corpus dependency), plus a guarded smoke test over the real share drive, where 青嶺不動産
アセットマネジメントのスケジュール carries genuine orange row fills.
"""
from __future__ import annotations

from pathlib import Path

import pytest

from src.rag.corpus import walk
from src.rag.tools import ContractError, highlight_extract, highlight_words, is_contract


# --------------------------------------------------------------------------- fixture builders
def _xlsx(path: Path, cells: dict[str, tuple[str, str]]) -> Path:
    """Write a 1-sheet xlsx placing each ``coord -> (value, argb_fill)``; empty fill = no highlight."""
    import openpyxl
    from openpyxl.styles import PatternFill

    wb = openpyxl.Workbook()
    ws = wb.active
    for coord, (value, argb) in cells.items():
        ws[coord] = value
        if argb:
            ws[coord].fill = PatternFill(fill_type="solid", fgColor=argb)
    wb.save(path)
    return path


def _pptx_highlight(path: Path, runs: list[tuple[str, str | None]]) -> Path:
    """Write a 1-slide pptx; each ``(text, srgb_hex)`` becomes a run, hex=None → no highlight mark."""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tf = box.text_frame
    para = tf.paragraphs[0]
    for text, srgb in runs:
        run = para.add_run()
        run.text = text
        if srgb:
            rpr = run._r.get_or_add_rPr()
            hl = rpr.makeelement(qn("a:highlight"), {})
            clr = hl.makeelement(qn("a:srgbClr"), {"val": srgb})
            hl.append(clr)
            rpr.append(hl)
    prs.save(path)
    return path


def _docx_highlight(path: Path, runs: list[tuple[str, str | None]]) -> Path:
    """Write a docx; each ``(text, wd_color_name)`` becomes a run, name=None → no highlight."""
    import docx
    from docx.enum.text import WD_COLOR_INDEX

    d = docx.Document()
    para = d.add_paragraph()
    for text, name in runs:
        run = para.add_run(text)
        if name:
            run.font.highlight_color = getattr(WD_COLOR_INDEX, name)
    d.save(path)
    return path


def _pdf(path: Path, *, highlight_word: str, colored: tuple[str, tuple[float, float, float]] | None):
    """Write a 1-page PDF: plain text with one Highlight annotation and optional coloured text."""
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 100), f"plain {highlight_word} plain", fontsize=14)
    for rect in page.search_for(highlight_word):
        page.add_highlight_annot(rect)
    if colored is not None:
        text, rgb = colored
        page.insert_text((72, 160), text, fontsize=14, color=rgb)
    doc.save(str(path))
    doc.close()
    return path


def _contract_ok(result) -> bool:
    """Outer result and every per-item hit are valid ``{value, evidence, method}`` contracts."""
    return is_contract(result) and all(is_contract(item) for item in result["value"])


# --------------------------------------------------------------------------- xlsx
def test_xlsx_extracts_filled_cells_in_order(tmp_path):
    p = _xlsx(tmp_path / "s.xlsx", {
        "A1": ("keep", ""),
        "B1": ("alpha", "FFFF00"),   # yellow
        "A2": ("beta", "ED7D31"),    # orange
        "C2": ("gamma", ""),         # no fill
    })
    result = highlight_extract(p)
    assert _contract_ok(result)
    values = [it["value"] for it in result["value"]]
    assert values == ["alpha", "beta"]                       # document (row/col) order, unfilled dropped
    colors = {it["value"]: it["method"]["color"] for it in result["value"]}
    assert colors == {"alpha": "黄", "beta": "オレンジ"}
    assert result["value"][0]["evidence"]["cell"] == "B1"
    assert result["value"][0]["method"]["engine"] == "xlsx"


def test_xlsx_color_filter(tmp_path):
    p = _xlsx(tmp_path / "s.xlsx", {"A1": ("alpha", "FFFF00"), "A2": ("beta", "ED7D31")})
    assert highlight_words(p, color="オレンジ") == ["beta"]
    assert highlight_words(p, color="yellow") == ["alpha"]   # english synonym
    assert highlight_words(p, color="オレンジ色") == ["beta"]  # trailing 色 stripped


# --------------------------------------------------------------------------- pptx
def test_pptx_extracts_highlighted_runs(tmp_path):
    p = _pptx_highlight(tmp_path / "d.pptx",
                        [("normal ", None), ("marked", "FFFF00"), (" tail", None)])
    result = highlight_extract(p)
    assert _contract_ok(result)
    hits = [it for it in result["value"] if it["method"]["kind"] == "run_highlight"]
    assert [h["value"] for h in hits] == ["marked"]
    assert hits[0]["method"]["color"] == "黄"
    assert hits[0]["evidence"]["slide"] == 1


# --------------------------------------------------------------------------- docx
def test_docx_extracts_highlighted_runs(tmp_path):
    p = _docx_highlight(tmp_path / "d.docx",
                        [("plain ", None), ("hit", "YELLOW"), (" more ", None), ("two", "GREEN")])
    result = highlight_extract(p)
    assert _contract_ok(result)
    got = {it["value"]: it["method"]["color"] for it in result["value"]}
    assert got == {"hit": "黄", "two": "緑"}
    assert highlight_words(p, color="緑") == ["two"]


# --------------------------------------------------------------------------- pdf
def test_pdf_annotation_and_colored_text(tmp_path):
    p = _pdf(tmp_path / "d.pdf", highlight_word="TARGET", colored=("REDWORD", (1.0, 0.0, 0.0)))
    result = highlight_extract(p)
    assert _contract_ok(result)
    kinds = {it["method"]["kind"] for it in result["value"]}
    assert "annotation" in kinds
    ann = [it for it in result["value"] if it["method"]["kind"] == "annotation"]
    assert any("TARGET" in it["value"] for it in ann)
    colored = [it for it in result["value"] if it["method"]["kind"] == "colored_text"]
    assert any("REDWORD" in it["value"] and it["method"]["color"] == "赤" for it in colored)


def test_image_only_pdf_uses_generic_vision_fallback(tmp_path, monkeypatch):
    """Raster PDFs are interpreted from pixels; the response is not tied to corpus words."""
    import fitz
    from src.rag import llm

    p = tmp_path / "scan.pdf"
    png = tmp_path / "page.png"
    doc = fitz.open()
    page = doc.new_page(width=300, height=200)
    page.draw_rect((30, 60, 180, 100), color=(.7, .7, .7), fill=(.7, .7, .7))
    page.insert_text((40, 88), "NOVEL_MARKER", fontsize=14)
    pix = page.get_pixmap(alpha=False)
    pix.save(png)
    doc.close()
    scanned = fitz.open()
    target = scanned.new_page(width=300, height=200)
    target.insert_image(target.rect, filename=str(png))
    scanned.save(p)
    scanned.close()

    monkeypatch.setattr(llm, "generate", lambda *a, **k:
                        '[{"page": 1, "value": "NOVEL_MARKER", "color": "gray"}]')
    result = highlight_extract(p)
    assert _contract_ok(result)
    assert [x["value"] for x in result["value"]] == ["NOVEL_MARKER"]
    assert result["value"][0]["method"] == {
        "engine": "pdf_vision", "color": "灰", "kind": "image_marker"}
    assert result["value"][0]["evidence"]["page"] == 1


# --------------------------------------------------------------------------- portability / errors
def test_arbitrary_color_and_word_is_generic(tmp_path):
    """A colour+word that is not any corpus fact still extracts — the tool hard-codes no fixture."""
    p = _xlsx(tmp_path / "s.xlsx", {"A1": ("ZZZ_novel_token", "7030A0")})  # purple
    assert highlight_words(p, color="紫") == ["ZZZ_novel_token"]


def test_unsupported_extension_raises(tmp_path):
    bad = tmp_path / "x.txt"
    bad.write_text("nope", encoding="utf-8")
    with pytest.raises(ContractError):
        highlight_extract(bad)


# --------------------------------------------------------------------------- real-corpus smoke
def _corpus_ref(substr: str, ext: str):
    for r in walk():
        if substr in r.rel and r.ext == ext:
            return r
    return None


@pytest.mark.skipif(_corpus_ref("スケジュール_r2", "xlsx") is None,
                    reason="share drive not present")
def test_real_corpus_orange_rows_smoke():
    ref = _corpus_ref("スケジュール_r2", "xlsx")
    words = highlight_words(ref, color="オレンジ")
    # the orange-filled schedule rows include the kickoff task name
    assert any("プロジェクトキックオフ" in w for w in words)
